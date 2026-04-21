from __future__ import annotations

import json
import shutil
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

ROOT_DIR = Path(__file__).resolve().parents[2]
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

from backend.app.tools.handlers import (  # noqa: E402
    HANDLERS,
    build_epub_from_text,
    build_odt_from_text,
    build_rtf_from_text,
    create_job_workspace,
)


@dataclass
class Case:
    slug: str
    files: list[str]
    payload: dict[str, Any]
    optional: bool = False


def _copy_if_exists(source: Path, target: Path) -> bool:
    if not source.exists():
        return False
    shutil.copy(source, target)
    return True


def _create_demo_image(path: Path, text: str, color: tuple[int, int, int]) -> None:
    image = Image.new("RGB", (820, 520), color)
    draw = ImageDraw.Draw(image)
    draw.text((24, 24), text, fill=(255, 255, 255))
    image.save(path)


def _create_demo_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.setTitle("ISHU TOOLS Smoke Test PDF")
    width, height = A4
    lines = [
        "ISHU TOOLS SMOKE TEST PDF",
        "This generated PDF validates merge, split, convert, text extraction,",
        "metadata, watermark, OCR fallback, and export workflows.",
        "Sample table: name Ishu score 95.",
    ]
    y = height - 72
    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(72, y, lines[0])
    pdf.setFont("Helvetica", 11)
    for line in lines[1:]:
        y -= 24
        pdf.drawString(72, y, line)
    pdf.showPage()
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(72, height - 72, "Second page for page operation checks")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(72, height - 104, "Redaction keyword: Ishu. Page order and extraction smoke content.")
    pdf.save()


def _create_ocr_text_image(path: Path) -> None:
    image = Image.new("RGB", (1800, 900), (255, 255, 255))
    draw = ImageDraw.Draw(image)

    font = None
    try:
        from PIL import ImageFont

        try:
            font = ImageFont.truetype("arial.ttf", 120)
        except Exception:
            font = ImageFont.truetype("DejaVuSans.ttf", 120)
    except Exception:
        font = None

    lines = [
        "ISHU TOOLS OCR TEST",
        "PDF IMAGE TEXT 12345",
        "SMOKE MATRIX VALIDATION",
    ]

    y = 120
    line_step = 220 if font is not None else 90
    for line in lines:
        draw.text((80, y), line, fill=(0, 0, 0), font=font)
        y += line_step

    image.save(path)


def prepare_fixtures(fixtures_dir: Path) -> dict[str, Path]:
    fixtures_dir.mkdir(parents=True, exist_ok=True)

    assets_main = ROOT_DIR / "tmp_smoke_assets"
    assets_two = ROOT_DIR / "tmp_smoke_assets_2"
    assets_three = ROOT_DIR / "tmp_smoke_assets_3"

    fixture_paths: dict[str, Path] = {}

    pdf_path = fixtures_dir / "sample.pdf"
    if not _copy_if_exists(assets_two / "sample.pdf", pdf_path):
        _create_demo_pdf(pdf_path)
    fixture_paths["pdf"] = pdf_path

    pdf_b = fixtures_dir / "sample-b.pdf"
    shutil.copy(pdf_path, pdf_b)
    fixture_paths["pdf_b"] = pdf_b

    png_path = fixtures_dir / "sample.png"
    if not _copy_if_exists(assets_two / "photo.png", png_path):
        _create_demo_image(png_path, "SMOKE SAMPLE PNG", (52, 92, 185))
    fixture_paths["png"] = png_path

    logo_path = fixtures_dir / "logo.png"
    if not _copy_if_exists(assets_main / "logo.png", logo_path):
        _create_demo_image(logo_path, "LOGO", (18, 108, 73))
    fixture_paths["logo"] = logo_path

    signature_path = fixtures_dir / "signature.png"
    if not _copy_if_exists(assets_two / "sign.png", signature_path):
        _create_demo_image(signature_path, "SIGN", (26, 80, 152))
    fixture_paths["signature"] = signature_path

    jpg_path = fixtures_dir / "sample.jpg"
    if not _copy_if_exists(assets_main / "lena.jpg", jpg_path):
        Image.open(png_path).convert("RGB").save(jpg_path, "JPEG", quality=92)
    fixture_paths["jpg"] = jpg_path

    jpeg_path = fixtures_dir / "sample.jpeg"
    Image.open(jpg_path).convert("RGB").save(jpeg_path, "JPEG", quality=92)
    fixture_paths["jpeg"] = jpeg_path

    ocr_png = fixtures_dir / "ocr-text.png"
    _create_ocr_text_image(ocr_png)
    fixture_paths["ocr_png"] = ocr_png

    ocr_jpg = fixtures_dir / "ocr-text.jpg"
    Image.open(ocr_png).convert("RGB").save(ocr_jpg, "JPEG", quality=95)
    fixture_paths["ocr_jpg"] = ocr_jpg

    ocr_pdf = fixtures_dir / "ocr-scan.pdf"
    Image.open(ocr_png).convert("RGB").save(ocr_pdf, "PDF", resolution=300)
    fixture_paths["ocr_pdf"] = ocr_pdf

    heic_path = fixtures_dir / "sample.heic"
    if _copy_if_exists(assets_three / "sample.heic", heic_path):
        fixture_paths["heic"] = heic_path

    txt_path = fixtures_dir / "sample.txt"
    if not _copy_if_exists(assets_two / "note.txt", txt_path):
        txt_path.write_text("Ishu tools smoke test text sample.", encoding="utf-8")
    fixture_paths["txt"] = txt_path

    md_path = fixtures_dir / "sample.md"
    md_path.write_text("# Smoke Test\n\nThis is markdown content for smoke testing.", encoding="utf-8")
    fixture_paths["md"] = md_path

    json_path = fixtures_dir / "sample.json"
    json_path.write_text(json.dumps([{"name": "Ishu", "score": 95}], indent=2), encoding="utf-8")
    fixture_paths["json"] = json_path

    xml_path = fixtures_dir / "sample.xml"
    xml_path.write_text("<root><name>Ishu</name><score>95</score></root>", encoding="utf-8")
    fixture_paths["xml"] = xml_path

    csv_path = fixtures_dir / "sample.csv"
    csv_path.write_text("name,score\nIshu,95\nAvi,88\n", encoding="utf-8")
    fixture_paths["csv"] = csv_path

    docx_path = fixtures_dir / "sample.docx"
    docx = Document()
    docx.add_heading("Smoke Test DOCX", level=1)
    docx.add_paragraph("This DOCX file is generated for conversion smoke tests.")
    docx.save(str(docx_path))
    fixture_paths["docx"] = docx_path

    xlsx_path = fixtures_dir / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    if sheet is None:
        sheet = workbook.create_sheet("Data")
    sheet.title = "Data"
    sheet.append(["name", "score"])
    sheet.append(["Ishu", 95])
    sheet.append(["Avi", 88])
    workbook.save(xlsx_path)
    fixture_paths["xlsx"] = xlsx_path

    pptx_path = fixtures_dir / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Smoke Test PPTX"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.9), Inches(8.5), Inches(1.6))
    textbox.text_frame.text = "Generated slide content for conversion testing."
    presentation.save(str(pptx_path))
    fixture_paths["pptx"] = pptx_path

    rtf_path = fixtures_dir / "sample.rtf"
    build_rtf_from_text("Smoke test RTF content.", rtf_path)
    fixture_paths["rtf"] = rtf_path

    odt_path = fixtures_dir / "sample.odt"
    build_odt_from_text("Smoke test ODT content.", odt_path, "Smoke ODT")
    fixture_paths["odt"] = odt_path

    epub_path = fixtures_dir / "sample.epub"
    build_epub_from_text("Smoke test EPUB content.", epub_path, "Smoke EPUB")
    fixture_paths["epub"] = epub_path

    svg_path = fixtures_dir / "sample.svg"
    svg_path.write_text(
        "<svg xmlns='http://www.w3.org/2000/svg' width='300' height='180'>"
        "<rect width='300' height='180' fill='#1f2937'/>"
        "<text x='20' y='96' font-size='24' fill='#f9fafb'>Smoke SVG</text>"
        "</svg>",
        encoding="utf-8",
    )
    fixture_paths["svg"] = svg_path

    eml_path = fixtures_dir / "sample.eml"
    if not _copy_if_exists(assets_two / "mail.eml", eml_path):
        eml_path.write_text(
            "Subject: Smoke Test Email\nFrom: smoke@example.com\nTo: test@example.com\n\nHello from smoke test.",
            encoding="utf-8",
        )
    fixture_paths["eml"] = eml_path

    fb2_path = fixtures_dir / "sample.fb2"
    fb2_path.write_text(
        "<FictionBook><body><section><title><p>Smoke FB2</p></title><p>Sample body text.</p></section></body></FictionBook>",
        encoding="utf-8",
    )
    fixture_paths["fb2"] = fb2_path

    cbz_path = fixtures_dir / "sample.cbz"
    if not _copy_if_exists(assets_two / "comic.cbz", cbz_path):
        with zipfile.ZipFile(cbz_path, "w") as archive:
            archive.write(jpg_path, arcname="page-1.jpg")
    fixture_paths["cbz"] = cbz_path

    cbr_path = fixtures_dir / "sample.cbr"
    with zipfile.ZipFile(cbr_path, "w") as archive:
        archive.write(jpg_path, arcname="page-1.jpg")
        archive.write(png_path, arcname="page-2.png")
    fixture_paths["cbr"] = cbr_path

    zip_images = fixtures_dir / "images.zip"
    with zipfile.ZipFile(zip_images, "w") as archive:
        archive.write(jpg_path, arcname="one.jpg")
        archive.write(png_path, arcname="two.png")
    fixture_paths["zip_images"] = zip_images

    return fixture_paths


def build_cases() -> list[Case]:
    return [
        Case("merge-pdf", ["pdf", "pdf_b"], {}),
        Case("split-pdf", ["pdf"], {"pages": "1"}),
        Case("extract-pages", ["pdf"], {"pages": "1"}),
        Case("compress-pdf", ["pdf"], {}),
        Case("optimize-pdf", ["pdf"], {}),
        Case("rotate-pdf", ["pdf"], {"angle": 90}),
        Case("organize-pdf", ["pdf"], {"page_order": "1"}),
        Case("crop-pdf", ["pdf"], {"left": 1, "right": 1, "top": 1, "bottom": 1}),
        Case("watermark-pdf", ["pdf"], {"text": "ISHU TOOLS"}),
        Case("page-numbers-pdf", ["pdf"], {}),
        Case("protect-pdf", ["pdf"], {"password": "1234"}),
        Case("unlock-pdf", ["pdf"], {"password": ""}),
        Case("redact-pdf", ["pdf"], {"keywords": "Ishu"}),
        Case("compare-pdf", ["pdf", "pdf_b"], {}),
        Case("extract-text-pdf", ["pdf"], {}),
        Case("extract-images-pdf", ["pdf"], {}, optional=True),
        Case("pdf-to-jpg", ["pdf"], {}),
        Case("pdf-to-png", ["pdf"], {}),
        Case("jpg-to-pdf", ["jpg", "png"], {}),
        Case("jpeg-to-pdf", ["jpeg"], {}),
        Case("image-to-pdf", ["jpg", "png"], {}),
        Case("scan-to-pdf", ["jpg", "png"], {}),
        Case("repair-pdf", ["pdf"], {}),
        Case("translate-pdf", ["pdf"], {"target_lang": "en"}),
        Case("summarize-pdf", ["pdf"], {}),
        Case("pdf-to-docx", ["pdf"], {}),
        Case("docx-to-pdf", ["docx"], {}),
        Case("pdf-to-excel", ["pdf"], {}),
        Case("excel-to-pdf", ["xlsx"], {}),
        Case("pdf-to-pptx", ["pdf"], {}),
        Case("pdf-to-ppt", ["pdf"], {}),
        Case("pptx-to-pdf", ["pptx"], {}),
        Case("ppt-to-pdf", ["pptx"], {}),
        Case("pdf-to-word", ["pdf"], {}),
        Case("word-to-pdf", ["docx"], {}),
        Case("pdf-to-powerpoint", ["pdf"], {}),
        Case("powerpoint-to-pdf", ["pptx"], {}),
        Case("sign-pdf", ["pdf"], {"signer": "Ishu"}),
        Case("edit-metadata-pdf", ["pdf"], {"title": "Smoke Title", "author": "Ishu"}),
        Case("edit-metadata", ["jpg"], {"title": "Image Title", "author": "Ishu"}),
        Case("pdf-viewer", ["pdf"], {}),
        Case("pdf-intelligence", ["pdf"], {}),
        Case("pdf-to-pdfa", ["pdf"], {}),
        Case("annotate-pdf", ["pdf"], {"text": "Smoke note"}),
        Case("highlight-pdf", ["pdf"], {"page": 1, "x": 32, "y": 32, "width": 120, "height": 22}),
        Case("add-text-pdf", ["pdf"], {"text": "HELLO"}),
        Case("edit-pdf-text", ["pdf"], {"text": "HELLO"}),
        Case("add-text", ["pdf"], {"text": "HELLO"}),
        Case("add-image-to-pdf", ["pdf", "logo"], {}),
        Case("header-footer-pdf", ["pdf"], {"header": "H", "footer": "F"}),
        Case("resize-pages-pdf", ["pdf"], {"scale": 0.95}),
        Case("remove-pages", ["pdf"], {"pages": "1"}, optional=True),
        Case("add-page-numbers", ["pdf"], {}),
        Case("add-watermark", ["pdf"], {"text": "ISHU"}),
        Case("remove-metadata-pdf", ["pdf"], {}),
        Case("whiteout-pdf", ["pdf"], {"keywords": "Ishu"}),
        Case("grayscale-pdf", ["pdf"], {}),
        Case("create-pdf", [], {"text": "Create PDF from plain text"}),
        Case("html-to-pdf", [], {"html": "<h1>Smoke</h1><p>Test</p>"}),
        Case("url-to-pdf", [], {"url": "https://example.com"}, optional=True),
        Case("md-to-pdf", ["md"], {}),
        Case("txt-to-pdf", ["txt"], {}),
        Case("json-to-pdf", ["json"], {}),
        Case("xml-to-pdf", ["xml"], {}),
        Case("csv-to-pdf", ["csv"], {}),
        Case("convert-to-pdf", ["docx"], {}),
        Case("convert-from-pdf", ["pdf"], {"target_format": "png"}),
        Case("pdf-converter", ["pdf"], {"target_format": "jpg"}),
        Case("compress-image", ["jpg", "png"], {"quality": 72}),
        Case("resize-image", ["jpg"], {"width": 640, "height": 420}),
        Case("crop-image", ["jpg"], {"x": 0, "y": 0, "width": 200, "height": 200}),
        Case("rotate-image", ["jpg"], {"angle": 90}),
        Case("convert-image", ["jpg"], {"target_format": "png"}),
        Case("watermark-image", ["jpg"], {"text": "ISHU"}),
        Case("grayscale-image", ["jpg"], {}),
        Case("blur-image", ["jpg"], {"radius": 2}),
        Case("pixelate-image", ["jpg"], {"factor": 10}),
        Case("meme-generator", ["jpg"], {"top_text": "TOP", "bottom_text": "BOTTOM"}),
        Case("flip-image", ["jpg"], {"mode": "horizontal"}),
        Case("add-border-image", ["jpg"], {"border_size": 12, "color": "#ffffff"}),
        Case("thumbnail-image", ["jpg"], {"max_width": 240, "max_height": 240}),
        Case("image-collage", ["jpg", "png"], {"columns": 2, "cell_size": 220}),
        Case("sharpen-image", ["jpg"], {"factor": 1.6}),
        Case("brighten-image", ["jpg"], {"factor": 1.2}),
        Case("contrast-image", ["jpg"], {"factor": 1.2}),
        Case("invert-image", ["jpg"], {}),
        Case("posterize-image", ["jpg"], {"bits": 4}),
        Case("image-histogram", ["jpg"], {}),
        Case("upscale-image", ["jpg"], {"scale": 2}),
        Case("add-text-image", ["jpg"], {"text": "Smoke", "x": 24, "y": 24}),
        Case("add-logo-image", ["jpg", "logo"], {"position": "bottom-right"}),
        Case("join-images", ["jpg", "png"], {"direction": "horizontal"}),
        Case("split-image", ["jpg"], {"columns": 2, "rows": 2}),
        Case("circle-crop-image", ["jpg"], {"size": 256}),
        Case("square-crop-image", ["jpg"], {"size": 512}),
        Case("motion-blur-image", ["jpg"], {"direction": "horizontal", "strength": 9}),
        Case("check-image-dpi", ["jpg"], {}),
        Case("check-dpi", ["jpg"], {}),
        Case("convert-dpi", ["jpg"], {"dpi": 300}),
        Case("resize-image-in-cm", ["jpg"], {"width_cm": 3.5, "height_cm": 4.5, "dpi": 300}),
        Case("resize-image-in-mm", ["jpg"], {"width_mm": 35, "height_mm": 45, "dpi": 300}),
        Case("resize-image-in-inch", ["jpg"], {"width_inch": 2, "height_inch": 2, "dpi": 300}),
        Case("add-name-dob-image", ["jpg"], {"name": "Ishu", "dob": "01-01-2000"}),
        Case("merge-photo-signature", ["jpg", "logo"], {"direction": "vertical"}),
        Case("black-and-white-image", ["jpg"], {"threshold": 128}),
        Case("generate-signature", [], {"text": "Ishu Kumar"}),
        Case("picture-to-pixel-art", ["jpg"], {}),
        Case("censor-photo", ["jpg"], {}, optional=True),
        Case("remove-background", ["jpg"], {}, optional=True),
        Case("blur-background", ["jpg"], {"radius": 10}),
        Case("blur-face", ["jpg"], {"blur_strength": 35}, optional=True),
        Case("unblur-face", ["jpg"], {"strength": 1.8, "denoise": 4}, optional=True),
        Case("pixelate-face", ["jpg"], {"pixel_size": 12}, optional=True),
        Case("remove-image-object", ["jpg"], {"x": 60, "y": 60, "width": 120, "height": 120, "radius": 5}),
        Case("ocr-image", ["ocr_png"], {}, optional=True),
        Case("ocr-pdf", ["ocr_pdf"], {}, optional=True),
        Case("pdf-ocr", ["ocr_pdf"], {}, optional=True),
        Case("image-to-text", ["ocr_png"], {}, optional=True),
        Case("jpg-to-text", ["ocr_jpg"], {}, optional=True),
        Case("png-to-text", ["ocr_png"], {}, optional=True),
        Case("summarize-text", ["txt"], {}),
        Case("translate-text", ["txt"], {"target_lang": "en"}),
        Case("qr-code-generator", [], {"text": "https://ishu.tools"}),
        Case("extract-metadata", ["jpg"], {}),
        Case("remove-metadata-image", ["jpg"], {}),
        Case("word-count-text", ["txt"], {}),
        Case("case-converter-text", ["txt"], {"mode": "upper"}),
        Case("extract-keywords-text", ["txt"], {"top_n": 8}),
        Case("slug-generator-text", [], {"text": "Ishu Tools Platform"}),
        Case("remove-extra-spaces-text", [], {"text": "A   B    C"}),
        Case("sort-lines-text", [], {"text": "c\na\nb", "direction": "asc"}),
        Case("deduplicate-lines-text", [], {"text": "a\na\nb\na"}),
        Case("find-replace-text", [], {"text": "hello ishu", "find": "ishu", "replace": "tools"}),
        Case("split-text-file", [], {"text": "line1\nline2\nline3", "lines_per_file": 1}),
        Case("reading-time-text", [], {"text": "one two three four five", "wpm": 200}),
        Case("create-workflow", [], {"name": "Sample Workflow", "steps": ["merge-pdf", "compress-pdf"]}),
        Case("pdf-pages-to-zip", ["pdf"], {}),
        Case("zip-images-to-pdf", ["zip_images"], {}),
        Case("zip-to-pdf", ["zip_images"], {}),
        Case("images-to-zip", ["jpg", "png"], {}),
        Case("batch-convert-images", ["jpg", "png"], {"target_format": "webp"}),
        Case("merge-text-files", ["txt", "md"], {"separator": "\n---\n"}),
        Case("json-prettify", ["json"], {}),
        Case("csv-to-json", ["csv"], {}),
        Case("json-to-csv", ["json"], {}),
        Case("pdf-page-count", ["pdf"], {}),
        Case("reverse-pdf", ["pdf"], {}),
        Case("flatten-pdf", ["pdf"], {}),
        Case("pdf-to-html", ["pdf"], {}),
        Case("pdf-to-bmp", ["pdf"], {}),
        Case("pdf-to-gif", ["pdf"], {}),
        Case("png-to-pdf", ["png"], {}),
        Case("webp-to-pdf", ["jpg"], {}),
        Case("gif-to-pdf", ["jpg"], {}),
        Case("bmp-to-pdf", ["jpg"], {}),
        Case("jpg-to-png", ["jpg"], {}),
        Case("png-to-jpg", ["png"], {}),
        Case("image-to-webp", ["jpg"], {}),
        Case("pdf-to-txt", ["pdf"], {}),
        Case("pdf-to-markdown", ["pdf"], {}),
        Case("pdf-to-json", ["pdf"], {}),
        Case("pdf-to-csv", ["pdf"], {}),
        Case("pdf-to-tiff", ["pdf"], {}),
        Case("tiff-to-pdf", ["jpg"], {}),
        Case("pdf-to-svg", ["pdf"], {}),
        Case("svg-to-pdf", ["svg"], {}, optional=True),
        Case("pdf-to-rtf", ["pdf"], {}),
        Case("rtf-to-pdf", ["rtf"], {}),
        Case("pdf-to-odt", ["pdf"], {}),
        Case("odt-to-pdf", ["odt"], {}),
        Case("pdf-to-epub", ["pdf"], {}),
        Case("epub-to-pdf", ["epub"], {}),
        Case("eml-to-pdf", ["eml"], {}),
        Case("fb2-to-pdf", ["fb2"], {}),
        Case("cbz-to-pdf", ["cbz"], {}),
        Case("ebook-to-pdf", ["epub"], {}),
        Case("cbr-to-pdf", ["cbr"], {}),
        Case("html-to-image", [], {"html": "<h1>Smoke</h1>", "format": "png"}),
        Case("reduce-image-size-in-kb", ["jpg"], {"target_kb": 120}),
        Case("compress-to-kb", ["jpg"], {"target_kb": 80}),
        Case("passport-photo-maker", ["jpg"], {"preset": "3.5x4.5cm", "dpi": 300}),
        Case("passport-size-photo", ["jpg"], {"preset": "2x2inch", "dpi": 300}),
        Case("social-media-resize", ["jpg"], {"platform": "instagram"}),
        Case("resize-for-instagram", ["jpg"], {}),
        Case("resize-for-whatsapp", ["jpg"], {}),
        Case("resize-for-youtube", ["jpg"], {}),
        Case("instagram-grid", ["jpg"], {"rows": 2, "cols": 2}),
        Case("convert-to-jpg", ["png"], {}),
        Case("convert-from-jpg", ["jpg"], {"target_format": "png"}),
        Case("bulk-image-resizer", ["jpg", "png"], {"width": 512, "height": 512}),
        Case("heic-to-jpg", ["heic"], {}, optional=True),
        Case("webp-to-jpg", ["jpg"], {}),
        Case("jpeg-to-jpg", ["jpeg"], {}),
        Case("jpeg-to-png", ["jpg"], {}),
        Case("png-to-jpeg", ["png"], {}),
        Case("photo-editor", ["jpg"], {"mode": "enhance"}),
        Case("unblur-image", ["jpg"], {"factor": 2.0}),
        Case("increase-image-quality", ["jpg"], {}),
        Case("beautify-image", ["jpg"], {}),
        Case("retouch-image", ["jpg"], {}),
        Case("photo-blemish-remover", ["jpg"], {}),
        Case("super-resolution", ["jpg"], {"scale": 2}),
        Case("zoom-out-image", ["jpg"], {"border_percent": 20}),
        Case("add-white-border-image", ["jpg"], {"border_size": 20}),
        Case("freehand-crop", ["jpg"], {"x1": 10, "y1": 10, "x2": 200, "y2": 200}),
        Case("crop-png", ["png"], {"x": 0, "y": 0, "width": 100, "height": 100}),
        Case("image-splitter", ["jpg"], {"rows": 2, "columns": 2}),
        Case("resize-image-pixel", ["jpg"], {"width": 620, "height": 620}),
        Case("resize-signature", ["signature"], {"width_mm": 50, "height_mm": 20, "dpi": 300}),
        Case("resize-image-to-3.5cmx4.5cm", ["jpg"], {"dpi": 300}),
        Case("resize-image-to-6cmx2cm", ["jpg"], {"dpi": 300}),
        Case("resize-signature-to-50mmx20mm", ["signature"], {"dpi": 300}),
        Case("resize-image-to-35mmx45mm", ["jpg"], {"dpi": 300}),
        Case("resize-image-to-2x2", ["jpg"], {"dpi": 300}),
        Case("resize-image-to-3x4", ["jpg"], {"dpi": 300}),
        Case("resize-image-to-4x6", ["jpg"], {"dpi": 300}),
        Case("resize-image-to-600x600-pixel", ["jpg"], {}),
        Case("resize-image-for-whatsapp-dp", ["jpg"], {}),
        Case("resize-image-for-youtube-banner", ["jpg"], {}),
        Case("resize-image-to-a4-size", ["jpg"], {"dpi": 300}),
        Case("ssc-photo-resizer", ["jpg"], {"dpi": 300}),
        Case("resize-for-pan-card", ["jpg"], {"dpi": 300}),
        Case("resize-image-for-upsc", ["jpg"], {"dpi": 300}),
        Case("psc-photo-resize", ["jpg"], {"dpi": 300}),
        Case("a4-size-resize", ["jpg"], {"dpi": 300}),
        Case("color-code-from-image", ["jpg"], {"colors": 5}),
        Case("view-metadata", ["jpg"], {}),
        Case("remove-image-metadata", ["jpg"], {}),
        Case("increase-image-size-in-kb", ["jpg"], {"target_kb": 260}),
        Case("reduce-image-size-in-mb", ["jpg"], {"target_mb": 0.2}),
        Case("convert-image-from-mb-to-kb", ["jpg"], {"target_kb": 140}),
        Case("convert-image-size-kb-to-mb", ["jpg"], {"target_mb": 0.8}),
        Case("jpg-to-kb", ["jpg"], {"target_kb": 90}),
        Case("compress-to-5kb", ["jpg"], {}),
        Case("compress-image-to-5kb", ["jpg"], {}),
        Case("compress-to-10kb", ["jpg"], {}),
        Case("compress-jpeg-to-10kb", ["jpeg"], {}),
        Case("compress-to-15kb", ["jpg"], {}),
        Case("compress-image-to-15kb", ["jpg"], {}),
        Case("compress-to-20kb", ["jpg"], {}),
        Case("compress-image-to-20kb", ["jpg"], {}),
        Case("compress-jpeg-between-20kb-to-50kb", ["jpeg"], {}),
        Case("compress-to-25kb", ["jpg"], {}),
        Case("compress-jpeg-to-25kb", ["jpeg"], {}),
        Case("compress-to-30kb", ["jpg"], {}),
        Case("compress-jpeg-to-30kb", ["jpeg"], {}),
        Case("compress-to-40kb", ["jpg"], {}),
        Case("compress-jpeg-to-40kb", ["jpeg"], {}),
        Case("compress-to-50kb", ["jpg"], {}),
        Case("compress-image-to-50kb", ["jpg"], {}),
        Case("compress-to-100kb", ["jpg"], {}),
        Case("compress-image-to-100kb", ["jpg"], {}),
        Case("compress-to-150kb", ["jpg"], {}),
        Case("compress-jpeg-to-150kb", ["jpeg"], {}),
        Case("compress-to-200kb", ["jpg"], {}),
        Case("compress-image-to-200kb", ["jpg"], {}),
        Case("compress-to-300kb", ["jpg"], {}),
        Case("compress-jpeg-to-300kb", ["jpeg"], {}),
        Case("compress-to-500kb", ["jpg"], {}),
        Case("compress-jpeg-to-500kb", ["jpeg"], {}),
        Case("compress-to-1mb", ["jpg"], {}),
        Case("compress-image-to-1mb", ["jpg"], {}),
        Case("compress-to-2mb", ["jpg"], {}),
        Case("compress-image-to-2mb", ["jpg"], {}),
        Case("jpg-to-pdf-under-50kb", ["jpg"], {}),
        Case("jpg-to-pdf-under-100kb", ["jpg"], {}),
        Case("jpeg-to-pdf-under-200kb", ["jpeg"], {}),
        Case("jpg-to-pdf-under-300kb", ["jpg"], {}),
        Case("jpg-to-pdf-under-500kb", ["jpg"], {}),
        Case("user-agent-parser", [], {"user_agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36"}),
        Case("csv-cleaner", [], {"csv": "Name, Score\n Ishu , 95\n Tools , 100\n", "trim_cells": True, "remove_empty_rows": True}),
        Case("regex-replace", [], {"text": "Call 98765-43210 today", "pattern": r"\d{5}-\d{5}", "replacement": "[phone]"}),
        Case("weighted-gpa-calculator", [], {"courses": "Math,A,4\nPhysics,B+,3\nEnglish,A-,2", "scale": 4}),
        Case("grade-average-calculator", [], {"scores": "78,85,91,66", "max_marks": 100}),
        Case("syllabus-study-planner", [], {"topics": "Algebra, hard\nPhysics, medium\nEnglish, easy", "days": 5, "daily_hours": 2}),
        Case("citation-url-cleaner", [], {"urls": "https://example.com/article?utm_source=x&gclid=abc&id=42"}),
        Case("url-query-parser", [], {"url": "https://example.com/search?q=ishu&sort=new#top"}),
        Case("timestamp-converter", [], {"timestamp": "1700000000", "offset_minutes": 330}),
        Case("markdown-table-generator", [], {"data": "Name,Score\nIshu,95\nTools,100"}),
        Case("email-extractor", [], {"text": "Contact ishu@example.com and support@ishutools.com"}),
        Case("phone-number-extractor", [], {"text": "Call +91 98765 43210 or 011-2345-6789"}),
        Case("keyword-density-checker", [], {"text": "ISHU TOOLS free tools for students and developers. Free tools online."}),
        Case("robots-txt-generator", [], {"site_url": "https://ishutools.com", "disallow": "/api/"}),
        Case("meta-description-generator", [], {"topic": "Free PDF tools", "audience": "students", "keywords": "pdf tools, merge pdf"}),
        Case("utm-builder", [], {"url": "https://ishutools.com/tools", "source": "newsletter", "medium": "email", "campaign": "student_tools"}),
        Case("html-table-generator", [], {"data": "Name,Score\nIshu,95\nTools,100", "has_header": True}),
        Case("json-schema-generator", [], {"json": "{\"name\":\"Ishu\",\"score\":95,\"tags\":[\"tools\"]}"}),
        Case("css-clamp-generator", [], {"min_px": 16, "max_px": 32, "min_viewport": 360, "max_viewport": 1440}),
        Case("slug-bulk-generator", [], {"text": "Hello World\nISHU Tools Page"}),
        Case("table-to-csv-converter", [], {"table": "| Name | Score |\n| --- | --- |\n| Ishu | 95 |"}),
        Case("csv-column-extractor", [], {"csv": "Name,Score,City\nIshu,95,Patna\nTools,100,Web", "columns": "Name,Score"}),
        Case("css-specificity-calculator", [], {"selectors": "#app .card button:hover\nmain article h2"}),
        Case("http-status-code-lookup", [], {"code": 404}),
        Case("mime-type-lookup", [], {"extension": "pdf"}),
        Case("loan-comparison-calculator", [], {"principal": 500000, "years": 5, "rates": "8.5,10,12"}),
        Case("break-even-calculator", [], {"fixed_cost": 100000, "price_per_unit": 500, "variable_cost_per_unit": 250}),
        Case("profit-margin-calculator", [], {"cost": 100, "selling_price": 150}),
        Case("gst-reverse-calculator", [], {"inclusive_amount": 1180, "gst_rate": 18}),
        Case("discount-stack-calculator", [], {"price": 1000, "discounts": "10,5"}),
        Case("study-break-planner", [], {"total_minutes": 180, "focus_minutes": 50, "break_minutes": 10, "start_time": "09:00"}),
        Case("water-reminder-schedule", [], {"wake_time": "07:00", "sleep_time": "22:30", "glasses": 8}),
        Case("workout-plan-generator", [], {"goal": "fitness", "level": "beginner", "days_per_week": 4}),
        Case("meal-plan-generator", [], {"calories": 2000, "meals": 4, "diet": "balanced"}),
        Case("name-initials-generator", [], {"names": "Ishu Kumar\nIndian Student Hub"}),
        Case("acronym-generator", [], {"phrase": "Indian Student Hub University Tools"}),
        Case("username-generator", [], {"name": "Ishu Kumar", "keyword": "tools", "count": 12}),
        Case("youtube-thumbnail-url-generator", [], {"url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ"}),
        Case("youtube-embed-code-generator", [], {"url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ", "responsive": True}),
        Case("video-aspect-ratio-calculator", [], {"width": 1920, "height": 1080}),
        Case("video-bitrate-calculator", [], {"size_mb": 100, "duration": "00:03:00", "audio_kbps": 128}),
        Case("video-file-size-estimator", [], {"duration": "00:03:00", "video_kbps": 2500, "audio_kbps": 128}),
        Case("video-duration-calculator", [], {"duration": "01:23:45", "fps": 30}),
        Case("serp-snippet-preview", [], {"title": "Free Online Tools for Students", "description": "Use fast PDF, image, developer, finance, and study tools for free.", "url": "https://ishutools.com/tools"}),
        Case("keyword-cluster-generator", [], {"keywords": "pdf merge\nmerge pdf online\ncompress pdf\nimage compressor"}),
        Case("headline-analyzer", [], {"headline": "Best Free Online Tools for Students"}),
        Case("json-path-extractor", [], {"json": "{\"user\":{\"name\":\"Ishu\",\"scores\":[95,100]}}", "path": "user.name"}),
        Case("uuid-validator", [], {"uuid": "550e8400-e29b-41d4-a716-446655440000"}),
        Case("ulid-generator", [], {"count": 5}),
        Case("ics-calendar-generator", [], {"title": "Study Session", "start": "2026-01-01 10:00", "duration_minutes": 90}),
        Case("vcard-generator", [], {"name": "Ishu Kumar", "email": "ishu@example.com", "phone": "+91 98765 43210"}),
        Case("salary-to-hourly-calculator", [], {"annual_salary": 600000, "hours_per_week": 40, "weeks_per_year": 52}),
        Case("hourly-to-salary-calculator", [], {"hourly_rate": 500, "hours_per_week": 40, "weeks_per_year": 52}),
        Case("debt-payoff-planner", [], {"debts": "Card,25000,24\nLoan,100000,12", "extra_payment": 5000, "strategy": "avalanche"}),
        Case("goal-progress-calculator", [], {"start": 0, "current": 45, "target": 100, "daily_rate": 5}),
        Case("habit-streak-calculator", [], {"dates": "2026-01-01\n2026-01-02\n2026-01-03", "today": "2026-01-03"}),
        Case("exam-timetable-generator", [], {"subjects": "Math,Physics,English", "days": 6, "daily_hours": 3}),
        Case("flashcard-csv-generator", [], {"notes": "Photosynthesis: Plants make food using sunlight\nGravity: Force that attracts objects"}),
        Case("notes-to-quiz-generator", [], {"notes": "Photosynthesis uses sunlight to make glucose. Gravity attracts objects toward Earth."}),
        Case("simple-rubric-generator", [], {"assignment": "Essay", "criteria": "Content,Organization,Grammar,References", "points_each": 10}),
    ]


def run_case(case: Case, fixture_paths: dict[str, Path]) -> tuple[str, str]:
    handler = HANDLERS.get(case.slug)
    if not handler:
        return "fail", "handler not registered"

    missing_fixtures = [name for name in case.files if name not in fixture_paths]
    if missing_fixtures:
        if case.optional:
            return "skip", f"missing fixture(s): {', '.join(missing_fixtures)}"
        return "fail", f"missing fixture(s): {', '.join(missing_fixtures)}"

    _job_id, input_dir, output_dir = create_job_workspace()
    local_files: list[Path] = []
    for fixture_key in case.files:
        source = fixture_paths[fixture_key]
        destination = input_dir / source.name
        shutil.copy(source, destination)
        local_files.append(destination)

    try:
        result = handler(local_files, dict(case.payload), output_dir)
        if result.kind == "file":
            if not result.output_path or not result.output_path.exists():
                raise RuntimeError("missing output file")
            if result.output_path.stat().st_size <= 0 and case.slug not in {"compare-pdf"}:
                raise RuntimeError("output file is empty")
        elif result.kind == "json":
            if not isinstance(result.data, dict):
                raise RuntimeError("json result data missing")
        else:
            raise RuntimeError(f"unexpected result kind: {result.kind}")

        return "pass", result.message
    except Exception as exc:  # noqa: BLE001
        if case.optional:
            return "skip", str(exc)
        return "fail", str(exc)


def main() -> int:
    fixtures_dir = ROOT_DIR / "backend" / "storage" / "smoke-fixtures"
    fixture_paths = prepare_fixtures(fixtures_dir)
    cases = build_cases()

    passed = 0
    failed = 0
    skipped = 0
    failures: list[tuple[str, str]] = []

    print(f"Running {len(cases)} smoke cases...")
    for case in cases:
        status, detail = run_case(case, fixture_paths)
        if status == "pass":
            passed += 1
            print(f"[PASS] {case.slug}")
        elif status == "skip":
            skipped += 1
            print(f"[SKIP] {case.slug} -> {detail}")
        else:
            failed += 1
            failures.append((case.slug, detail))
            print(f"[FAIL] {case.slug} -> {detail}")

    print("\nSmoke summary")
    print(f"  Passed:  {passed}")
    print(f"  Failed:  {failed}")
    print(f"  Skipped: {skipped}")

    if failures:
        print("\nFailure details")
        for slug, detail in failures:
            print(f"  - {slug}: {detail}")

    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
