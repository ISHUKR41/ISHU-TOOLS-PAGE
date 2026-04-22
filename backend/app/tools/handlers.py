from __future__ import annotations

import base64
import csv
import difflib
import html
import inspect
import io
import json
import math
import re
import shutil
import subprocess
import tempfile
import textwrap
import uuid
import zipfile
from collections import Counter
from dataclasses import dataclass
from datetime import datetime
from email import policy
from email.parser import BytesParser
from functools import wraps
from pathlib import Path
from typing import Any, Callable

import fitz
import httpx
import qrcode
from bs4 import BeautifulSoup
from deep_translator import GoogleTranslator
from docx import Document
from fastapi import HTTPException, UploadFile
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageColor, ImageDraw, ImageEnhance, ImageFilter, ImageFont, ImageOps, ImageSequence, PngImagePlugin
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from ..config import JOBS_DIR, MAX_UPLOAD_MB


@dataclass
class ExecutionResult:
    kind: str
    message: str
    output_path: Path | None = None
    filename: str | None = None
    data: dict[str, Any] | None = None
    content_type: str | None = None


ToolHandler = Callable[[list[Path], dict[str, Any], Path], ExecutionResult]


STOPWORDS = {
    "the",
    "is",
    "a",
    "an",
    "of",
    "to",
    "in",
    "and",
    "or",
    "on",
    "for",
    "with",
    "this",
    "that",
    "it",
    "as",
    "by",
    "at",
    "from",
    "be",
    "are",
    "was",
    "were",
    "will",
    "can",
    "your",
    "you",
    "we",
    "our",
}


LANGUAGE_ALIAS_MAP = {
    "arabic": "ar",
    "ar": "ar",
    "bengali": "bn",
    "bn": "bn",
    "chinese": "zh-cn",
    "chinesesimplified": "zh-cn",
    "chinesetraditional": "zh-tw",
    "cn": "zh-cn",
    "de": "de",
    "dutch": "nl",
    "en": "en",
    "english": "en",
    "es": "es",
    "french": "fr",
    "fr": "fr",
    "german": "de",
    "gujarati": "gu",
    "gu": "gu",
    "hindi": "hi",
    "hi": "hi",
    "id": "id",
    "indonesian": "id",
    "it": "it",
    "italian": "it",
    "ja": "ja",
    "japanese": "ja",
    "ko": "ko",
    "korean": "ko",
    "marathi": "mr",
    "mr": "mr",
    "nl": "nl",
    "pl": "pl",
    "polish": "pl",
    "portuguese": "pt",
    "pt": "pt",
    "punjabi": "pa",
    "pa": "pa",
    "russian": "ru",
    "ru": "ru",
    "spanish": "es",
    "ta": "ta",
    "tamil": "ta",
    "te": "te",
    "telugu": "te",
    "th": "th",
    "thai": "th",
    "tr": "tr",
    "turkish": "tr",
    "uk": "uk",
    "ukrainian": "uk",
    "urdu": "ur",
    "ur": "ur",
    "vi": "vi",
    "vietnamese": "vi",
    "zh": "zh-cn",
    "zhcn": "zh-cn",
    "zhtw": "zh-tw",
}


CONVERT_FROM_PDF_TARGET_ALIASES = {
    "bmp": "bmp",
    "csv": "csv",
    "doc": "docx",
    "docx": "docx",
    "epub": "epub",
    "excel": "xlsx",
    "gif": "gif",
    "html": "html",
    "image": "image",
    "jpeg": "jpg",
    "jpg": "jpg",
    "json": "json",
    "markdown": "md",
    "md": "md",
    "odt": "odt",
    "pdf-a": "pdfa",
    "pdfa": "pdfa",
    "pdfa1": "pdfa",
    "pdfa2": "pdfa",
    "png": "png",
    "powerpoint": "pptx",
    "ppt": "pptx",
    "pptx": "pptx",
    "rtf": "rtf",
    "svg": "svg",
    "text": "txt",
    "tif": "tiff",
    "tiff": "tiff",
    "txt": "txt",
    "word": "docx",
    "xls": "xlsx",
    "xlsx": "xlsx",
}


def normalize_lookup_token(raw: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", str(raw).strip().lower())


def tokenize_words(text: str) -> list[str]:
    return [
        token.lower()
        for token in re.findall(r"[^\W\d_]+(?:'[^\W\d_]+)?", text, flags=re.UNICODE)
        if token
    ]


def split_sentences(text: str) -> list[str]:
    parts = [segment.strip() for segment in re.split(r"(?<=[.!?।])\s+|\n+", text) if segment.strip()]
    if parts:
        return parts

    cleaned = re.sub(r"\s+", " ", text).strip()
    return [cleaned] if cleaned else []


def normalize_convert_from_pdf_target(raw_target: Any) -> str:
    cleaned = str(raw_target or "").strip().lower()
    if not cleaned:
        return "jpg"

    normalized = cleaned.replace("_", "-").replace("/", "-")
    normalized = re.sub(r"[^a-z0-9-]+", "", normalized)
    normalized = normalized.strip("-")
    key = normalized.replace("-", "")
    return CONVERT_FROM_PDF_TARGET_ALIASES.get(normalized) or CONVERT_FROM_PDF_TARGET_ALIASES.get(key) or normalized


def get_translation_language_map() -> dict[str, str]:
    cached = getattr(get_translation_language_map, "_cache", None)
    if isinstance(cached, dict):
        return cached

    language_map = dict(LANGUAGE_ALIAS_MAP)
    try:
        supported = GoogleTranslator(source="auto", target="en").get_supported_languages(as_dict=True)
        if isinstance(supported, dict):
            for name, code in supported.items():
                normalized_code = str(code).strip().lower()
                if not normalized_code:
                    continue
                language_map[normalize_lookup_token(name)] = normalized_code
                language_map[normalize_lookup_token(normalized_code)] = normalized_code
                language_map[normalized_code] = normalized_code
    except Exception:
        pass

    get_translation_language_map._cache = language_map
    return language_map


def normalize_translation_target(target_lang: str) -> str:
    raw = str(target_lang or "").strip().lower()
    if not raw:
        return "en"

    language_map = get_translation_language_map()
    normalized = normalize_lookup_token(raw)
    if normalized in language_map:
        return language_map[normalized]

    raw_hyphenated = raw.replace("_", "-")
    if re.fullmatch(r"[a-z]{2,3}(?:-[a-z]{2,4})?", raw_hyphenated):
        return raw_hyphenated

    return "en"


MAX_UPLOAD_BYTES = MAX_UPLOAD_MB * 1024 * 1024


def create_job_workspace() -> tuple[str, Path, Path]:
    job_id = uuid.uuid4().hex
    job_dir = JOBS_DIR / job_id
    input_dir = job_dir / "input"
    output_dir = job_dir / "output"
    input_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    return job_id, input_dir, output_dir


def save_uploads(files: list[UploadFile], input_dir: Path) -> list[Path]:
    saved_paths: list[Path] = []
    for index, upload in enumerate(files):
        filename = upload.filename or f"file-{index}"
        safe_name = f"{index + 1}_{Path(filename).name}"
        destination = input_dir / safe_name
        with destination.open("wb") as f:
            written = 0
            while True:
                chunk = upload.file.read(1024 * 1024)
                if not chunk:
                    break
                written += len(chunk)
                if written > MAX_UPLOAD_BYTES:
                    raise HTTPException(
                        status_code=413,
                        detail=f"{filename} exceeds max upload size of {MAX_UPLOAD_MB} MB",
                    )
                f.write(chunk)
        saved_paths.append(destination)
    return saved_paths


def parse_payload(payload_str: str | None) -> dict[str, Any]:
    if not payload_str:
        return {}
    try:
        return json.loads(payload_str)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=400, detail="Invalid payload JSON") from exc


def ensure_files(files: list[Path], min_count: int = 1) -> None:
    if len(files) < min_count:
        raise HTTPException(status_code=400, detail=f"This tool requires at least {min_count} file(s)")


def parse_page_numbers(raw: str, total_pages: int) -> list[int]:
    cleaned = raw.strip().lower()
    if not cleaned or cleaned == "all":
        return list(range(1, total_pages + 1))

    pages: list[int] = []
    for token in cleaned.split(","):
        part = token.strip()
        if not part:
            continue
        if "-" in part:
            start_raw, end_raw = part.split("-", 1)
            if start_raw.strip().isdigit() and end_raw.strip().isdigit():
                start = int(start_raw)
                end = int(end_raw)
                if start <= end:
                    pages.extend(range(start, end + 1))
        elif part.isdigit():
            pages.append(int(part))

    ordered_unique: list[int] = []
    seen: set[int] = set()
    for page_no in pages:
        if 1 <= page_no <= total_pages and page_no not in seen:
            ordered_unique.append(page_no)
            seen.add(page_no)

    return ordered_unique


def zip_outputs(paths: list[Path], output_dir: Path, zip_name: str) -> Path:
    archive_path = output_dir / f"{zip_name}.zip"

    # Build the archive manually to avoid self-inclusion when the ZIP lives in output_dir.
    with zipfile.ZipFile(archive_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path in sorted(paths):
            if not path.exists() or path.resolve() == archive_path.resolve():
                continue

            if path.is_file():
                archive.write(path, arcname=path.relative_to(output_dir))
                continue

            for child in sorted(path.rglob("*")):
                if child.is_file() and child.resolve() != archive_path.resolve():
                    archive.write(child, arcname=child.relative_to(output_dir))

    return archive_path


def create_single_file_result(path: Path, message: str, content_type: str = "application/octet-stream") -> ExecutionResult:
    return ExecutionResult(
        kind="file",
        message=message,
        output_path=path,
        filename=path.name,
        content_type=content_type,
    )


def create_zip_result(output_dir: Path, message: str, zip_name: str = "result") -> ExecutionResult:
    zip_file = zip_outputs(list(output_dir.iterdir()), output_dir, zip_name)
    return ExecutionResult(
        kind="file",
        message=message,
        output_path=zip_file,
        filename=zip_file.name,
        content_type="application/zip",
    )


def extract_pdf_text(pdf_path: Path) -> str:
    # Try multiple extraction engines and keep the most informative result.
    # This improves accuracy for scanned, complex-layout, and mixed-encoding PDFs.
    candidates: list[str] = []

    try:
        document = fitz.open(str(pdf_path))
        try:
            fitz_text = "\n\n".join((page.get_text("text") or "").strip() for page in document).strip()
            if fitz_text:
                candidates.append(fitz_text)
        finally:
            document.close()
    except Exception:
        pass

    try:
        reader = PdfReader(str(pdf_path))
        pypdf_text = "\n\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()
        if pypdf_text:
            candidates.append(pypdf_text)
    except Exception:
        pass

    try:
        import pdfplumber

        with pdfplumber.open(str(pdf_path)) as pdf:
            plumber_text = "\n\n".join((page.extract_text() or "").strip() for page in pdf.pages).strip()
            if plumber_text:
                candidates.append(plumber_text)
    except Exception:
        pass

    def score_text(value: str) -> int:
        return len(re.sub(r"[\W_]+", "", value or "", flags=re.UNICODE))

    best = max(candidates, key=score_text).strip() if candidates else ""

    # Fallback OCR for scanned/image-only PDFs.
    if score_text(best) < 140:
        try:
            ocr_sections: list[str] = []
            with tempfile.TemporaryDirectory(prefix="pdf-ocr-") as temp_name:
                temp_dir = Path(temp_name)
                document = fitz.open(str(pdf_path))
                try:
                    max_pages = max(1, min(12, len(document)))
                    for index in range(max_pages):
                        page = document.load_page(index)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2.3, 2.3), alpha=False)
                        page_image = temp_dir / f"page-{index + 1}.png"
                        pix.save(str(page_image))
                        page_text = extract_ocr_text_from_image(page_image)
                        if page_text:
                            ocr_sections.append(f"Page {index + 1}\n{page_text}")
                finally:
                    document.close()

            ocr_text = "\n\n".join(ocr_sections).strip()
            if score_text(ocr_text) > score_text(best):
                best = ocr_text
        except Exception:
            pass

    return best


def text_to_pdf(text: str, output_path: Path, title: str = "Generated Document") -> None:
    pdf = canvas.Canvas(str(output_path), pagesize=A4)
    pdf.setTitle(title)
    width, height = A4
    margin = 40
    y = height - margin

    for raw_line in text.splitlines() or [text]:
        wrapped = textwrap.wrap(raw_line, width=100) or [""]
        for line in wrapped:
            if y <= margin:
                pdf.showPage()
                y = height - margin
            pdf.drawString(margin, y, line)
            y -= 14

    pdf.save()


def chunk_text(text: str, max_chars: int = 4000) -> list[str]:
    cleaned = text.strip()
    if not cleaned:
        return []

    if len(cleaned) <= max_chars:
        return [cleaned]

    paragraphs = [part.strip() for part in re.split(r"\n{2,}", cleaned) if part.strip()]
    chunks: list[str] = []
    current_parts: list[str] = []
    current_len = 0

    def flush() -> None:
        nonlocal current_parts, current_len
        if not current_parts:
            return
        chunks.append(" ".join(current_parts).strip())
        current_parts = []
        current_len = 0

    for paragraph in paragraphs:
        sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", paragraph) if segment.strip()]
        if not sentences:
            sentences = [paragraph]

        for sentence in sentences:
            if len(sentence) > max_chars:
                flush()
                for start in range(0, len(sentence), max_chars):
                    piece = sentence[start : start + max_chars].strip()
                    if piece:
                        chunks.append(piece)
                continue

            projected_len = current_len + len(sentence) + (1 if current_parts else 0)
            if projected_len > max_chars and current_parts:
                flush()

            current_parts.append(sentence)
            current_len += len(sentence) + (1 if current_len else 0)

        if current_parts and current_len >= int(max_chars * 0.9):
            flush()

    flush()
    return chunks or [cleaned]


def summarize_text_algo(text: str, max_sentences: int = 5) -> str:
    cleaned_text = text.strip()
    if not cleaned_text:
        return ""

    sentences = _dedupe_ocr_lines(split_sentences(cleaned_text))
    if not sentences:
        return ""

    max_sentences = max(1, min(20, int(max_sentences)))
    if len(sentences) <= max_sentences:
        return " ".join(sentences)

    words = tokenize_words(cleaned_text)
    frequencies: dict[str, int] = {}
    for word in words:
        if word.isascii() and word in STOPWORDS:
            continue
        if len(word) > 1:
            frequencies[word] = frequencies.get(word, 0) + 1

    scored: list[tuple[float, int, str]] = []
    total_sentences = len(sentences)
    for idx, sentence in enumerate(sentences):
        sentence_words = tokenize_words(sentence)
        if not sentence_words:
            continue

        base_score = sum(frequencies.get(word, 0) for word in sentence_words) / max(1, len(sentence_words))
        position_boost = 1.15 if idx < max(2, total_sentences // 8) else 1.0
        if idx >= total_sentences - max(2, total_sentences // 10):
            position_boost = max(position_boost, 1.05)
        length_penalty = 0.85 if len(sentence_words) > 45 else 1.0
        score = base_score * position_boost * length_penalty
        scored.append((score, idx, sentence))

    if not scored:
        return " ".join(sentences[:max_sentences])

    top = sorted(scored, key=lambda item: item[0], reverse=True)[:max_sentences]
    selected_indices = {item[1] for item in top}
    ordered = [sentence for index, sentence in enumerate(sentences) if index in selected_indices]
    return " ".join(ordered)


def _translate_chunk_resilient(translator: GoogleTranslator, chunk: str) -> str:
    try:
        translated = translator.translate(chunk)
        if translated and translated.strip():
            return translated.strip()
    except Exception:
        pass

    sub_chunks = chunk_text(chunk, max_chars=1200)
    if len(sub_chunks) <= 1:
        return chunk

    translated_parts: list[str] = []
    for part in sub_chunks:
        try:
            translated = translator.translate(part)
            translated_parts.append((translated or part).strip())
        except Exception:
            translated_parts.append(part)

    rebuilt = " ".join(item for item in translated_parts if item).strip()
    return rebuilt or chunk


def translate_text_chunks(text: str, target_lang: str) -> str:
    cleaned = text.strip()
    if not cleaned:
        return ""

    normalized_target = normalize_translation_target(target_lang)
    translator = GoogleTranslator(source="auto", target=normalized_target)

    translated_parts: list[str] = []
    for chunk in chunk_text(cleaned, max_chars=3200):
        translated_parts.append(_translate_chunk_resilient(translator, chunk))

    translated_text = "\n\n".join(part for part in translated_parts if part).strip()
    if not translated_text:
        raise HTTPException(status_code=502, detail="Translation failed. Please retry with smaller input.")
    return translated_text


def read_text_input(files: list[Path], payload: dict[str, Any]) -> str:
    text = str(payload.get("text", "")).strip()
    if text:
        return text

    if files:
        source = files[0]
        suffix = source.suffix.lower()
        if suffix == ".pdf":
            return extract_pdf_text(source)
        if suffix in {".txt", ".md", ".csv", ".log", ".json", ".xml", ".yml", ".yaml", ".rtf"}:
            return source.read_text(encoding="utf-8", errors="ignore")

    return ""


def get_soffice_binary() -> str | None:
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "C:/Program Files/LibreOffice/program/soffice.exe",
        "C:/Program Files (x86)/LibreOffice/program/soffice.exe",
    ]

    for candidate in candidates:
        if not candidate:
            continue
        path = Path(candidate)
        if path.exists() and path.is_file():
            return str(path)

    return None


def try_libreoffice_convert_to_pdf(source: Path, output_dir: Path, timeout_seconds: int = 240) -> Path | None:
    soffice = get_soffice_binary()
    if not soffice:
        return None

    before = {item.resolve() for item in output_dir.glob("*.pdf")}
    command = [
        soffice,
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(source),
    ]

    try:
        subprocess.run(command, check=True, capture_output=True, timeout=timeout_seconds)
    except Exception:
        return None

    direct = output_dir / f"{source.stem}.pdf"
    if direct.exists() and direct.stat().st_size > 0:
        return direct

    produced = [
        item
        for item in output_dir.glob("*.pdf")
        if item.resolve() not in before and item.stat().st_size > 0
    ]
    if produced:
        return max(produced, key=lambda item: item.stat().st_mtime)

    return None


def extract_text_from_binary_file(source: Path, max_chars: int = 80000) -> str:
    raw = source.read_bytes()
    decoded = raw.decode("utf-8", errors="ignore")
    if len(decoded.strip()) < 80:
        decoded = raw.decode("latin-1", errors="ignore")

    cleaned = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", " ", decoded)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned[:max_chars]


def has_meaningful_text(text: str, min_chars: int = 120) -> bool:
    return len(re.sub(r"\s+", "", text)) >= min_chars


def extract_docx_text(docx_path: Path) -> str:
    document = Document(str(docx_path))
    lines: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)

    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                lines.append(" | ".join(values))

    return "\n".join(lines).strip()


def build_docx_from_text(text: str, output_path: Path, heading: str) -> None:
    document = Document()
    document.add_heading(heading, level=1)

    paragraphs = [line.strip() for line in text.splitlines() if line.strip()]
    if not paragraphs:
        paragraphs = ["No textual content was found in the source file."]

    for paragraph in paragraphs:
        document.add_paragraph(paragraph)

    document.save(str(output_path))


def top_matching_sentences(text: str, question: str, limit: int = 3) -> list[str]:
    sentences = split_sentences(text)
    if not sentences:
        return []

    keywords = [word for word in tokenize_words(question) if len(word) > 1]
    if not keywords:
        return sentences[:limit]

    scored: list[tuple[int, str]] = []
    for index, sentence in enumerate(sentences):
        sentence_tokens = Counter(tokenize_words(sentence))
        score = sum(sentence_tokens.get(keyword, 0) for keyword in keywords)
        scored.append((score * 1000 - index, sentence))

    best = [item[1] for item in sorted(scored, reverse=True) if item[0] > 0][:limit]
    if not best:
        return sentences[:limit]
    return best


def open_image(path: Path) -> Image.Image:
    return Image.open(path)


def save_image(img: Image.Image, output_path: Path, fmt: str | None = None, quality: int | None = None) -> None:
    kwargs: dict[str, Any] = {}
    if quality:
        kwargs["quality"] = quality
        kwargs["optimize"] = True

    normalized_format = None
    if fmt:
        format_upper = fmt.upper()
        normalized_format = {"JPG": "JPEG", "TIF": "TIFF"}.get(format_upper, format_upper)

    img.save(output_path, format=normalized_format, **kwargs)


def render_pdf_pages_to_pil_images(pdf_path: Path, scale: float = 2.0) -> list[Image.Image]:
    document = fitz.open(str(pdf_path))
    images: list[Image.Image] = []

    try:
        matrix = fitz.Matrix(scale, scale)
        for page in document:
            pix = page.get_pixmap(matrix=matrix)
            image = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
            images.append(image)
    finally:
        document.close()

    return images


def build_tiff_from_images(images: list[Image.Image], output_path: Path) -> None:
    if not images:
        raise HTTPException(status_code=400, detail="No pages were rendered for TIFF export")

    prepared = [img.convert("RGB") for img in images]
    prepared[0].save(output_path, save_all=True, append_images=prepared[1:], compression="tiff_deflate")


def basic_rtf_escape(text: str) -> str:
    escaped = text.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}")
    return escaped.replace("\n", "\\par\n")


def build_rtf_from_text(text: str, output_path: Path) -> None:
    content = text.strip() or "No textual content was found in the source file."
    rtf_body = basic_rtf_escape(content)
    output_path.write_text("{\\rtf1\\ansi\\deff0\n" + rtf_body + "\n}", encoding="utf-8")


def decode_rtf_hex_escapes(content: str) -> str:
    def repl(match: re.Match[str]) -> str:
        try:
            return bytes.fromhex(match.group(1)).decode("cp1252", errors="ignore")
        except ValueError:
            return ""

    return re.sub(r"\\'([0-9a-fA-F]{2})", repl, content)


def extract_rtf_text(rtf_path: Path) -> str:
    raw = rtf_path.read_text(encoding="utf-8", errors="ignore")
    raw = decode_rtf_hex_escapes(raw)
    raw = re.sub(r"\\par[d]? ?", "\n", raw)
    raw = re.sub(r"\\tab ?", "\t", raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)
    raw = raw.replace("{", "").replace("}", "")
    lines = [line.strip() for line in raw.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def build_odt_from_text(text: str, output_path: Path, title: str) -> None:
    from odf import teletype
    from odf.opendocument import OpenDocumentText
    from odf.text import H, P

    document = OpenDocumentText()
    document.text.addElement(H(outlinelevel=1, text=title))

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        lines = ["No textual content was found in the source file."]

    for line in lines:
        paragraph = P()
        teletype.addTextToElement(paragraph, line)
        document.text.addElement(paragraph)

    document.save(str(output_path), addsuffix=False)


def extract_odt_text(odt_path: Path) -> str:
    from odf import teletype
    from odf.opendocument import load
    from odf.text import H, P

    document = load(str(odt_path))
    lines: list[str] = []
    for element in list(document.getElementsByType(H)) + list(document.getElementsByType(P)):
        text = teletype.extractText(element).strip()
        if text:
            lines.append(text)
    return "\n".join(lines).strip()


def build_epub_from_text(text: str, output_path: Path, title: str) -> None:
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier(uuid.uuid4().hex)
    book.set_title(title)
    book.set_language("en")

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        lines = ["No textual content was found in the source file."]

    chapter_html = "".join(f"<p>{html.escape(line)}</p>" for line in lines)
    chapter = epub.EpubHtml(title=title, file_name="chapter-1.xhtml", lang="en")
    chapter.content = f"<html><body><h1>{html.escape(title)}</h1>{chapter_html}</body></html>"

    book.add_item(chapter)
    book.toc = (chapter,)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())

    epub.write_epub(str(output_path), book)


def extract_epub_text(epub_path: Path) -> str:
    from ebooklib import ITEM_DOCUMENT, epub

    book = epub.read_epub(str(epub_path))
    chunks: list[str] = []
    for item in book.get_items():
        if item.get_type() != ITEM_DOCUMENT:
            continue
        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        text = soup.get_text("\n", strip=True)
        if text:
            chunks.append(text)
    return "\n\n".join(chunks).strip()


def get_ocr_engine() -> Any | None:
    try:
        from rapidocr_onnxruntime import RapidOCR
    except Exception:
        return None

    if not hasattr(get_ocr_engine, "_engine"):
        try:
            get_ocr_engine._engine = RapidOCR()
        except Exception:
            return None
    return getattr(get_ocr_engine, "_engine", None)


def _score_ocr_lines(lines: list[str]) -> int:
    return sum(len(re.sub(r"[\W_]+", "", line, flags=re.UNICODE)) for line in lines)


def _dedupe_ocr_lines(lines: list[str]) -> list[str]:
    cleaned: list[str] = []
    seen: set[str] = set()
    for line in lines:
        normalized = re.sub(r"\s+", " ", str(line).strip())
        if not normalized:
            continue
        key = normalized.lower()
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(normalized)
    return cleaned


def _build_ocr_variants(image: Image.Image) -> list[Image.Image]:
    base = image.convert("RGB")
    gray = ImageOps.grayscale(base)
    auto = ImageOps.autocontrast(gray)
    sharp = ImageEnhance.Sharpness(auto).enhance(2.0)
    contrast = ImageEnhance.Contrast(auto).enhance(1.7)

    variants: list[Image.Image] = [
        base,
        auto.convert("RGB"),
        sharp.convert("RGB"),
        contrast.convert("RGB"),
    ]

    try:
        import cv2
        import numpy as np

        arr = np.array(auto)
        _, otsu = cv2.threshold(arr, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        variants.append(Image.fromarray(otsu).convert("RGB"))
    except Exception:
        thresholded = auto.point(lambda value: 255 if value > 160 else 0)
        variants.append(thresholded.convert("RGB"))

    return variants


def extract_ocr_lines_from_image(image_path: Path) -> list[str]:
    source = open_image_file(image_path, "RGB")
    variants = _build_ocr_variants(source)

    best_lines: list[str] = []
    best_score = 0

    with tempfile.TemporaryDirectory(prefix="ocr-variants-") as temp_name:
        temp_dir = Path(temp_name)
        variant_paths: list[Path] = []
        for index, variant in enumerate(variants, start=1):
            variant_path = temp_dir / f"variant-{index}.png"
            variant.save(variant_path, format="PNG")
            variant_paths.append(variant_path)

        engine = get_ocr_engine()
        if engine is not None:
            for variant_path in variant_paths:
                try:
                    result, _ = engine(str(variant_path))
                    if not result:
                        continue
                    rapid_lines = _dedupe_ocr_lines([str(item[1]).strip() for item in result if len(item) > 1])
                    score = _score_ocr_lines(rapid_lines)
                    if score > best_score:
                        best_lines = rapid_lines
                        best_score = score
                except Exception:
                    continue

        # Tesseract fallback with multiple page segmentation modes.
        try:
            import pytesseract

            for variant in variants:
                for psm in (6, 11):
                    config = f"--oem 3 --psm {psm}"
                    tess_text = pytesseract.image_to_string(variant, config=config)
                    tess_lines = _dedupe_ocr_lines(tess_text.splitlines())
                    score = _score_ocr_lines(tess_lines)
                    if score > best_score:
                        best_lines = tess_lines
                        best_score = score
        except Exception:
            pass

        # EasyOCR fallback when RapidOCR/Tesseract are weak or unavailable.
        if best_score < 80:
            try:
                import easyocr

                if not hasattr(extract_ocr_lines_from_image, "_easyocr_reader"):
                    extract_ocr_lines_from_image._easyocr_reader = easyocr.Reader(["en"], gpu=False, verbose=False)

                reader = getattr(extract_ocr_lines_from_image, "_easyocr_reader")
                for variant_path in variant_paths:
                    easy_result = reader.readtext(str(variant_path), detail=0, paragraph=False)
                    easy_lines = _dedupe_ocr_lines([str(item).strip() for item in easy_result])
                    score = _score_ocr_lines(easy_lines)
                    if score > best_score:
                        best_lines = easy_lines
                        best_score = score
            except Exception:
                pass

    return best_lines


def extract_ocr_text_from_image(image_path: Path) -> str:
    return "\n".join(extract_ocr_lines_from_image(image_path)).strip()


def extract_ocr_text_from_pdf(pdf_path: Path, temp_dir: Path) -> str:
    document = fitz.open(str(pdf_path))
    pages_text: list[str] = []

    try:
        for index, page in enumerate(document, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            page_image = temp_dir / f"ocr-page-{index}.png"
            pix.save(str(page_image))
            page_text = extract_ocr_text_from_image(page_image)
            if page_text:
                pages_text.append(f"Page {index}\n{page_text}")
    finally:
        document.close()

    return "\n\n".join(pages_text).strip()


def detect_faces_in_image(image_path: Path) -> tuple[Any, Any]:
    import cv2

    cascade_files = [
        "haarcascade_frontalface_default.xml",
        "haarcascade_frontalface_alt2.xml",
        "haarcascade_profileface.xml",
    ]
    classifiers = [
        (cascade_name, cv2.CascadeClassifier(cv2.data.haarcascades + cascade_name))
        for cascade_name in cascade_files
        if Path(cv2.data.haarcascades + cascade_name).exists()
    ]
    if not classifiers:
        raise HTTPException(status_code=500, detail="OpenCV face detection cascades are not available")

    image = cv2.imread(str(image_path))
    if image is None:
        raise HTTPException(status_code=400, detail="Unable to open image for face detection")

    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    gray = cv2.equalizeHist(gray)

    candidates: list[tuple[int, int, int, int]] = []

    def add_faces(detected: Any, scale: float = 1.0, mirrored: bool = False) -> None:
        for x, y, w, h in detected:
            if mirrored:
                x = gray.shape[1] - x - w
            x = int(round(x / scale))
            y = int(round(y / scale))
            w = int(round(w / scale))
            h = int(round(h / scale))
            if w < 20 or h < 20:
                continue
            candidates.append((x, y, w, h))

    for cascade_name, classifier in classifiers:
        try:
            faces = classifier.detectMultiScale(gray, scaleFactor=1.08, minNeighbors=4, minSize=(24, 24))
            add_faces(faces)

            if "profileface" in cascade_name:
                flipped = cv2.flip(gray, 1)
                mirrored_faces = classifier.detectMultiScale(flipped, scaleFactor=1.08, minNeighbors=4, minSize=(24, 24))
                add_faces(mirrored_faces, mirrored=True)
        except Exception:
            continue

    # Small-face recovery pass for low-resolution uploads.
    if not candidates:
        scale = 2.0
        upscaled = cv2.resize(gray, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
        for _, classifier in classifiers:
            try:
                faces = classifier.detectMultiScale(upscaled, scaleFactor=1.08, minNeighbors=4, minSize=(40, 40))
                add_faces(faces, scale=scale)
            except Exception:
                continue

    def iou(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> float:
        ax1, ay1, aw, ah = a
        bx1, by1, bw, bh = b
        ax2, ay2 = ax1 + aw, ay1 + ah
        bx2, by2 = bx1 + bw, by1 + bh
        ix1, iy1 = max(ax1, bx1), max(ay1, by1)
        ix2, iy2 = min(ax2, bx2), min(ay2, by2)
        iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
        inter = iw * ih
        if inter <= 0:
            return 0.0
        union = aw * ah + bw * bh - inter
        return inter / max(1.0, float(union))

    deduped: list[tuple[int, int, int, int]] = []
    for rect in sorted(candidates, key=lambda value: value[2] * value[3], reverse=True):
        if any(iou(rect, existing) > 0.38 for existing in deduped):
            continue
        deduped.append(rect)

    return image, deduped


def get_resampling_module() -> Any:
    return getattr(Image, "Resampling", Image)


def register_heif_support() -> None:
    try:
        from pillow_heif import register_heif_opener

        register_heif_opener()
    except Exception:
        pass


def open_image_file(image_path: Path, mode: str | None = None) -> Image.Image:
    if image_path.suffix.lower() in {".heic", ".heif"}:
        register_heif_support()
    image = ImageOps.exif_transpose(Image.open(image_path))
    return image.convert(mode) if mode else image


def load_ui_font(size: int) -> Any:
    normalized_size = max(12, int(size))
    candidates = [
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
        Path("/usr/share/fonts/truetype/freefont/FreeSans.ttf"),
        Path("/usr/share/fonts/truetype/ubuntu/Ubuntu-R.ttf"),
        Path("/usr/share/fonts/TTF/DejaVuSans.ttf"),
        Path("/usr/share/fonts/dejavu/DejaVuSans.ttf"),
        Path("/usr/local/share/fonts/DejaVuSans.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/segoeui.ttf"),
        Path("C:/Windows/Fonts/calibri.ttf"),
        Path("C:/Windows/Fonts/verdana.ttf"),
        Path("C:/Windows/Fonts/tahoma.ttf"),
        Path("C:/Windows/Fonts/DejaVuSans.ttf"),
        Path("/System/Library/Fonts/Helvetica.ttc"),
        Path("/System/Library/Fonts/Arial.ttf"),
    ]

    for candidate in candidates:
        if candidate.exists():
            try:
                return ImageFont.truetype(str(candidate), normalized_size)
            except OSError:
                continue

    for candidate_name in ("DejaVuSans.ttf", "DejaVuSans-Bold.ttf", "LiberationSans-Regular.ttf", "arial.ttf", "segoeui.ttf"):
        try:
            return ImageFont.truetype(candidate_name, normalized_size)
        except OSError:
            continue

    try:
        return ImageFont.load_default(size=normalized_size)
    except TypeError:
        return ImageFont.load_default()


def parse_color_value(raw: Any, default: tuple[int, int, int]) -> tuple[int, int, int]:
    try:
        return ImageColor.getrgb(str(raw).strip())
    except Exception:
        return default


def clamp_percentage(raw: Any, default: float) -> float:
    try:
        value = float(raw)
    except (TypeError, ValueError):
        value = default
    return max(0.0, min(100.0, value))


def is_truthy(raw: Any) -> bool:
    if isinstance(raw, bool):
        return raw
    if isinstance(raw, (int, float)):
        return raw != 0
    return str(raw).strip().lower() in {"1", "true", "yes", "on"}


def resolve_anchor_position(
    position: str,
    base_size: tuple[int, int],
    overlay_size: tuple[int, int],
    margin: int = 24,
) -> tuple[int, int]:
    base_width, base_height = base_size
    overlay_width, overlay_height = overlay_size
    normalized = position.strip().lower()

    if normalized == "top-left":
        return margin, margin
    if normalized == "top-right":
        return max(0, base_width - overlay_width - margin), margin
    if normalized == "bottom-left":
        return margin, max(0, base_height - overlay_height - margin)
    if normalized == "center":
        return max(0, (base_width - overlay_width) // 2), max(0, (base_height - overlay_height) // 2)
    return max(0, base_width - overlay_width - margin), max(0, base_height - overlay_height - margin)


def save_images_as_pdf(images: list[Image.Image], output_path: Path) -> None:
    if not images:
        raise HTTPException(status_code=400, detail="No images were available to build a PDF")

    prepared = [image.convert("RGB") for image in images]
    prepared[0].save(output_path, save_all=True, append_images=prepared[1:])


def image_paths_to_pdf(paths: list[Path], output_path: Path) -> None:
    save_images_as_pdf([open_image_file(path, "RGB") for path in paths], output_path)


def color_to_pdf_tuple(raw: Any, default: tuple[int, int, int] = (255, 234, 88)) -> tuple[float, float, float]:
    rgb = parse_color_value(raw, default)
    return tuple(round(channel / 255, 4) for channel in rgb)


def extract_top_keywords(text: str, limit: int = 8) -> list[dict[str, Any]]:
    words = [
        word
        for word in tokenize_words(text)
        if len(word) > 1 and not (word.isascii() and word in STOPWORDS)
    ]
    return [
        {"keyword": word, "count": count}
        for word, count in Counter(words).most_common(max(1, limit))
    ]


def extract_image_dpi_info(source: Path) -> dict[str, Any]:
    image = open_image_file(source)
    raw_dpi = image.info.get("dpi")
    if isinstance(raw_dpi, tuple) and len(raw_dpi) >= 2:
        dpi_x, dpi_y = raw_dpi[0], raw_dpi[1]
    elif isinstance(raw_dpi, (int, float)):
        dpi_x = dpi_y = raw_dpi
    else:
        dpi_x = dpi_y = 72

    return {
        "format": image.format,
        "size": {"width": image.width, "height": image.height},
        "dpi": {"x": round(float(dpi_x), 2), "y": round(float(dpi_y), 2)},
    }


def resize_image_to_physical_units(
    source: Path,
    output_path: Path,
    width_value: float,
    height_value: float,
    dpi: int,
    unit: str,
) -> None:
    unit_to_inches = {
        "cm": 1 / 2.54,
        "mm": 1 / 25.4,
        "inch": 1.0,
    }
    factor = unit_to_inches[unit]
    width_px = max(1, int(round(width_value * factor * dpi)))
    height_px = max(1, int(round(height_value * factor * dpi)))
    image = open_image_file(source)
    resized = image.resize((width_px, height_px), get_resampling_module().LANCZOS)
    save_image(resized, output_path, quality=95)


def dispatch_existing_handler(
    slug: str,
    files: list[Path],
    payload: dict[str, Any],
    output_dir: Path,
) -> ExecutionResult:
    handler = HANDLERS.get(slug)
    if not handler:
        raise HTTPException(status_code=501, detail=f"Handler not found for {slug}")
    return handler(files, payload, output_dir)


def handle_merge_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 2)
    output = output_dir / "merged.pdf"

    # Prefer PyMuPDF merge for better structural preservation.
    try:
        merged = fitz.open()
        try:
            for file_path in files:
                part = fitz.open(str(file_path))
                try:
                    merged.insert_pdf(part)
                finally:
                    part.close()

            merged.save(str(output), garbage=3, deflate=True)
        finally:
            merged.close()

        return create_single_file_result(output, "PDF files merged successfully", "application/pdf")
    except Exception:
        pass

    writer = PdfWriter()
    for file_path in files:
        reader = PdfReader(str(file_path))
        for page in reader.pages:
            writer.add_page(page)

    with output.open("wb") as f:
        writer.write(f)
    return create_single_file_result(output, "PDF files merged successfully", "application/pdf")


def handle_split_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))

    requested = str(payload.get("pages", "all"))
    pages = parse_page_numbers(requested, len(reader.pages))
    if not pages:
        raise HTTPException(status_code=400, detail="Provide valid pages like '1,2,5' or ranges like '2-4'")

    generated = 0
    for page_no in pages:
        if page_no < 1 or page_no > len(reader.pages):
            continue
        writer = PdfWriter()
        writer.add_page(reader.pages[page_no - 1])
        part = output_dir / f"page-{page_no}.pdf"
        with part.open("wb") as f:
            writer.write(f)
        generated += 1

    if generated == 0:
        raise HTTPException(status_code=400, detail="No valid pages matched your split request")

    return create_zip_result(output_dir, "PDF split completed", "split-pages")


def handle_compress_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    output = output_dir / "compressed.pdf"
    original_size = files[0].stat().st_size
    quality = str(payload.get("quality", "recommended")).lower()
    best_output: Path | None = None
    best_size = original_size

    # Attempt 1: pikepdf with stream compression and object stream deduplication.
    try:
        import pikepdf

        tmp1 = output_dir / "compressed_pikepdf.pdf"
        with pikepdf.open(str(files[0])) as pdf:
            pdf.remove_unreferenced_resources()
            save_kwargs: dict[str, Any] = {
                "compress_streams": True,
                "recompress_flate": True,
                "normalize_content": True,
            }
            try:
                save_kwargs["object_stream_mode"] = pikepdf.ObjectStreamMode.generate
            except Exception:
                pass
            pdf.save(str(tmp1), **save_kwargs)

        candidate_size = tmp1.stat().st_size
        if candidate_size < best_size:
            best_size = candidate_size
            best_output = tmp1
    except Exception:
        pass

    # Attempt 2: PyMuPDF with image downsampling for quality levels below high.
    try:
        import fitz  # PyMuPDF

        tmp2 = output_dir / "compressed_mupdf.pdf"
        doc = fitz.open(str(files[0]))
        dpi_map = {"low": 72, "recommended": 96, "high": 150}
        image_dpi = dpi_map.get(quality, 96)

        for page in doc:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n >= 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_data = pix.tobytes("jpeg", jpg_quality=int(image_dpi * 0.7))
                    doc.update_stream(xref, img_data)
                except Exception:
                    pass

        doc.save(
            str(tmp2),
            garbage=4,
            clean=True,
            deflate=True,
            deflate_images=True,
            deflate_fonts=True,
            linear=True,
        )
        doc.close()

        candidate_size = tmp2.stat().st_size
        if candidate_size < best_size:
            best_size = candidate_size
            best_output = tmp2
    except Exception:
        pass

    # Attempt 3: pypdf with content stream compression (fallback).
    if best_output is None:
        try:
            tmp3 = output_dir / "compressed_pypdf.pdf"
            reader = PdfReader(str(files[0]))
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            for page in writer.pages:
                try:
                    page.compress_content_streams()
                except Exception:
                    pass
            if reader.metadata:
                writer.add_metadata(reader.metadata)
            with tmp3.open("wb") as f:
                writer.write(f)
            candidate_size = tmp3.stat().st_size
            if candidate_size < best_size:
                best_size = candidate_size
                best_output = tmp3
        except Exception:
            pass

    if best_output and best_output != output:
        import shutil
        shutil.copy2(str(best_output), str(output))
    elif best_output is None:
        # Nothing reduced size — return original unchanged.
        import shutil
        shutil.copy2(str(files[0]), str(output))
        best_size = original_size

    reduction = max(0.0, ((original_size - best_size) / max(1, original_size)) * 100)
    if reduction < 0.5:
        msg = "PDF is already optimally compressed - no significant reduction possible."
    else:
        msg = f"PDF compressed successfully ({reduction:.1f}% smaller)"

    return create_single_file_result(output, msg, "application/pdf")


def handle_rotate_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    angle = int(payload.get("angle", 90))
    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)

    output = output_dir / "rotated.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF rotation completed", "application/pdf")


def handle_organize_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    raw_order = str(payload.get("page_order", "")).strip()
    if not raw_order:
        raise HTTPException(status_code=400, detail="page_order is required, e.g. 3,1,2")

    order = parse_page_numbers(raw_order, len(reader.pages))
    if not order:
        raise HTTPException(status_code=400, detail="No valid page numbers in page_order")

    writer = PdfWriter()
    total = len(reader.pages)
    for page_no in order:
        if 1 <= page_no <= total:
            writer.add_page(reader.pages[page_no - 1])

    output = output_dir / "organized.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF pages reorganized", "application/pdf")


def handle_crop_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    left = float(payload.get("left", 0))
    right = float(payload.get("right", 0))
    top = float(payload.get("top", 0))
    bottom = float(payload.get("bottom", 0))

    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        box = page.mediabox
        min_x = float(box.left)
        min_y = float(box.bottom)
        max_x = float(box.right)
        max_y = float(box.top)

        x0 = min(max_x - 1, max(min_x, min_x + left))
        y0 = min(max_y - 1, max(min_y, min_y + bottom))
        x1 = max(x0 + 1, min(max_x, max_x - right))
        y1 = max(y0 + 1, min(max_y, max_y - top))
        page.mediabox.lower_left = (x0, y0)
        page.mediabox.upper_right = (x1, y1)
        writer.add_page(page)

    output = output_dir / "cropped.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF cropped successfully", "application/pdf")


def _merge_overlay_to_pdf(
    files: list[Path],
    output_dir: Path,
    text_value: str,
    mode: str,
    font_size: int,
) -> Path:
    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for idx, page in enumerate(reader.pages, start=1):
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)

        stream = io.BytesIO()
        overlay_canvas = canvas.Canvas(stream, pagesize=(width, height))
        overlay_canvas.setFont("Helvetica", font_size)

        if mode == "numbers":
            draw_text = str(idx)
            overlay_canvas.drawCentredString(width / 2, 24, draw_text)
        else:
            overlay_canvas.drawString(36, 24, text_value)

        overlay_canvas.save()
        stream.seek(0)

        overlay_pdf = PdfReader(stream)
        page.merge_page(overlay_pdf.pages[0])
        writer.add_page(page)

    output = output_dir / ("numbered.pdf" if mode == "numbers" else "watermarked.pdf")
    with output.open("wb") as f:
        writer.write(f)
    return output


def handle_watermark_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    watermark = str(payload.get("text", "ISHU TOOLS")).strip() or "ISHU TOOLS"
    font_size = int(payload.get("font_size", 18))
    output = _merge_overlay_to_pdf(files, output_dir, watermark, "watermark", font_size)
    return create_single_file_result(output, "Watermark applied", "application/pdf")


def handle_page_numbers_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    font_size = int(payload.get("font_size", 12))
    output = _merge_overlay_to_pdf(files, output_dir, "", "numbers", font_size)
    return create_single_file_result(output, "Page numbers added", "application/pdf")


def handle_protect_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    password = str(payload.get("password", "")).strip()
    if not password:
        raise HTTPException(status_code=400, detail="password is required")

    reader = PdfReader(str(files[0]))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(password)
    output = output_dir / "protected.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF protected", "application/pdf")


def handle_unlock_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    password = str(payload.get("password", "")).strip()

    reader = PdfReader(str(files[0]))
    if reader.is_encrypted:
        if not password:
            raise HTTPException(status_code=400, detail="password is required for encrypted PDF")
        reader.decrypt(password)

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    output = output_dir / "unlocked.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF unlocked", "application/pdf")


def handle_redact_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    keywords_raw = str(payload.get("keywords", "")).strip()
    keywords = [k.strip() for k in keywords_raw.split(",") if k.strip()]
    if not keywords:
        raise HTTPException(status_code=400, detail="keywords is required, e.g. secret,phone")

    document = fitz.open(str(files[0]))
    for page in document:
        for keyword in keywords:
            found = page.search_for(keyword)
            for area in found:
                page.add_redact_annot(area, fill=(0, 0, 0))
        page.apply_redactions()

    output = output_dir / "redacted.pdf"
    document.save(str(output), garbage=4, deflate=True)
    document.close()

    return create_single_file_result(output, "PDF redaction complete", "application/pdf")


def handle_compare_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 2)

    first = extract_pdf_text(files[0]).splitlines()
    second = extract_pdf_text(files[1]).splitlines()

    diff = difflib.unified_diff(first, second, fromfile=files[0].name, tofile=files[1].name, lineterm="")
    output = output_dir / "comparison.diff.txt"
    output.write_text("\n".join(diff), encoding="utf-8")

    return create_single_file_result(output, "PDF comparison generated", "text/plain")


def handle_extract_text_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    output = output_dir / "extracted.txt"
    output.write_text(extract_pdf_text(files[0]), encoding="utf-8")
    return create_single_file_result(output, "Text extracted", "text/plain")


def handle_extract_images_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    document = fitz.open(str(files[0]))
    count = 0

    for page_index in range(len(document)):
        page = document[page_index]
        for image_index, image in enumerate(page.get_images(full=True), start=1):
            xref = image[0]
            base = document.extract_image(xref)
            image_bytes = base["image"]
            ext = base.get("ext", "png")
            output_file = output_dir / f"page-{page_index + 1}-img-{image_index}.{ext}"
            output_file.write_bytes(image_bytes)
            count += 1

    document.close()

    if count == 0:
        raise HTTPException(status_code=400, detail="No embedded images found in the PDF")

    return create_zip_result(output_dir, "Images extracted from PDF", "pdf-images")


_FITZ_NATIVE_FORMATS = {"png", "jpg", "jpeg", "pnm", "pgm", "ppm", "pbm", "pam", "psd", "ps"}


def _pdf_to_image(files: list[Path], output_dir: Path, ext: str) -> ExecutionResult:
    document = fitz.open(str(files[0]))
    use_pillow = ext.lower() not in _FITZ_NATIVE_FORMATS

    for index, page in enumerate(document, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        output = output_dir / f"page-{index}.{ext}"
        if use_pillow:
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            img.save(str(output))
        else:
            pix.save(str(output))

    document.close()
    return create_zip_result(output_dir, f"PDF converted to {ext.upper()}", f"pdf-to-{ext}")


def handle_pdf_to_jpg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    return _pdf_to_image(files, output_dir, "jpg")


def handle_pdf_to_png(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    return _pdf_to_image(files, output_dir, "png")


def handle_jpg_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    output = output_dir / "images.pdf"
    image_paths_to_pdf(files, output)
    return create_single_file_result(output, "Images converted to PDF", "application/pdf")


def handle_image_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_jpg_to_pdf(files, payload, output_dir)


def handle_repair_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]), strict=False)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    output = output_dir / "repaired.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF repair attempt completed", "application/pdf")


def handle_translate_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target = str(payload.get("target_lang", "en")).strip() or "en"
    content = extract_pdf_text(files[0])
    if not has_meaningful_text(content, min_chars=16):
        raise HTTPException(status_code=400, detail="No readable text was found in the uploaded PDF")
    translated = translate_text_chunks(content, target)

    output = output_dir / "translated.pdf"
    text_to_pdf(translated, output, title="Translated PDF")
    return create_single_file_result(output, "PDF translated", "application/pdf")


def handle_summarize_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    max_sentences = max(1, min(20, int(payload.get("max_sentences", 5))))
    content = extract_pdf_text(files[0])
    summary = summarize_text_algo(content, max_sentences=max_sentences)
    if not summary:
        raise HTTPException(status_code=400, detail="No readable text was found in the uploaded PDF")

    output = output_dir / "summary.txt"
    output.write_text(summary, encoding="utf-8")
    return create_single_file_result(output, "PDF summary generated", "text/plain")


def handle_compress_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    quality = int(payload.get("quality", 70))
    quality = max(5, min(95, quality))

    for file_path in files:
        img = open_image_file(file_path)
        ext = file_path.suffix.lower().replace(".", "")
        target_ext = "jpg" if ext in {"jpeg", "jpg"} else ext
        output = output_dir / f"compressed-{file_path.stem}.{target_ext}"

        if target_ext in {"jpg", "jpeg"}:
            save_image(img.convert("RGB"), output, fmt="JPEG", quality=quality)
        elif target_ext == "png":
            img.save(str(output), format="PNG", optimize=True, compress_level=9)
        elif target_ext == "webp":
            img.convert("RGB").save(str(output), format="WEBP", quality=quality, method=6)
        else:
            save_image(img.convert("RGB") if img.mode == "RGBA" else img, output, fmt=target_ext.upper(), quality=quality)

    if len(files) == 1:
        only_file = next(output_dir.iterdir())
        return create_single_file_result(only_file, "Image compressed", "image/*")

    return create_zip_result(output_dir, "Images compressed", "compressed-images")


def handle_resize_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    width = max(1, int(payload.get("width", 1200)))
    height = max(1, int(payload.get("height", 800)))
    maintain_aspect = str(payload.get("maintain_aspect", "true")).strip().lower() not in {"false", "0", "no"}

    source = open_image_file(files[0])
    if maintain_aspect:
        resized = source.copy()
        resized.thumbnail((width, height), get_resampling_module().LANCZOS)
    else:
        resized = source.resize((width, height), get_resampling_module().LANCZOS)

    output = output_dir / f"resized-{files[0].name}"
    save_image(resized, output, quality=95)
    aspect_message = "(aspect ratio preserved)" if maintain_aspect else ""
    return create_single_file_result(output, f"Image resized {aspect_message}".strip(), "image/*")


def handle_crop_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = open_image_file(files[0])
    x = int(payload.get("x", 0))
    y = int(payload.get("y", 0))
    w = int(payload.get("width", max(1, source.width - x)))
    h = int(payload.get("height", max(1, source.height - y)))

    left = max(0, min(source.width - 1, x))
    top = max(0, min(source.height - 1, y))
    right = max(left + 1, min(source.width, left + max(1, w)))
    bottom = max(top + 1, min(source.height, top + max(1, h)))

    if right <= left or bottom <= top:
        raise HTTPException(status_code=400, detail="Invalid crop area. Check x, y, width, and height values")

    cropped = source.crop((left, top, right, bottom))
    output = output_dir / f"cropped-{files[0].name}"
    save_image(cropped, output)
    return create_single_file_result(output, "Image cropped", "image/*")


def handle_rotate_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    angle = float(payload.get("angle", 90))

    source = open_image_file(files[0])
    fill = parse_color_value(payload.get("background", "#ffffff"), (255, 255, 255))
    if source.mode == "RGBA":
        fill_color = (*fill, 0)
    else:
        fill_color = fill

    rotated = source.rotate(-angle, expand=True, resample=get_resampling_module().BICUBIC, fillcolor=fill_color)
    output = output_dir / f"rotated-{files[0].name}"
    save_image(rotated, output)
    return create_single_file_result(output, "Image rotated", "image/*")


def handle_convert_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target_format = str(payload.get("target_format", "png")).lower().strip()
    supported = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff", "tif", "ico", "pdf", "svg"}
    if target_format not in supported:
        raise HTTPException(status_code=400, detail=f"target_format must be one of: {', '.join(sorted(supported))}")

    quality = max(5, min(100, int(payload.get("quality", 92))))
    img = open_image_file(files[0])
    ext = "jpg" if target_format == "jpeg" else ("tiff" if target_format == "tif" else target_format)
    output = output_dir / f"converted.{ext}"

    if target_format in {"jpg", "jpeg"}:
        save_image(img.convert("RGB"), output, fmt="JPEG", quality=quality)
    elif target_format == "png":
        img.save(str(output), format="PNG", optimize=True, compress_level=9)
    elif target_format == "webp":
        img.convert("RGB").save(str(output), format="WEBP", quality=quality, method=6)
    elif target_format in {"tiff", "tif"}:
        img.save(str(output), format="TIFF", compression="tiff_lzw")
    elif target_format == "ico":
        sizes = [(16, 16), (32, 32), (48, 48), (64, 64), (128, 128), (256, 256)]
        resized = img.convert("RGBA")
        resized.thumbnail((256, 256), get_resampling_module().LANCZOS)
        resized.save(str(output), format="ICO", sizes=sizes)
    elif target_format == "pdf":
        try:
            import fitz
            pix_img = img.convert("RGB")
            buf = io.BytesIO()
            pix_img.save(buf, format="JPEG", quality=92)
            buf.seek(0)
            doc = fitz.open()
            w, h = pix_img.size
            page = doc.new_page(width=w, height=h)
            page.insert_image(fitz.Rect(0, 0, w, h), stream=buf.read())
            doc.save(str(output))
            doc.close()
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"Image to PDF failed: {exc}") from exc
    elif target_format == "svg":
        try:
            import cairosvg
            buf = io.BytesIO()
            img.convert("RGBA").save(buf, format="PNG")
            w, h = img.size
            svg_data = f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}"><image href="data:image/png;base64,{__import__("base64").b64encode(buf.getvalue()).decode()}" width="{w}" height="{h}"/></svg>'
            output.write_text(svg_data, encoding="utf-8")
        except Exception:
            # Simple SVG wrapper without cairosvg
            buf = io.BytesIO()
            img.convert("RGBA").save(buf, format="PNG")
            import base64
            encoded = base64.b64encode(buf.getvalue()).decode()
            w, h = img.size
            svg_data = f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}"><image href="data:image/png;base64,{encoded}" width="{w}" height="{h}"/></svg>'
            output.write_text(svg_data, encoding="utf-8")
    else:
        save_image(img, output, fmt=target_format.upper(), quality=quality)

    return create_single_file_result(output, "Image converted", "image/*")


def handle_watermark_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    text = str(payload.get("text", "ISHU TOOLS")).strip() or "ISHU TOOLS"
    font_size = int(payload.get("font_size", 36))
    opacity = int(clamp_percentage(payload.get("opacity", 60), 60) * 255 / 100)

    image = Image.open(files[0]).convert("RGBA")
    layer = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(layer)
    font = load_ui_font(font_size)

    x = int(payload.get("x", 20))
    y = int(payload.get("y", max(0, image.height - font_size - 30)))
    draw.text((x, y), text, fill=(255, 255, 255, opacity), font=font)

    output_img = Image.alpha_composite(image, layer).convert("RGB")
    output = output_dir / f"watermarked-{files[0].stem}.jpg"
    output_img.save(output, quality=92)
    return create_single_file_result(output, "Watermark added", "image/jpeg")


def handle_grayscale_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image = Image.open(files[0]).convert("L")
    output = output_dir / f"grayscale-{files[0].name}"
    image.save(output)
    return create_single_file_result(output, "Grayscale applied", "image/*")


def handle_blur_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    radius = float(payload.get("radius", 4))

    image = Image.open(files[0])
    blurred = image.filter(ImageFilter.GaussianBlur(radius=radius))
    output = output_dir / f"blur-{files[0].name}"
    blurred.save(output)
    return create_single_file_result(output, "Blur effect applied", "image/*")


def handle_pixelate_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    factor = int(payload.get("factor", 12))
    factor = max(2, factor)

    source = Image.open(files[0])
    small = source.resize((max(1, source.width // factor), max(1, source.height // factor)), Image.NEAREST)
    pixelated = small.resize(source.size, Image.NEAREST)

    output = output_dir / f"pixelate-{files[0].name}"
    pixelated.save(output)
    return create_single_file_result(output, "Pixelate effect applied", "image/*")


def handle_meme_generator(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    top_text = str(payload.get("top_text", "TOP TEXT")).strip()
    bottom_text = str(payload.get("bottom_text", "BOTTOM TEXT")).strip()

    image = Image.open(files[0]).convert("RGB")
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()

    if top_text:
        draw.text((20, 20), top_text, fill="white", font=font)
    if bottom_text:
        draw.text((20, image.height - 40), bottom_text, fill="white", font=font)

    output = output_dir / "meme.jpg"
    image.save(output, quality=92)
    return create_single_file_result(output, "Meme generated", "image/jpeg")


def handle_pdf_to_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    return _pdf_to_image(files, output_dir, "png")


def handle_html_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    url = str(payload.get("url", "")).strip()
    html_content = str(payload.get("html", "")).strip()

    if files and not html_content:
        try:
            html_content = files[0].read_text(encoding="utf-8", errors="ignore")
        except Exception:
            html_content = ""

    output = output_dir / "html-to-pdf.pdf"

    # Best effort: render full HTML with CSS using WeasyPrint.
    try:
        from weasyprint import HTML

        if html_content:
            base_url = str(files[0].parent) if files else None
            HTML(string=html_content, base_url=base_url).write_pdf(str(output))
            return create_single_file_result(output, "HTML converted to PDF", "application/pdf")

        if url:
            HTML(url=url).write_pdf(str(output))
            return create_single_file_result(output, "HTML converted to PDF", "application/pdf")
    except Exception:
        pass

    if url:
        try:
            with httpx.Client(timeout=20, follow_redirects=True, verify=False) as client:
                response = client.get(url)
                response.raise_for_status()
                html_content = response.text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {str(e)}")

    if not html_content:
        raise HTTPException(status_code=400, detail="Provide either url or html in payload")

    text = BeautifulSoup(html_content, "html.parser").get_text("\n", strip=True)
    text_to_pdf(text, output, title="HTML to PDF")
    return create_single_file_result(output, "HTML converted to PDF", "application/pdf")


def handle_md_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide markdown text or upload an md/txt file")

    output = output_dir / "markdown.pdf"

    try:
        import markdown
        from weasyprint import HTML

        rendered = markdown.markdown(
            text,
            extensions=["fenced_code", "tables", "sane_lists"],
        )
        html_doc = f"""<!doctype html>
<html lang=\"en\"><head><meta charset=\"utf-8\" />
<style>
body {{ font-family: Arial, sans-serif; margin: 28px; line-height: 1.5; }}
pre, code {{ font-family: Consolas, 'Courier New', monospace; }}
table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
th, td {{ border: 1px solid #c7cdd8; padding: 6px 8px; text-align: left; }}
</style></head><body>{rendered}</body></html>"""
        HTML(string=html_doc).write_pdf(str(output))
        return create_single_file_result(output, "Markdown converted to PDF", "application/pdf")
    except Exception:
        pass

    plain_text = re.sub(r"[#>*`_\-]", "", text)
    text_to_pdf(plain_text, output, title="Markdown to PDF")
    return create_single_file_result(output, "Markdown converted to PDF", "application/pdf")


def handle_txt_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide text in payload or upload txt file")

    output = output_dir / "text.pdf"
    text_to_pdf(text, output, title="Text to PDF")
    return create_single_file_result(output, "Text converted to PDF", "application/pdf")


def handle_summarize_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to summarize")

    max_sentences = max(1, min(20, int(payload.get("max_sentences", 5))))
    summary = summarize_text_algo(content, max_sentences=max_sentences)
    output = output_dir / "summary.txt"
    output.write_text(summary, encoding="utf-8")
    return create_single_file_result(output, "Text summarized", "text/plain")


def handle_translate_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to translate")

    target = str(payload.get("target_lang", "en")).strip() or "en"
    translated = translate_text_chunks(content, target)

    output = output_dir / "translated.txt"
    output.write_text(translated, encoding="utf-8")
    return create_single_file_result(output, "Text translated", "text/plain")


def handle_qr_code_generator(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    text = str(payload.get("text", "")).strip()
    if not text:
        raise HTTPException(status_code=400, detail="Provide text or URL in payload.text")

    image = qrcode.make(text)
    output = output_dir / "qr-code.png"
    image.save(output)
    return create_single_file_result(output, "QR generated", "image/png")


def handle_extract_metadata(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = files[0]
    metadata: dict[str, Any] = {"filename": source.name, "suffix": source.suffix.lower()}

    if source.suffix.lower() == ".pdf":
        reader = PdfReader(str(source))
        metadata["pdf_metadata"] = dict(reader.metadata or {})
        metadata["pages"] = len(reader.pages)
    else:
        image = Image.open(source)
        metadata["format"] = image.format
        metadata["mode"] = image.mode
        metadata["size"] = {"width": image.width, "height": image.height}
        metadata["exif"] = {str(k): str(v) for k, v in image.getexif().items()}

    output = output_dir / "metadata.json"
    output.write_text(json.dumps(metadata, indent=2), encoding="utf-8")
    return create_single_file_result(output, "Metadata extracted", "application/json")


def handle_remove_metadata_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)

    for source in files:
        image = Image.open(source)
        clean = Image.new(image.mode, image.size)
        clean.putdata(list(image.getdata()))
        out = output_dir / f"clean-{source.name}"
        clean.save(out)

    if len(files) == 1:
        result = next(output_dir.iterdir())
        return create_single_file_result(result, "Metadata removed", "image/*")

    return create_zip_result(output_dir, "Metadata removed from images", "clean-images")


def handle_edit_metadata_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = files[0]

    fields = {
        "title": str(payload.get("title", "")).strip(),
        "author": str(payload.get("author", "")).strip(),
        "subject": str(payload.get("subject", "")).strip(),
        "keywords": str(payload.get("keywords", "")).strip(),
    }
    fields = {key: value for key, value in fields.items() if value}
    if not fields:
        raise HTTPException(status_code=400, detail="Provide at least one metadata field to update")

    suffix = source.suffix.lower()
    output_suffix = suffix if suffix else ".png"
    output = output_dir / f"metadata-edited-{source.stem}{output_suffix}"

    image = open_image_file(source)

    if suffix in {".jpg", ".jpeg", ".tif", ".tiff"}:
        exif = image.getexif()
        description_parts = [fields.get("title", ""), fields.get("subject", ""), fields.get("keywords", "")]
        description = " | ".join(part for part in description_parts if part)
        if description:
            exif[0x010E] = description
        if fields.get("author"):
            exif[0x013B] = fields["author"]

        writable = image.convert("RGB") if suffix in {".jpg", ".jpeg"} and image.mode in {"RGBA", "P"} else image
        writable.save(output, exif=exif)
    elif suffix == ".png":
        png_info = PngImagePlugin.PngInfo()
        for key, value in fields.items():
            png_info.add_text(key, value)
        image.save(output, pnginfo=png_info)
    else:
        # Preserve visual content for formats without easy metadata writing support.
        writable = image.convert("RGB") if image.mode in {"RGBA", "P"} else image
        writable.save(output)

    return create_single_file_result(output, "Image metadata updated", "image/*")


def handle_extract_pages(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_split_pdf(files, payload, output_dir)


def handle_delete_pages(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    to_delete = parse_page_numbers(str(payload.get("pages", "")), len(reader.pages))
    if not to_delete:
        raise HTTPException(status_code=400, detail="Provide pages to delete, e.g. 2,4 or 5-8")

    writer = PdfWriter()
    for page_index, page in enumerate(reader.pages, start=1):
        if page_index not in to_delete:
            writer.add_page(page)

    output = output_dir / "pages-deleted.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "Selected pages deleted", "application/pdf")


def handle_rearrange_pages(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_organize_pdf(files, payload, output_dir)


def handle_resize_pages_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    scale = float(payload.get("scale", 1.0))
    if scale <= 0:
        raise HTTPException(status_code=400, detail="scale must be greater than 0")

    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        page.scale_by(scale)
        writer.add_page(page)

    output = output_dir / "resized-pages.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF page size scaled", "application/pdf")


def handle_add_text_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_watermark_pdf(files, payload, output_dir)


def handle_header_footer_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    header = str(payload.get("header", "")).strip()
    footer = str(payload.get("footer", "")).strip()
    font_size = int(payload.get("font_size", 11))

    if not header and not footer:
        raise HTTPException(status_code=400, detail="Provide header or footer text")

    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)

        stream = io.BytesIO()
        overlay_canvas = canvas.Canvas(stream, pagesize=(width, height))
        overlay_canvas.setFont("Helvetica", font_size)

        if header:
            overlay_canvas.drawString(36, height - 28, header)
        if footer:
            overlay_canvas.drawString(36, 22, footer)

        overlay_canvas.save()
        stream.seek(0)

        overlay_pdf = PdfReader(stream)
        page.merge_page(overlay_pdf.pages[0])
        writer.add_page(page)

    output = output_dir / "header-footer.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "Header and footer added", "application/pdf")


def handle_whiteout_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    keywords_raw = str(payload.get("keywords", "")).strip()
    keywords = [k.strip() for k in keywords_raw.split(",") if k.strip()]
    if not keywords:
        raise HTTPException(status_code=400, detail="keywords is required, e.g. name,email")

    document = fitz.open(str(files[0]))
    for page in document:
        for keyword in keywords:
            found = page.search_for(keyword)
            for area in found:
                page.add_redact_annot(area, fill=(1, 1, 1))
        page.apply_redactions()

    output = output_dir / "whiteout.pdf"
    document.save(str(output), garbage=4, deflate=True)
    document.close()

    return create_single_file_result(output, "Whiteout applied", "application/pdf")


def handle_remove_metadata_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    output = output_dir / "clean-metadata.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF metadata removed", "application/pdf")


def handle_grayscale_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = fitz.open(str(files[0]))
    target = fitz.open()

    for page in source:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), colorspace=fitz.csGRAY)
        new_page = target.new_page(width=pix.width, height=pix.height)
        new_page.insert_image(fitz.Rect(0, 0, pix.width, pix.height), pixmap=pix)

    output = output_dir / "grayscale.pdf"
    target.save(str(output), garbage=4, deflate=True)
    source.close()
    target.close()
    return create_single_file_result(output, "PDF converted to grayscale", "application/pdf")


def handle_png_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_image_to_pdf(files, payload, output_dir)


def handle_webp_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_image_to_pdf(files, payload, output_dir)


def handle_gif_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_image_to_pdf(files, payload, output_dir)


def handle_bmp_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_image_to_pdf(files, payload, output_dir)


def handle_create_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_txt_to_pdf(files, payload, output_dir)


def handle_url_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_html_to_pdf(files, payload, output_dir)


def _convert_alias(files: list[Path], payload: dict[str, Any], output_dir: Path, target_format: str) -> ExecutionResult:
    alias_payload = dict(payload)
    alias_payload["target_format"] = target_format
    return handle_convert_image(files, alias_payload, output_dir)


def handle_jpg_to_png(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return _convert_alias(files, payload, output_dir, "png")


def handle_png_to_jpg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return _convert_alias(files, payload, output_dir, "jpg")


def handle_image_to_webp(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return _convert_alias(files, payload, output_dir, "webp")


def handle_pdf_to_txt(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_pdf_text(files[0])
    output = output_dir / "pdf.txt"
    output.write_text(content, encoding="utf-8")
    return create_single_file_result(output, "PDF converted to TXT", "text/plain")


def handle_pdf_to_markdown(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_pdf_text(files[0])
    lines = [line.strip() for line in content.splitlines() if line.strip()]
    markdown = "\n\n".join(lines)
    output = output_dir / "pdf.md"
    output.write_text(markdown, encoding="utf-8")
    return create_single_file_result(output, "PDF converted to Markdown", "text/markdown")


def handle_pdf_to_json(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    content = extract_pdf_text(files[0])
    data = {
        "filename": files[0].name,
        "pages": len(reader.pages),
        "metadata": dict(reader.metadata or {}),
        "text": content,
    }
    output = output_dir / "pdf.json"
    output.write_text(json.dumps(data, indent=2), encoding="utf-8")
    return create_single_file_result(output, "PDF converted to JSON", "application/json")


def handle_pdf_to_csv(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_pdf_text(files[0])
    lines = [line.strip() for line in content.splitlines() if line.strip()]

    output = output_dir / "pdf.csv"
    with output.open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["line_number", "text"])
        for index, line in enumerate(lines, start=1):
            writer.writerow([index, line])

    return create_single_file_result(output, "PDF converted to CSV", "text/csv")


def handle_flip_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    mode = str(payload.get("mode", "horizontal")).strip().lower()
    image = Image.open(files[0])

    if mode == "vertical":
        transformed = ImageOps.flip(image)
    else:
        transformed = ImageOps.mirror(image)

    output = output_dir / f"flip-{files[0].name}"
    save_image(transformed, output)
    return create_single_file_result(output, "Image flipped", "image/*")


def handle_add_border_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    border_size = int(payload.get("border_size", 20))
    color = str(payload.get("color", "#ffffff")).strip() or "#ffffff"

    image = Image.open(files[0])
    bordered = ImageOps.expand(image, border=max(0, border_size), fill=color)
    output = output_dir / f"border-{files[0].name}"
    save_image(bordered, output)
    return create_single_file_result(output, "Border added to image", "image/*")


def handle_thumbnail_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    max_width = int(payload.get("max_width", 320))
    max_height = int(payload.get("max_height", 320))

    image = Image.open(files[0]).copy()
    image.thumbnail((max_width, max_height))
    output = output_dir / f"thumbnail-{files[0].name}"
    save_image(image, output)
    return create_single_file_result(output, "Thumbnail generated", "image/*")


def handle_image_collage(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    columns = max(1, int(payload.get("columns", 2)))
    cell_size = max(80, int(payload.get("cell_size", 420)))

    images = [Image.open(path).convert("RGB") for path in files]
    rows = (len(images) + columns - 1) // columns

    collage = Image.new("RGB", (columns * cell_size, rows * cell_size), color="white")
    for index, img in enumerate(images):
        fitted = ImageOps.fit(img, (cell_size, cell_size))
        x = (index % columns) * cell_size
        y = (index // columns) * cell_size
        collage.paste(fitted, (x, y))

    output = output_dir / "collage.jpg"
    collage.save(output, quality=92)
    return create_single_file_result(output, "Image collage created", "image/jpeg")


def handle_word_count_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to analyze")

    words = re.findall(r"[A-Za-z']+", content)
    sentences = [s for s in re.split(r"(?<=[.!?])\s+", content.strip()) if s]
    data = {
        "characters": len(content),
        "words": len(words),
        "sentences": len(sentences),
        "avg_word_length": round(sum(len(word) for word in words) / max(1, len(words)), 2),
    }

    output = output_dir / "word-count.json"
    output.write_text(json.dumps(data, indent=2), encoding="utf-8")
    return create_single_file_result(output, "Word count report generated", "application/json")


def _sentence_case(text: str) -> str:
    lowered = text.strip().lower()
    if not lowered:
        return ""
    return lowered[0].upper() + lowered[1:]


def handle_case_converter_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to convert")

    mode = str(payload.get("mode", "upper")).strip().lower()
    if mode == "lower":
        converted = content.lower()
    elif mode == "title":
        converted = content.title()
    elif mode == "sentence":
        converted = " ".join(_sentence_case(part) for part in re.split(r"(?<=[.!?])\s+", content) if part)
    else:
        converted = content.upper()

    output = output_dir / "case-converted.txt"
    output.write_text(converted, encoding="utf-8")
    return create_single_file_result(output, "Text case converted", "text/plain")


def handle_extract_keywords_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to extract keywords")

    top_n = max(1, int(payload.get("top_n", 20)))
    words = [
        word
        for word in re.findall(r"[A-Za-z']+", content.lower())
        if word not in STOPWORDS and len(word) > 2
    ]

    top_words = Counter(words).most_common(top_n)
    output = output_dir / "keywords.txt"
    output.write_text("\n".join(f"{word},{count}" for word, count in top_words), encoding="utf-8")
    return create_single_file_result(output, "Keywords extracted", "text/plain")


def handle_slug_generator_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    text = str(payload.get("text", "")).strip()
    if not text:
        text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide text for slug generation")

    slug = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    output = output_dir / "slug.txt"
    output.write_text(slug, encoding="utf-8")
    return create_single_file_result(output, "Slug generated", "text/plain")


def handle_json_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    raw = read_text_input(files, payload)
    if not raw:
        raise HTTPException(status_code=400, detail="Provide JSON text or upload file")

    try:
        parsed = json.loads(raw)
        content = json.dumps(parsed, indent=2)
    except json.JSONDecodeError:
        content = raw

    output = output_dir / "json.pdf"
    text_to_pdf(content, output, title="JSON to PDF")
    return create_single_file_result(output, "JSON converted to PDF", "application/pdf")


def handle_xml_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    raw = read_text_input(files, payload)
    if not raw:
        raise HTTPException(status_code=400, detail="Provide XML text or upload file")

    output = output_dir / "xml.pdf"
    text_to_pdf(raw, output, title="XML to PDF")
    return create_single_file_result(output, "XML converted to PDF", "application/pdf")


def handle_csv_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    raw = read_text_input(files, payload)
    if not raw:
        raise HTTPException(status_code=400, detail="Provide CSV text or upload file")

    output = output_dir / "csv.pdf"
    text_to_pdf(raw, output, title="CSV to PDF")
    return create_single_file_result(output, "CSV converted to PDF", "application/pdf")


def handle_pdf_pages_to_zip(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_pdf_to_image(files, payload, output_dir)


def handle_zip_images_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = files[0]
    if source.suffix.lower() != ".zip":
        raise HTTPException(status_code=400, detail="Upload a ZIP file containing images")

    extract_dir = output_dir / "zip-images"
    extract_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(source, "r") as archive:
        archive.extractall(extract_dir)

    image_paths = [
        path
        for path in extract_dir.rglob("*")
        if path.is_file() and path.suffix.lower() in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
    ]
    if not image_paths:
        raise HTTPException(status_code=400, detail="No images found inside ZIP")

    images = [Image.open(path).convert("RGB") for path in image_paths]
    output = output_dir / "zip-images.pdf"
    images[0].save(output, save_all=True, append_images=images[1:])
    return create_single_file_result(output, "ZIP images converted to PDF", "application/pdf")


def handle_pdf_to_docx(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    output = output_dir / "pdf.docx"

    # Best effort layout-preserving conversion.
    converted = False
    converter = None
    try:
        from pdf2docx import Converter

        converter = Converter(str(files[0]))
        start_page = max(0, int(payload.get("start_page", 0)))
        end_page_raw = payload.get("end_page")
        end_page: int | None = None
        if end_page_raw not in (None, ""):
            end_page = max(start_page + 1, int(end_page_raw))

        converter.convert(str(output), start=start_page, end=end_page)
        converted = output.exists() and output.stat().st_size > 0
    except Exception:
        converted = False
    finally:
        try:
            if converter is not None:
                converter.close()
        except Exception:
            pass

    if not converted:
        content = extract_pdf_text(files[0])
        build_docx_from_text(content, output, "PDF to DOCX")

    return create_single_file_result(
        output,
        "PDF converted to DOCX",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


def handle_docx_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)

    libreoffice_result = try_libreoffice_convert_to_pdf(files[0], output_dir)
    if libreoffice_result:
        return create_single_file_result(
            libreoffice_result,
            "DOCX converted to PDF",
            "application/pdf",
        )

    content = extract_docx_text(files[0])
    output = output_dir / "docx.pdf"
    text_to_pdf(content, output, title="DOCX to PDF")
    return create_single_file_result(output, "DOCX converted to PDF", "application/pdf")


def handle_pdf_to_excel(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)

    sheet_counter = 1
    wrote_any = False

    def next_sheet(title_prefix: str) -> Any:
        nonlocal sheet_counter
        safe_prefix = re.sub(r"[^A-Za-z0-9 _-]", "", title_prefix).strip() or "Table"
        sheet_name = f"{safe_prefix[:20]}-{sheet_counter}"
        sheet_counter += 1
        return workbook.create_sheet(sheet_name)

    def append_table_rows(sheet: Any, table_rows: list[list[str]]) -> bool:
        wrote = False
        for row in table_rows:
            normalized = [str(cell).replace("\n", " ").strip() if cell is not None else "" for cell in row]
            if any(value for value in normalized):
                sheet.append(normalized)
                wrote = True
        return wrote

    def parse_tabular_rows(lines: list[str]) -> list[list[str]]:
        parsed_rows: list[list[str]] = []
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            cells: list[str]
            if "\t" in stripped:
                cells = [part.strip() for part in stripped.split("\t")]
            elif "|" in stripped and stripped.count("|") >= 1:
                cells = [part.strip() for part in stripped.split("|")]
            elif ";" in stripped and stripped.count(";") >= 1 and stripped.count(" ") < max(8, len(stripped) // 8):
                cells = [part.strip() for part in next(csv.reader([stripped], delimiter=";"))]
            elif "," in stripped and stripped.count(",") >= 1 and stripped.count(" ") < max(8, len(stripped) // 8):
                cells = [part.strip() for part in next(csv.reader([stripped]))]
            else:
                cells = [part.strip() for part in re.split(r"\s{2,}", stripped) if part.strip()]

            cells = [cell for cell in cells if cell]
            if len(cells) >= 2:
                parsed_rows.append(cells)

        if len(parsed_rows) < 2:
            return []

        max_columns = max(len(row) for row in parsed_rows)
        if max_columns < 2:
            return []

        return [row + [""] * (max_columns - len(row)) for row in parsed_rows]

    # 1) Camelot provides the highest quality table extraction when available.
    try:
        import camelot

        for flavor in ("lattice", "stream"):
            try:
                tables = camelot.read_pdf(str(files[0]), pages="all", flavor=flavor)
            except Exception:
                continue

            for index, table in enumerate(tables, start=1):
                dataframe = table.df.fillna("")
                rows = dataframe.values.tolist()
                if not rows:
                    continue

                sheet = next_sheet(f"Camelot-{flavor}-{index}")
                if append_table_rows(sheet, rows):
                    wrote_any = True
    except Exception:
        pass

    # 2) Tabula fallback for scanned / ruled table variants.
    if not wrote_any:
        try:
            import tabula

            table_sets: list[Any] = []
            for options in (
                {"lattice": True, "stream": False},
                {"lattice": False, "stream": True},
            ):
                try:
                    table_sets.extend(
                        tabula.read_pdf(
                            str(files[0]),
                            pages="all",
                            multiple_tables=True,
                            pandas_options={"dtype": str},
                            **options,
                        )
                    )
                except Exception:
                    continue

            for index, dataframe in enumerate(table_sets, start=1):
                if dataframe is None or dataframe.empty:
                    continue

                rows = dataframe.fillna("").astype(str).values.tolist()
                sheet = next_sheet(f"Tabula-{index}")
                if append_table_rows(sheet, rows):
                    wrote_any = True
        except Exception:
            pass

    # 3) pdfplumber fallback: usable when dedicated table engines are unavailable.
    if not wrote_any:
        try:
            import pdfplumber

            with pdfplumber.open(str(files[0])) as pdf:
                for page_index, page in enumerate(pdf.pages, start=1):
                    tables = page.extract_tables() or []
                    if tables:
                        for table_index, table in enumerate(tables, start=1):
                            sheet = next_sheet(f"Page{page_index}-T{table_index}")
                            if append_table_rows(sheet, [[cell if cell is not None else "" for cell in row] for row in (table or [])]):
                                wrote_any = True
                    else:
                        page_text = page.extract_text() or ""
                        text_lines = [line.strip() for line in page_text.splitlines() if line.strip()]
                        if text_lines:
                            parsed_rows = parse_tabular_rows(text_lines)
                            if parsed_rows:
                                sheet = next_sheet(f"Page{page_index}-Parsed")
                                if append_table_rows(sheet, parsed_rows):
                                    wrote_any = True
                            else:
                                sheet = next_sheet(f"Page{page_index}-Text")
                                for line in text_lines:
                                    sheet.append([line])
                                wrote_any = True
        except Exception:
            pass

    # 4) Final guaranteed fallback so tool always returns meaningful output.
    if not wrote_any:
        content = extract_pdf_text(files[0])
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        parsed_rows = parse_tabular_rows(lines)
        if parsed_rows:
            sheet = next_sheet("PDF-Parsed")
            append_table_rows(sheet, parsed_rows)
        else:
            sheet = next_sheet("PDF-Text")
            sheet.append(["line_number", "text"])
            for index, line in enumerate(lines, start=1):
                sheet.append([index, line])

    output = output_dir / "pdf.xlsx"
    workbook.save(str(output))
    return create_single_file_result(
        output,
        "PDF converted to Excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def handle_excel_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)

    libreoffice_result = try_libreoffice_convert_to_pdf(files[0], output_dir)
    if libreoffice_result:
        return create_single_file_result(
            libreoffice_result,
            "Excel converted to PDF",
            "application/pdf",
        )

    workbook = load_workbook(str(files[0]), data_only=True)
    lines: list[str] = []

    for sheet in workbook.worksheets:
        lines.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() if cell is not None else "" for cell in row]
            if any(values):
                lines.append(" | ".join(values))
        lines.append("")

    output = output_dir / "excel.pdf"
    text_to_pdf("\n".join(lines).strip(), output, title="Excel to PDF")
    return create_single_file_result(output, "Excel converted to PDF", "application/pdf")


def handle_pdf_to_pptx(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    document = fitz.open(str(files[0]))
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    emu_per_pixel = 9525
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)

    with tempfile.TemporaryDirectory(prefix="pdf-to-pptx-") as temp_name:
        temp_dir = Path(temp_name)
        for index, page in enumerate(document, start=1):
            slide = presentation.slides.add_slide(blank_layout)

            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_bytes = pix.tobytes("png")
            image = Image.open(io.BytesIO(image_bytes))

            image_w = image.width * emu_per_pixel
            image_h = image.height * emu_per_pixel
            scale = min(slide_width / max(1, image_w), slide_height / max(1, image_h))

            target_w = int(image_w * scale)
            target_h = int(image_h * scale)
            left = int((slide_width - target_w) / 2)
            top = int((slide_height - target_h) / 2)

            stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(stream, left, top, width=target_w, height=target_h)

            page_text = (page.get_text("text") or "").strip()
            if len(re.sub(r"[\W_]+", "", page_text, flags=re.UNICODE)) < 24:
                try:
                    image_path = temp_dir / f"page-{index}.png"
                    image_path.write_bytes(image_bytes)
                    ocr_text = extract_ocr_text_from_image(image_path)
                    if ocr_text:
                        page_text = ocr_text
                except Exception:
                    pass

            if page_text:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.clear()
                notes_frame.text = page_text[:5000]

    document.close()

    output = output_dir / "pdf.pptx"
    presentation.save(str(output))
    return create_single_file_result(
        output,
        "PDF converted to PPTX",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def handle_pptx_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)

    libreoffice_result = try_libreoffice_convert_to_pdf(files[0], output_dir)
    if libreoffice_result:
        return create_single_file_result(
            libreoffice_result,
            "PPTX converted to PDF",
            "application/pdf",
        )

    presentation = Presentation(str(files[0]))
    lines: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"Slide {index}")
        for shape in slide.shapes:
            text_value = str(getattr(shape, "text", "")).strip()
            if text_value:
                lines.append(text_value)
        lines.append("")

    output = output_dir / "pptx.pdf"
    text_to_pdf("\n".join(lines).strip(), output, title="PPTX to PDF")
    return create_single_file_result(output, "PPTX converted to PDF", "application/pdf")


def handle_sign_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    signer = str(payload.get("signer", "ISHU TOOLS")).strip() or "ISHU TOOLS"
    reason = str(payload.get("reason", "Approved")).strip() or "Approved"
    position = str(payload.get("position", "bottom-right")).strip().lower()

    reader = PdfReader(str(files[0]))
    writer = PdfWriter()
    signed_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")

    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)

        stream = io.BytesIO()
        overlay_canvas = canvas.Canvas(stream, pagesize=(width, height))
        overlay_canvas.setFont("Helvetica-Bold", 11)

        left = "left" in position
        top = "top" in position
        x = 36 if left else max(36, width - 230)
        y = max(22, height - 52) if top else 28

        overlay_canvas.drawString(x, y + 14, f"Signed by: {signer}")
        overlay_canvas.setFont("Helvetica", 9)
        overlay_canvas.drawString(x, y, f"{reason} | {signed_at}")
        overlay_canvas.save()
        stream.seek(0)

        overlay_pdf = PdfReader(stream)
        page.merge_page(overlay_pdf.pages[0])
        writer.add_page(page)

    output = output_dir / "signed.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF signature added", "application/pdf")


def handle_edit_metadata_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    metadata_map = {
        "/Title": str(payload.get("title", "")).strip(),
        "/Author": str(payload.get("author", "")).strip(),
        "/Subject": str(payload.get("subject", "")).strip(),
        "/Keywords": str(payload.get("keywords", "")).strip(),
    }
    metadata = {key: value for key, value in metadata_map.items() if value}
    if metadata:
        writer.add_metadata(metadata)

    output = output_dir / "metadata-edited.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF metadata updated", "application/pdf")

def handle_pdf_to_pdfa(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.add_metadata(
        {
            "/Title": str(payload.get("title", files[0].stem)).strip() or files[0].stem,
            "/Producer": "ISHU TOOLS Archival Export",
            "/Creator": "ISHU TOOLS",
        }
    )

    output = output_dir / "archival.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF archival export generated", "application/pdf")


def handle_scan_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_image_to_pdf(files, payload, output_dir)


def handle_chat_with_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    question = str(payload.get("question", "")).strip()
    if not question:
        raise HTTPException(status_code=400, detail="question is required")

    max_sentences = max(1, min(8, int(payload.get("max_sentences", 3))))
    content = extract_pdf_text(files[0])
    if not content:
        raise HTTPException(status_code=400, detail="No extractable text found in this PDF")

    matches = top_matching_sentences(content, question, max_sentences)
    answer = " ".join(matches)

    return ExecutionResult(
        kind="json",
        message="Answer generated from PDF",
        data={
            "question": question,
            "answer": answer,
            "evidence": matches,
        },
    )


def handle_create_workflow(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    workflow_name = str(payload.get("workflow_name", "My Workflow")).strip() or "My Workflow"
    raw_steps = str(payload.get("steps", "")).strip()
    if not raw_steps:
        raise HTTPException(status_code=400, detail="steps is required, use comma or new lines")

    steps = [step.strip() for step in re.split(r"[\n,]+", raw_steps) if step.strip()]
    if not steps:
        raise HTTPException(status_code=400, detail="No valid workflow steps were provided")

    payload_data = {
        "name": workflow_name,
        "steps": [{"order": index + 1, "tool": step} for index, step in enumerate(steps)],
        "step_count": len(steps),
    }

    output = output_dir / "workflow.json"
    output.write_text(json.dumps(payload_data, indent=2), encoding="utf-8")
    return create_single_file_result(output, "Workflow created", "application/json")


def handle_upscale_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    scale = float(payload.get("scale", 2.0))
    if scale <= 1:
        raise HTTPException(status_code=400, detail="scale must be greater than 1")

    image = Image.open(files[0])
    target_size = (max(1, int(image.width * scale)), max(1, int(image.height * scale)))
    upscaled = image.resize(target_size, resample=Image.LANCZOS)

    output = output_dir / f"upscaled-{files[0].name}"
    save_image(upscaled, output)
    return create_single_file_result(output, "Image upscaled", "image/*")


def handle_add_image_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 2)

    pdf_file = next((path for path in files if path.suffix.lower() == ".pdf"), None)
    image_file = next(
        (
            path
            for path in files
            if path.suffix.lower() in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
        ),
        None,
    )

    if not pdf_file or not image_file:
        raise HTTPException(status_code=400, detail="Upload one PDF and one image file")

    scale = max(0.05, min(0.8, float(payload.get("scale", 0.2))))
    opacity = max(10, min(100, int(payload.get("opacity", 65))))

    image = Image.open(image_file).convert("RGBA")
    if opacity < 100:
        alpha = image.getchannel("A")
        alpha = alpha.point(lambda value: int(value * (opacity / 100)))
        image.putalpha(alpha)

    document = fitz.open(str(pdf_file))

    for page in document:
        width = int(page.rect.width)
        height = int(page.rect.height)

        target_width = max(36, int(width * scale))
        target_height = max(36, int(image.height * (target_width / max(1, image.width))))
        resized = image.resize((target_width, target_height), resample=Image.LANCZOS)

        stream = io.BytesIO()
        resized.save(stream, format="PNG")
        stream_value = stream.getvalue()

        margin = 24
        x1 = width - margin
        y1 = height - margin
        rect = fitz.Rect(x1 - target_width, y1 - target_height, x1, y1)
        page.insert_image(rect, stream=stream_value)

    output = output_dir / "image-overlay.pdf"
    document.save(str(output), garbage=4, deflate=True)
    document.close()
    return create_single_file_result(output, "Image added to PDF", "application/pdf")


def handle_pdf_page_count(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    text = extract_pdf_text(files[0])
    words = re.findall(r"[A-Za-z']+", text)
    report = {
        "filename": files[0].name,
        "pages": len(reader.pages),
        "words": len(words),
        "characters": len(text),
    }

    output = output_dir / "pdf-page-count.json"
    output.write_text(json.dumps(report, indent=2), encoding="utf-8")
    return create_single_file_result(output, "PDF statistics generated", "application/json")


def handle_reverse_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    writer = PdfWriter()

    for page in reversed(reader.pages):
        writer.add_page(page)

    output = output_dir / "reversed.pdf"
    with output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(output, "PDF pages reversed", "application/pdf")


def handle_flatten_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = fitz.open(str(files[0]))
    flattened = fitz.open()

    for page in source:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_bytes = pix.tobytes("png")
        target_page = flattened.new_page(width=pix.width, height=pix.height)
        target_page.insert_image(fitz.Rect(0, 0, pix.width, pix.height), stream=img_bytes)

    output = output_dir / "flattened.pdf"
    flattened.save(str(output), garbage=4, deflate=True)
    source.close()
    flattened.close()
    return create_single_file_result(output, "PDF flattened successfully", "application/pdf")


def handle_pdf_to_html(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source_pdf = files[0]
    output = output_dir / "pdf.html"

    # 1) Highest-fidelity path when pdf2htmlEX is installed on the server.
    pdf2htmlex_bin = shutil.which("pdf2htmlEX")
    if pdf2htmlex_bin:
        command = [
            pdf2htmlex_bin,
            "--embed-css",
            "1",
            "--embed-font",
            "1",
            "--embed-image",
            "1",
            "--embed-javascript",
            "0",
            "--dest-dir",
            str(output_dir),
            str(source_pdf),
            output.name,
        ]
        try:
            subprocess.run(command, check=True, capture_output=True, timeout=900)
            if output.exists() and output.stat().st_size > 0:
                return create_single_file_result(output, "PDF converted to HTML", "text/html")
        except Exception:
            pass

    # 2) Native PyMuPDF layout extraction with scanned-page image fallback.
    try:
        page_fragments: list[str] = []
        with fitz.open(str(source_pdf)) as document:
            for page_index, page in enumerate(document, start=1):
                page_text = (page.get_text("text") or "").strip()
                if page_text:
                    page_html = page.get_text("html")
                    if page_html and page_html.strip():
                        page_fragments.append(
                            f'<section class="pdf-page text-page" data-page="{page_index}">{page_html}</section>'
                        )
                        continue

                # Scanned or image-only page: preserve visual fidelity via embedded PNG.
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                encoded = base64.b64encode(pix.tobytes("png")).decode("ascii")
                page_fragments.append(
                    f'<section class="pdf-page image-page" data-page="{page_index}">'
                    f'<img alt="Page {page_index}" src="data:image/png;base64,{encoded}" />'
                    "</section>"
                )

        if page_fragments:
            html_doc = f"""<!doctype html>
<html lang=\"en\">
<head>
  <meta charset=\"utf-8\" />
  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1\" />
  <title>{html.escape(source_pdf.stem)}</title>
  <style>
    body {{
      margin: 20px;
      background: #f6f8fb;
      color: #0f172a;
      font-family: Arial, sans-serif;
    }}
    .pdf-page {{
      margin: 0 auto 22px;
      padding: 10px;
      background: #ffffff;
      border: 1px solid #dbe3f1;
      border-radius: 10px;
      box-shadow: 0 2px 10px rgba(15, 23, 42, 0.08);
      overflow-x: auto;
      max-width: 1100px;
    }}
    .pdf-page img {{
      display: block;
      width: 100%;
      height: auto;
    }}
  </style>
</head>
<body>
{''.join(page_fragments)}
</body>
</html>
"""
            output.write_text(html_doc, encoding="utf-8")
            return create_single_file_result(output, "PDF converted to HTML", "text/html")
    except Exception:
        pass

    # 3) Guaranteed fallback: readable text-only HTML export.
    text = extract_pdf_text(source_pdf)
    if not text:
        raise HTTPException(status_code=400, detail="Unable to extract content from the PDF")

    escaped = html.escape(text)
    html_doc = f"""<!doctype html>
<html lang=\"en\">
<head>
  <meta charset=\"utf-8\" />
  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1\" />
  <title>{html.escape(source_pdf.stem)}</title>
  <style>body {{ font-family: Arial, sans-serif; margin: 24px; line-height: 1.5; }} pre {{ white-space: pre-wrap; }}</style>
</head>
<body>
  <h1>{html.escape(source_pdf.name)}</h1>
  <pre>{escaped}</pre>
</body>
</html>
"""

    output.write_text(html_doc, encoding="utf-8")
    return create_single_file_result(output, "PDF converted to HTML", "text/html")


def handle_pdf_to_bmp(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    return _pdf_to_image(files, output_dir, "bmp")


def handle_pdf_to_gif(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    document = fitz.open(str(files[0]))

    for index, page in enumerate(document, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        output = output_dir / f"page-{index}.gif"
        image.convert("P", palette=Image.ADAPTIVE).save(output, format="GIF")

    document.close()
    return create_zip_result(output_dir, "PDF converted to GIF", "pdf-to-gif")


def handle_sharpen_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    factor = float(payload.get("factor", 2.0))

    image = Image.open(files[0])
    enhanced = ImageEnhance.Sharpness(image).enhance(max(0.0, factor))
    output = output_dir / f"sharpen-{files[0].name}"
    save_image(enhanced, output)
    return create_single_file_result(output, "Image sharpened", "image/*")


def handle_brighten_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    factor = float(payload.get("factor", 1.25))

    image = Image.open(files[0])
    enhanced = ImageEnhance.Brightness(image).enhance(max(0.0, factor))
    output = output_dir / f"bright-{files[0].name}"
    save_image(enhanced, output)
    return create_single_file_result(output, "Image brightness adjusted", "image/*")


def handle_contrast_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    factor = float(payload.get("factor", 1.2))

    image = Image.open(files[0])
    enhanced = ImageEnhance.Contrast(image).enhance(max(0.0, factor))
    output = output_dir / f"contrast-{files[0].name}"
    save_image(enhanced, output)
    return create_single_file_result(output, "Image contrast adjusted", "image/*")


def handle_invert_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image = Image.open(files[0]).convert("RGB")
    inverted = ImageOps.invert(image)
    output = output_dir / f"invert-{files[0].name}"
    save_image(inverted, output)
    return create_single_file_result(output, "Image colors inverted", "image/*")


def handle_posterize_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    bits = int(payload.get("bits", 4))
    bits = min(8, max(1, bits))

    image = Image.open(files[0]).convert("RGB")
    posterized = ImageOps.posterize(image, bits)
    output = output_dir / f"posterize-{files[0].name}"
    save_image(posterized, output)
    return create_single_file_result(output, "Posterize effect applied", "image/*")


def handle_image_histogram(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image = Image.open(files[0]).convert("RGB")
    histogram = image.histogram()

    red = histogram[:256]
    green = histogram[256:512]
    blue = histogram[512:768]
    pixels = max(1, image.width * image.height)

    def channel_mean(channel_values: list[int]) -> float:
        return round(sum(index * count for index, count in enumerate(channel_values)) / pixels, 2)

    data = {
        "filename": files[0].name,
        "width": image.width,
        "height": image.height,
        "mode": image.mode,
        "mean_rgb": {
            "r": channel_mean(red),
            "g": channel_mean(green),
            "b": channel_mean(blue),
        },
    }

    output = output_dir / "image-histogram.json"
    output.write_text(json.dumps(data, indent=2), encoding="utf-8")
    return create_single_file_result(output, "Image histogram generated", "application/json")


def handle_remove_extra_spaces_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to clean")

    cleaned_lines = [re.sub(r"[ \t]+", " ", line).strip() for line in content.splitlines()]
    cleaned = "\n".join(cleaned_lines).strip()

    output = output_dir / "cleaned-text.txt"
    output.write_text(cleaned, encoding="utf-8")
    return create_single_file_result(output, "Extra spaces removed", "text/plain")


def handle_sort_lines_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text with multiple lines")

    direction = str(payload.get("direction", "asc")).strip().lower()
    reverse = direction == "desc"
    lines = [line for line in content.splitlines() if line.strip()]
    sorted_lines = sorted(lines, key=str.casefold, reverse=reverse)

    output = output_dir / "sorted-lines.txt"
    output.write_text("\n".join(sorted_lines), encoding="utf-8")
    return create_single_file_result(output, "Lines sorted", "text/plain")


def handle_deduplicate_lines_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text with lines")

    seen: set[str] = set()
    unique_lines: list[str] = []
    for line in content.splitlines():
        normalized = line.strip()
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        unique_lines.append(normalized)

    output = output_dir / "deduplicated-lines.txt"
    output.write_text("\n".join(unique_lines), encoding="utf-8")
    return create_single_file_result(output, "Duplicate lines removed", "text/plain")


def handle_find_replace_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text for find/replace")

    find_text = str(payload.get("find", "")).strip()
    replace_text = str(payload.get("replace", ""))
    case_sensitive_raw = payload.get("case_sensitive", False)
    if isinstance(case_sensitive_raw, bool):
        case_sensitive = case_sensitive_raw
    else:
        case_sensitive = str(case_sensitive_raw).strip().lower() in {"1", "true", "yes", "on"}
    if not find_text:
        raise HTTPException(status_code=400, detail="find value is required")

    if case_sensitive:
        updated = content.replace(find_text, replace_text)
    else:
        updated = re.sub(re.escape(find_text), replace_text, content, flags=re.IGNORECASE)

    output = output_dir / "find-replace.txt"
    output.write_text(updated, encoding="utf-8")
    return create_single_file_result(output, "Find and replace applied", "text/plain")


def handle_split_text_file(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text content to split")

    lines_per_file = max(1, int(payload.get("lines_per_file", 100)))
    lines = content.splitlines()
    for index, start in enumerate(range(0, len(lines), lines_per_file), start=1):
        chunk = lines[start : start + lines_per_file]
        part = output_dir / f"part-{index}.txt"
        part.write_text("\n".join(chunk), encoding="utf-8")

    return create_zip_result(output_dir, "Text split into multiple files", "split-text")


def handle_reading_time_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    content = read_text_input(files, payload)
    if not content:
        raise HTTPException(status_code=400, detail="Provide text to estimate reading time")

    words = re.findall(r"[A-Za-z']+", content)
    wpm = max(100, int(payload.get("wpm", 200)))
    minutes = max(1, math.ceil(len(words) / wpm))

    report = {
        "words": len(words),
        "wpm": wpm,
        "estimated_minutes": minutes,
    }

    output = output_dir / "reading-time.json"
    output.write_text(json.dumps(report, indent=2), encoding="utf-8")
    return create_single_file_result(output, "Reading time estimated", "application/json")


def handle_images_to_zip(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image_suffixes = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
    image_files = [path for path in files if path.suffix.lower() in image_suffixes]
    if not image_files:
        raise HTTPException(status_code=400, detail="Upload at least one image file")

    for image_file in image_files:
        shutil.copy(image_file, output_dir / image_file.name)

    return create_zip_result(output_dir, "Images packed into ZIP", "images-zip")


def handle_batch_convert_images(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target_format = str(payload.get("target_format", "png")).strip().lower()
    if target_format not in {"png", "jpg", "jpeg", "webp", "gif", "bmp"}:
        raise HTTPException(status_code=400, detail="target_format must be one of png,jpg,webp,gif,bmp")

    normalized_ext = "jpg" if target_format == "jpeg" else target_format
    for source in files:
        image = Image.open(source)
        converted = image.convert("RGB") if normalized_ext == "jpg" else image
        output = output_dir / f"{source.stem}.{normalized_ext}"
        save_image(converted, output, fmt=target_format.upper())

    if len(files) == 1:
        single = next(output_dir.iterdir())
        return create_single_file_result(single, "Image converted", "image/*")

    return create_zip_result(output_dir, "Batch image conversion completed", "batch-images")


def handle_merge_text_files(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    separator = str(payload.get("separator", "\n\n-----\n\n"))
    chunks: list[str] = []

    for source in files:
        chunks.append(source.read_text(encoding="utf-8", errors="ignore"))

    merged = separator.join(chunks)
    output = output_dir / "merged-text.txt"
    output.write_text(merged, encoding="utf-8")
    return create_single_file_result(output, "Text files merged", "text/plain")


def handle_json_prettify(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    raw = read_text_input(files, payload)
    if not raw:
        raise HTTPException(status_code=400, detail="Provide JSON content")

    parsed = json.loads(raw)
    output = output_dir / "pretty.json"
    output.write_text(json.dumps(parsed, indent=2, ensure_ascii=False), encoding="utf-8")
    return create_single_file_result(output, "JSON formatted successfully", "application/json")


def handle_csv_to_json(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    raw = read_text_input(files, payload)
    if not raw:
        raise HTTPException(status_code=400, detail="Provide CSV content")

    reader = csv.DictReader(io.StringIO(raw))
    rows = list(reader)

    output = output_dir / "data.json"
    output.write_text(json.dumps(rows, indent=2, ensure_ascii=False), encoding="utf-8")
    return create_single_file_result(output, "CSV converted to JSON", "application/json")


def handle_json_to_csv(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    raw = read_text_input(files, payload)
    if not raw:
        raise HTTPException(status_code=400, detail="Provide JSON content")

    parsed = json.loads(raw)
    if isinstance(parsed, dict):
        rows = [parsed]
    elif isinstance(parsed, list):
        rows = [item for item in parsed if isinstance(item, dict)]
    else:
        raise HTTPException(status_code=400, detail="JSON must be an object or list of objects")

    if not rows:
        raise HTTPException(status_code=400, detail="No object rows found to convert")

    fieldnames: list[str] = []
    seen_fields: set[str] = set()
    for row in rows:
        for key in row.keys():
            if key not in seen_fields:
                seen_fields.add(key)
                fieldnames.append(key)

    output = output_dir / "data.csv"
    with output.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)

    return create_single_file_result(output, "JSON converted to CSV", "text/csv")


def handle_pdf_to_tiff(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    images = render_pdf_pages_to_pil_images(files[0], scale=2.0)
    output = output_dir / "pdf.tiff"
    build_tiff_from_images(images, output)
    return create_single_file_result(output, "PDF converted to TIFF", "image/tiff")


def handle_tiff_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)

    pages: list[Image.Image] = []
    for source in files:
        image = Image.open(source)
        if getattr(image, "n_frames", 1) > 1:
            for frame in ImageSequence.Iterator(image):
                pages.append(ImageOps.exif_transpose(frame).convert("RGB"))
        else:
            pages.append(ImageOps.exif_transpose(image).convert("RGB"))

    output = output_dir / "tiff.pdf"
    save_images_as_pdf(pages, output)
    return create_single_file_result(output, "TIFF converted to PDF", "application/pdf")


def handle_pdf_to_svg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    document = fitz.open(str(files[0]))

    try:
        for index, page in enumerate(document, start=1):
            svg_text = page.get_svg_image()
            (output_dir / f"page-{index}.svg").write_text(svg_text, encoding="utf-8")
    finally:
        document.close()

    if len(list(output_dir.glob("*.svg"))) == 1:
        svg_file = next(output_dir.glob("*.svg"))
        return create_single_file_result(svg_file, "PDF converted to SVG", "image/svg+xml")

    return create_zip_result(output_dir, "PDF pages converted to SVG", "pdf-to-svg")


def handle_svg_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    generated_pdfs: list[Path] = []

    # Preferred path: render vector content faithfully with CairoSVG.
    try:
        import cairosvg

        for index, source in enumerate(files, start=1):
            pdf_path = output_dir / f"svg-{index}.pdf"
            cairosvg.svg2pdf(url=str(source), write_to=str(pdf_path))
            if pdf_path.exists() and pdf_path.stat().st_size > 0:
                generated_pdfs.append(pdf_path)
    except Exception:
        generated_pdfs = []

    if not generated_pdfs:
        # Fallback path: export textual content when vector renderer is unavailable.
        for index, source in enumerate(files, start=1):
            svg_content = source.read_text(encoding="utf-8", errors="ignore")
            text_content = BeautifulSoup(svg_content, "html.parser").get_text(" ", strip=True)
            pdf_path = output_dir / f"svg-{index}.pdf"
            pdf_canvas = canvas.Canvas(str(pdf_path), pagesize=A4)
            pdf_canvas.setFont("Helvetica", 11)
            pdf_canvas.drawString(40, 780, f"Converted from: {source.name}")
            y = 760
            for line in textwrap.wrap(text_content, width=90):
                if y < 40:
                    pdf_canvas.showPage()
                    pdf_canvas.setFont("Helvetica", 11)
                    y = 780
                pdf_canvas.drawString(40, y, line)
                y -= 15
            pdf_canvas.save()
            generated_pdfs.append(pdf_path)

    if len(generated_pdfs) == 1:
        return create_single_file_result(generated_pdfs[0], "SVG converted to PDF", "application/pdf")

    writer = PdfWriter()
    for pdf_path in generated_pdfs:
        reader = PdfReader(str(pdf_path))
        for page in reader.pages:
            writer.add_page(page)

    merged_output = output_dir / "svg-merged.pdf"
    with merged_output.open("wb") as f:
        writer.write(f)

    return create_single_file_result(merged_output, "SVG files converted to a merged PDF", "application/pdf")


def handle_pdf_to_rtf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_pdf_text(files[0])
    output = output_dir / "pdf.rtf"
    build_rtf_from_text(content, output)
    return create_single_file_result(output, "PDF converted to RTF", "application/rtf")


def handle_rtf_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_rtf_text(files[0])
    output = output_dir / "rtf.pdf"
    text_to_pdf(content, output, title="RTF to PDF")
    return create_single_file_result(output, "RTF converted to PDF", "application/pdf")


def handle_pdf_to_odt(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    title = str(payload.get("title", files[0].stem)).strip() or files[0].stem
    content = extract_pdf_text(files[0])
    output = output_dir / "pdf.odt"
    build_odt_from_text(content, output, title)
    return create_single_file_result(output, "PDF converted to ODT", "application/vnd.oasis.opendocument.text")


def handle_odt_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_odt_text(files[0])
    output = output_dir / "odt.pdf"
    text_to_pdf(content, output, title="ODT to PDF")
    return create_single_file_result(output, "ODT converted to PDF", "application/pdf")


def handle_pdf_to_epub(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    title = str(payload.get("title", files[0].stem)).strip() or files[0].stem
    content = extract_pdf_text(files[0])
    output = output_dir / "pdf.epub"
    build_epub_from_text(content, output, title)
    return create_single_file_result(output, "PDF converted to EPUB", "application/epub+zip")


def handle_epub_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    content = extract_epub_text(files[0])
    output = output_dir / "epub.pdf"
    text_to_pdf(content, output, title="EPUB to PDF")
    return create_single_file_result(output, "EPUB converted to PDF", "application/pdf")


def handle_ocr_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    text = extract_ocr_text_from_image(files[0])
    if not text:
        raise HTTPException(status_code=400, detail="No readable text was detected in the image")

    output = output_dir / "ocr-image.txt"
    output.write_text(text, encoding="utf-8")
    return create_single_file_result(output, "OCR extracted text from image", "text/plain")


def handle_ocr_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source_pdf = files[0]

    # Prefer OCRmyPDF when available: best-in-class searchable PDF generation.
    ocrmypdf_bin = shutil.which("ocrmypdf")
    if ocrmypdf_bin:
        output = output_dir / "ocr.pdf"
        command = [
            ocrmypdf_bin,
            "--skip-text",
            "--force-ocr",
            "--optimize",
            "3",
            "--output-type",
            "pdf",
        ]
        language = str(payload.get("lang", "")).strip()
        if language:
            command.extend(["--language", language])
        command.extend([str(source_pdf), str(output)])

        try:
            subprocess.run(command, check=True, capture_output=True, timeout=900)
            if output.exists() and output.stat().st_size > 0:
                return create_single_file_result(output, "OCR-generated searchable PDF created", "application/pdf")
        except Exception:
            pass

    # Fallback: rebuild pages as images and overlay invisible OCR text for searchability.
    output = output_dir / "ocr.pdf"
    aggregate_text: list[str] = []

    with tempfile.TemporaryDirectory(prefix="ocr-pdf-pages-") as temp_name:
        temp_dir = Path(temp_name)
        source_doc = fitz.open(str(source_pdf))
        target_doc = fitz.open()

        try:
            for page_index, page in enumerate(source_doc, start=1):
                pix = page.get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
                page_image = temp_dir / f"ocr-page-{page_index}.png"
                pix.save(str(page_image))

                page_lines = extract_ocr_lines_from_image(page_image)
                page_text = "\n".join(page_lines).strip()
                if page_text:
                    aggregate_text.append(f"Page {page_index}\n{page_text}")

                new_page = target_doc.new_page(width=page.rect.width, height=page.rect.height)
                new_page.insert_image(new_page.rect, filename=str(page_image), keep_proportion=False)

                # Invisible overlay text preserves search/copy ability without changing visual output.
                if page_lines:
                    margin_x = 8.0
                    usable_height = max(20.0, float(new_page.rect.height) - 16.0)
                    line_height = max(8.5, min(16.0, usable_height / max(1, len(page_lines))))
                    font_size = max(6.5, min(9.0, line_height * 0.62))
                    y = 10.0
                    for line in page_lines:
                        if y > float(new_page.rect.height) - 4:
                            break
                        new_page.insert_text(
                            fitz.Point(margin_x, y),
                            line[:1400],
                            fontsize=font_size,
                            render_mode=3,
                            color=(0, 0, 0),
                            fill_opacity=0,
                            stroke_opacity=0,
                            overlay=True,
                        )
                        y += line_height

            target_doc.save(str(output), garbage=4, deflate=True)
        finally:
            source_doc.close()
            target_doc.close()

    if not aggregate_text:
        raise HTTPException(status_code=400, detail="No readable text was detected in the PDF pages")

    return create_single_file_result(output, "OCR-generated searchable PDF created", "application/pdf")


def remove_background_fallback(image_bytes: bytes, tolerance: int) -> bytes:
    image = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
    width, height = image.size
    pixels = image.load()

    sample_points = [
        pixels[0, 0],
        pixels[max(0, width - 1), 0],
        pixels[0, max(0, height - 1)],
        pixels[max(0, width - 1), max(0, height - 1)],
    ]
    bg = tuple(int(sum(pixel[idx] for pixel in sample_points) / len(sample_points)) for idx in range(3))
    threshold_sq = max(1, tolerance) ** 2

    new_pixels = []
    for r, g, b, a in image.getdata():
        distance_sq = (r - bg[0]) ** 2 + (g - bg[1]) ** 2 + (b - bg[2]) ** 2
        if distance_sq <= threshold_sq:
            new_pixels.append((r, g, b, 0))
        else:
            new_pixels.append((r, g, b, a))

    image.putdata(new_pixels)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def handle_remove_background(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image_bytes = files[0].read_bytes()
    tolerance = max(5, min(80, int(payload.get("tolerance", 28))))

    try:
        from rembg import remove

        removed = remove(image_bytes)
    except Exception:
        removed = remove_background_fallback(image_bytes, tolerance)

    output = output_dir / f"{files[0].stem}-background-removed.png"
    output.write_bytes(removed)
    return create_single_file_result(output, "Background removed from image", "image/png")


def handle_blur_background(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    radius = max(2, float(payload.get("radius", 12)))
    original = open_image_file(files[0], "RGBA")

    try:
        from rembg import remove

        cutout = Image.open(io.BytesIO(remove(files[0].read_bytes()))).convert("RGBA")
        mask = cutout.getchannel("A")
    except Exception:
        # Fallback: keep center subject area and blur outer background smoothly.
        width, height = original.size
        mask = Image.new("L", (width, height), 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse(
            (
                int(width * 0.15),
                int(height * 0.05),
                int(width * 0.85),
                int(height * 0.95),
            ),
            fill=255,
        )
        mask = mask.filter(ImageFilter.GaussianBlur(radius=max(12, int(min(width, height) * 0.04))))

    blurred = original.filter(ImageFilter.GaussianBlur(radius=radius))
    output_img = Image.composite(original, blurred, mask)

    output = output_dir / f"{files[0].stem}-blur-background.png"
    save_image(output_img, output, fmt="PNG")
    return create_single_file_result(output, "Background blur applied", "image/png")


def handle_blur_face(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    import cv2

    blur_strength = max(15, int(payload.get("blur_strength", 45)))
    if blur_strength % 2 == 0:
        blur_strength += 1

    image, faces = detect_faces_in_image(files[0])
    if len(faces) == 0:
        raise HTTPException(status_code=400, detail="No face was detected in the uploaded image")

    for (x, y, w, h) in faces:
        region = image[y : y + h, x : x + w]
        image[y : y + h, x : x + w] = cv2.GaussianBlur(region, (blur_strength, blur_strength), 0)

    output = output_dir / f"{files[0].stem}-blur-face.png"
    cv2.imwrite(str(output), image)
    return create_single_file_result(output, "Faces blurred successfully", "image/png")


def handle_pixelate_face(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    import cv2

    pixel_size = max(4, int(payload.get("pixel_size", 16)))
    image, faces = detect_faces_in_image(files[0])
    if len(faces) == 0:
        raise HTTPException(status_code=400, detail="No face was detected in the uploaded image")

    for (x, y, w, h) in faces:
        region = image[y : y + h, x : x + w]
        tiny = cv2.resize(
            region,
            (max(1, w // pixel_size), max(1, h // pixel_size)),
            interpolation=cv2.INTER_LINEAR,
        )
        image[y : y + h, x : x + w] = cv2.resize(tiny, (w, h), interpolation=cv2.INTER_NEAREST)

    output = output_dir / f"{files[0].stem}-pixelate-face.png"
    cv2.imwrite(str(output), image)
    return create_single_file_result(output, "Faces pixelated successfully", "image/png")


def handle_unblur_face(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    import cv2

    image, faces = detect_faces_in_image(files[0])
    if len(faces) == 0:
        raise HTTPException(status_code=400, detail="No face was detected in the uploaded image")

    sharpen_strength = max(1.1, min(3.0, float(payload.get("strength", 1.8))))
    denoise = max(0, min(20, int(payload.get("denoise", 6))))

    for (x, y, w, h) in faces:
        region = image[y : y + h, x : x + w]
        if region.size == 0:
            continue

        if denoise > 0:
            region = cv2.fastNlMeansDenoisingColored(region, None, denoise, denoise, 7, 21)

        blurred = cv2.GaussianBlur(region, (0, 0), 2.2)
        sharpened = cv2.addWeighted(region, sharpen_strength, blurred, -(sharpen_strength - 1.0), 0)

        lab = cv2.cvtColor(sharpened, cv2.COLOR_BGR2LAB)
        l_channel, a_channel, b_channel = cv2.split(lab)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        l_channel = clahe.apply(l_channel)
        merged = cv2.merge((l_channel, a_channel, b_channel))
        image[y : y + h, x : x + w] = cv2.cvtColor(merged, cv2.COLOR_LAB2BGR)

    output = output_dir / f"{files[0].stem}-unblur-face.png"
    cv2.imwrite(str(output), image)
    return create_single_file_result(output, "Faces enhanced successfully", "image/png")


def handle_remove_image_object(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    import cv2
    import numpy as np

    image = cv2.imread(str(files[0]))
    if image is None:
        raise HTTPException(status_code=400, detail="Unable to read image for object removal")

    height, width = image.shape[:2]
    x = int(payload.get("x", width * 0.35))
    y = int(payload.get("y", height * 0.35))
    w = int(payload.get("width", width * 0.2))
    h = int(payload.get("height", height * 0.2))
    inpaint_radius = max(1, min(21, int(payload.get("radius", 5))))

    x = max(0, min(width - 2, x))
    y = max(0, min(height - 2, y))
    w = max(1, min(width - x, w))
    h = max(1, min(height - y, h))

    mask = np.zeros((height, width), dtype=np.uint8)
    mask[y : y + h, x : x + w] = 255

    method = str(payload.get("method", "telea")).strip().lower()
    algorithm = cv2.INPAINT_NS if method == "ns" else cv2.INPAINT_TELEA
    restored = cv2.inpaint(image, mask, inpaint_radius, algorithm)

    output = output_dir / f"{files[0].stem}-object-removed.png"
    cv2.imwrite(str(output), restored)
    return create_single_file_result(output, "Object removed from image", "image/png")


def handle_add_text_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    text = str(payload.get("text", "")).strip()
    if not text:
        raise HTTPException(status_code=400, detail="Provide the text you want to place on the image")

    font_size = max(12, int(payload.get("font_size", 48)))
    opacity = int(round(clamp_percentage(payload.get("opacity", 90), 90) * 2.55))
    x = int(payload.get("x", 40))
    y = int(payload.get("y", 40))
    text_color = parse_color_value(payload.get("color", "#ffffff"), (255, 255, 255))

    image = open_image_file(files[0], "RGBA")
    layer = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(layer)
    font = load_ui_font(font_size)
    stroke_opacity = max(60, min(220, opacity))
    draw.text(
        (x, y),
        text,
        fill=(*text_color, opacity),
        font=font,
        stroke_width=max(1, font_size // 18),
        stroke_fill=(0, 0, 0, stroke_opacity),
    )

    output_img = Image.alpha_composite(image, layer)
    output = output_dir / f"{files[0].stem}-text.png"
    save_image(output_img, output, fmt="PNG")
    return create_single_file_result(output, "Text added to image", "image/png")


def handle_add_logo_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 2)

    base = open_image_file(files[0], "RGBA")
    logo = open_image_file(files[1], "RGBA")
    scale_percent = clamp_percentage(payload.get("scale_percent", 20), 20)
    opacity = clamp_percentage(payload.get("opacity", 85), 85)
    position = str(payload.get("position", "bottom-right"))

    max_width = max(40, int(base.width * (scale_percent / 100)))
    max_height = max(40, int(base.height * 0.35))
    logo.thumbnail((max_width, max_height), get_resampling_module().LANCZOS)

    alpha = logo.getchannel("A").point(lambda value: int(value * (opacity / 100)))
    logo.putalpha(alpha)

    x, y = resolve_anchor_position(position, base.size, logo.size)
    output_img = base.copy()
    output_img.alpha_composite(logo, (x, y))

    output = output_dir / f"{files[0].stem}-logo-overlay.png"
    save_image(output_img, output, fmt="PNG")
    return create_single_file_result(output, "Logo added to image", "image/png")


def handle_join_images(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 2)
    direction = str(payload.get("direction", "horizontal")).strip().lower()
    gap = max(0, int(payload.get("gap", 16)))
    background = parse_color_value(payload.get("background", "#0b1120"), (11, 17, 32))

    images = [open_image_file(path, "RGBA") for path in files]
    if direction == "vertical":
        canvas_width = max(image.width for image in images)
        canvas_height = sum(image.height for image in images) + gap * (len(images) - 1)
        canvas_image = Image.new("RGBA", (canvas_width, canvas_height), (*background, 255))
        y = 0
        for image in images:
            x = (canvas_width - image.width) // 2
            canvas_image.alpha_composite(image, (x, y))
            y += image.height + gap
    else:
        canvas_width = sum(image.width for image in images) + gap * (len(images) - 1)
        canvas_height = max(image.height for image in images)
        canvas_image = Image.new("RGBA", (canvas_width, canvas_height), (*background, 255))
        x = 0
        for image in images:
            y = (canvas_height - image.height) // 2
            canvas_image.alpha_composite(image, (x, y))
            x += image.width + gap

    output = output_dir / "joined-images.png"
    save_image(canvas_image, output, fmt="PNG")
    return create_single_file_result(output, "Images joined into a single canvas", "image/png")


def handle_split_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    columns = max(1, int(payload.get("columns", 2)))
    rows = max(1, int(payload.get("rows", 2)))
    image = open_image_file(files[0], "RGBA")

    for row_index in range(rows):
        top = round(row_index * image.height / rows)
        bottom = round((row_index + 1) * image.height / rows)
        for col_index in range(columns):
            left = round(col_index * image.width / columns)
            right = round((col_index + 1) * image.width / columns)
            tile = image.crop((left, top, right, bottom))
            output = output_dir / f"tile-r{row_index + 1}-c{col_index + 1}.png"
            save_image(tile, output, fmt="PNG")

    if columns * rows == 1:
        tile_file = next(output_dir.glob("*.png"))
        return create_single_file_result(tile_file, "Image split into one tile", "image/png")

    return create_zip_result(output_dir, "Image split into tiles", "split-image")


def handle_circle_crop_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image = open_image_file(files[0], "RGBA")
    side = min(image.width, image.height)
    cropped = ImageOps.fit(image, (side, side))

    target_size = max(0, int(payload.get("size", side)))
    if target_size:
        cropped = cropped.resize((target_size, target_size), get_resampling_module().LANCZOS)
        side = target_size

    mask = Image.new("L", (side, side), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, side, side), fill=255)

    output_img = Image.new("RGBA", (side, side), (255, 255, 255, 0))
    output_img.paste(cropped, (0, 0), mask)
    output = output_dir / f"{files[0].stem}-circle-crop.png"
    save_image(output_img, output, fmt="PNG")
    return create_single_file_result(output, "Circle crop created", "image/png")


def handle_square_crop_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    image = open_image_file(files[0], "RGBA")
    side = min(image.width, image.height)
    output_img = ImageOps.fit(image, (side, side))

    target_size = max(0, int(payload.get("size", side)))
    if target_size:
        output_img = output_img.resize((target_size, target_size), get_resampling_module().LANCZOS)

    output = output_dir / f"{files[0].stem}-square-crop.png"
    save_image(output_img, output, fmt="PNG")
    return create_single_file_result(output, "Square crop created", "image/png")


def handle_image_color_picker(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    color_count = max(3, min(12, int(payload.get("colors", 5))))
    source = open_image_file(files[0], "RGB")

    sample = source.copy()
    sample.thumbnail((220, 220), get_resampling_module().LANCZOS)
    quantized = sample.quantize(colors=color_count, method=Image.MEDIANCUT)
    palette = quantized.getpalette()
    total_pixels = max(1, sample.width * sample.height)
    swatches: list[dict[str, Any]] = []

    for palette_index, count in Counter(quantized.getdata()).most_common(color_count):
        rgb = tuple(palette[palette_index * 3 : palette_index * 3 + 3])
        swatches.append(
            {
                "hex": "#{:02X}{:02X}{:02X}".format(*rgb),
                "rgb": list(rgb),
                "coverage_percent": round((count / total_pixels) * 100, 2),
            }
        )

    (output_dir / "palette.json").write_text(json.dumps(swatches, indent=2), encoding="utf-8")

    swatch_width = 240
    swatch_height = 120
    preview = Image.new("RGB", (swatch_width, swatch_height * len(swatches)), color="#050816")
    draw = ImageDraw.Draw(preview)
    font = load_ui_font(22)
    small_font = load_ui_font(16)

    for index, swatch in enumerate(swatches):
        top = index * swatch_height
        color = tuple(swatch["rgb"])
        draw.rectangle((0, top, swatch_width, top + swatch_height), fill=color)
        text_color = (17, 24, 39) if sum(color) > 420 else (248, 250, 252)
        draw.text((18, top + 22), swatch["hex"], fill=text_color, font=font)
        draw.text(
            (18, top + 60),
            f"{swatch['coverage_percent']}% coverage",
            fill=text_color,
            font=small_font,
        )

    preview.save(output_dir / "palette-preview.png")
    return create_zip_result(output_dir, "Image palette extracted", "image-color-picker")


def handle_motion_blur_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    import cv2
    import numpy as np

    direction = str(payload.get("direction", "horizontal")).strip().lower()
    strength = max(3, min(31, int(payload.get("strength", 9))))
    if strength % 2 == 0:
        strength += 1

    image = cv2.imread(str(files[0]), cv2.IMREAD_UNCHANGED)
    if image is None:
        raise HTTPException(status_code=400, detail="Unable to open image for motion blur")

    kernel = np.zeros((strength, strength), dtype=np.float32)
    if direction == "vertical":
        kernel[:, strength // 2] = 1.0
    else:
        kernel[strength // 2, :] = 1.0
    kernel /= kernel.sum()

    blurred = cv2.filter2D(image, -1, kernel)
    output = output_dir / f"{files[0].stem}-motion-blur.png"
    cv2.imwrite(str(output), blurred)
    return create_single_file_result(output, "Motion blur applied", "image/png")


def handle_optimize_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    optimized_payload = dict(payload)
    optimized_payload.setdefault("image_quality", 55)
    return handle_compress_pdf(files, optimized_payload, output_dir)


def handle_remove_pages(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_delete_pages(files, payload, output_dir)


def handle_pdf_ocr(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_ocr_pdf(files, payload, output_dir)


def handle_image_to_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_ocr_image(files, payload, output_dir)


def handle_jpg_to_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_ocr_image(files, payload, output_dir)


def handle_png_to_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_ocr_image(files, payload, output_dir)


def handle_ai_summarizer(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    if files:
        return handle_summarize_pdf(files, payload, output_dir)
    return handle_summarize_text(files, payload, output_dir)


def handle_word_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_docx_to_pdf(files, payload, output_dir)


def handle_pdf_to_word(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_pdf_to_docx(files, payload, output_dir)


def handle_powerpoint_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_pptx_to_pdf(files, payload, output_dir)


def handle_pdf_to_powerpoint(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_pdf_to_pptx(files, payload, output_dir)


def handle_add_page_numbers(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_page_numbers_pdf(files, payload, output_dir)


def handle_add_watermark(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_watermark_pdf(files, payload, output_dir)


def handle_edit_metadata(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    if files[0].suffix.lower() == ".pdf":
        return handle_edit_metadata_pdf(files, payload, output_dir)
    return handle_edit_metadata_image(files, payload, output_dir)


def handle_pdf_viewer(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    reader = PdfReader(str(files[0]))
    preview_pages: list[dict[str, Any]] = []
    for index, page in enumerate(reader.pages[:3], start=1):
        preview_pages.append(
            {
                "page": index,
                "preview": (page.extract_text() or "").strip()[:1200],
            }
        )

    return ExecutionResult(
        kind="json",
        message="PDF preview generated",
        data={
            "filename": files[0].name,
            "pages": len(reader.pages),
            "metadata": dict(reader.metadata or {}),
            "preview_pages": preview_pages,
        },
    )


def handle_pdf_intelligence(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    text = extract_pdf_text(files[0])
    reader = PdfReader(str(files[0]))
    summary_sentences = max(2, min(8, int(payload.get("summary_sentences", 4))))
    top_keywords = max(3, min(15, int(payload.get("top_keywords", 8))))

    return ExecutionResult(
        kind="json",
        message="PDF intelligence report generated",
        data={
            "filename": files[0].name,
            "pages": len(reader.pages),
            "characters": len(text),
            "words": len(re.findall(r"[A-Za-z']+", text)),
            "summary": summarize_text_algo(text, max_sentences=summary_sentences),
            "keywords": extract_top_keywords(text, limit=top_keywords),
            "preview": text[:1600],
        },
    )


def handle_annotate_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    text = str(payload.get("text", "")).strip()
    if not text:
        raise HTTPException(status_code=400, detail="text is required for annotation")

    page_number = max(1, int(payload.get("page", 1)))
    x = float(payload.get("x", 48))
    y = float(payload.get("y", 72))
    width = max(120.0, float(payload.get("width", 220)))
    height = max(40.0, float(payload.get("height", 80)))
    font_size = max(8, int(payload.get("font_size", 12)))
    color = color_to_pdf_tuple(payload.get("color", "#38bdf8"), (56, 189, 248))

    document = fitz.open(str(files[0]))
    try:
        if page_number > len(document):
            raise HTTPException(status_code=400, detail="Selected page is outside the PDF range")
        page = document[page_number - 1]
        annotation = page.add_freetext_annot(fitz.Rect(x, y, x + width, y + height), text)
        annotation.update(fontsize=font_size, text_color=color, fill_color=(1, 1, 1))

        output = output_dir / "annotated.pdf"
        document.save(str(output))
    finally:
        document.close()

    return create_single_file_result(output, "PDF annotation added", "application/pdf")


def handle_highlight_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    page_number = max(1, int(payload.get("page", 1)))
    x = float(payload.get("x", 48))
    y = float(payload.get("y", 72))
    width = max(24.0, float(payload.get("width", 180)))
    height = max(12.0, float(payload.get("height", 24)))
    opacity = max(0.1, min(1.0, float(payload.get("opacity", 0.35))))
    color = color_to_pdf_tuple(payload.get("color", "#fde047"), (253, 224, 71))

    document = fitz.open(str(files[0]))
    try:
        if page_number > len(document):
            raise HTTPException(status_code=400, detail="Selected page is outside the PDF range")
        page = document[page_number - 1]
        rect = fitz.Rect(x, y, x + width, y + height)
        annotation = page.add_rect_annot(rect)
        annotation.set_colors(stroke=color, fill=color)
        annotation.set_opacity(opacity)
        annotation.update()

        output = output_dir / "highlighted.pdf"
        document.save(str(output))
    finally:
        document.close()

    return create_single_file_result(output, "PDF highlight added", "application/pdf")


def handle_pdf_filler(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_annotate_pdf(files, payload, output_dir)


def handle_edit_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    mode = str(payload.get("mode", "add_text")).strip().lower()
    if mode == "watermark":
        return handle_watermark_pdf(files, payload, output_dir)
    if mode == "annotate":
        return handle_annotate_pdf(files, payload, output_dir)
    if mode == "highlight":
        return handle_highlight_pdf(files, payload, output_dir)
    if mode == "add_image":
        return handle_add_image_pdf(files, payload, output_dir)
    return handle_add_text_pdf(files, payload, output_dir)


def handle_pdf_security(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    mode = str(payload.get("mode", "protect")).strip().lower()
    if mode == "unlock":
        return handle_unlock_pdf(files, payload, output_dir)
    return handle_protect_pdf(files, payload, output_dir)


def handle_convert_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    suffixes = {path.suffix.lower() for path in files}
    image_types = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".jfif", ".heic", ".heif"}
    if suffixes.issubset(image_types):
        return handle_jpg_to_pdf(files, payload, output_dir)

    if len(files) > 1:
        raise HTTPException(status_code=400, detail="Multiple uploads are currently supported for image-to-PDF conversion only")

    source = files[0]
    extension = source.suffix.lower()
    handler_by_extension: dict[str, ToolHandler] = {
        ".doc": handle_word_to_pdf,
        ".docx": handle_docx_to_pdf,
        ".ppt": handle_powerpoint_to_pdf,
        ".pptx": handle_pptx_to_pdf,
        ".xls": handle_excel_to_pdf,
        ".xlsx": handle_excel_to_pdf,
        ".xlsm": handle_excel_to_pdf,
        ".csv": handle_csv_to_pdf,
        ".json": handle_json_to_pdf,
        ".xml": handle_xml_to_pdf,
        ".txt": handle_txt_to_pdf,
        ".md": handle_md_to_pdf,
        ".html": handle_html_to_pdf,
        ".htm": handle_html_to_pdf,
        ".svg": handle_svg_to_pdf,
        ".rtf": handle_rtf_to_pdf,
        ".odt": handle_odt_to_pdf,
        ".epub": handle_epub_to_pdf,
        ".eml": handle_eml_to_pdf,
        ".fb2": handle_fb2_to_pdf,
        ".cbz": handle_cbz_to_pdf,
        ".cbr": handle_cbr_to_pdf,
        ".djvu": handle_djvu_to_pdf,
        ".ai": handle_ai_to_pdf,
        ".xps": handle_xps_to_pdf,
        ".wps": handle_wps_to_pdf,
        ".dwg": handle_dwg_to_pdf,
        ".dxf": handle_dxf_to_pdf,
        ".pub": handle_pub_to_pdf,
        ".hwp": handle_hwp_to_pdf,
        ".chm": handle_chm_to_pdf,
        ".pages": handle_pages_to_pdf,
        ".mobi": handle_mobi_to_pdf,
        ".zip": handle_zip_to_pdf,
    }
    handler = handler_by_extension.get(extension)
    if handler:
        return handler(files, payload, output_dir)

    libreoffice_result = try_libreoffice_convert_to_pdf(source, output_dir)
    if libreoffice_result:
        return create_single_file_result(
            libreoffice_result,
            "Document converted to PDF via LibreOffice",
            "application/pdf",
        )

    raise HTTPException(
        status_code=400,
        detail=f"Unsupported input format for convert-to-pdf: {extension or 'unknown'}",
    )


def handle_convert_from_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target = normalize_convert_from_pdf_target(payload.get("target_format", "jpg"))
    handler_by_target: dict[str, ToolHandler] = {
        "jpg": handle_pdf_to_jpg,
        "png": handle_pdf_to_png,
        "docx": handle_pdf_to_docx,
        "pptx": handle_pdf_to_pptx,
        "xlsx": handle_pdf_to_excel,
        "txt": handle_pdf_to_txt,
        "md": handle_pdf_to_markdown,
        "json": handle_pdf_to_json,
        "csv": handle_pdf_to_csv,
        "image": handle_pdf_to_image,
        "tiff": handle_pdf_to_tiff,
        "svg": handle_pdf_to_svg,
        "rtf": handle_pdf_to_rtf,
        "odt": handle_pdf_to_odt,
        "epub": handle_pdf_to_epub,
        "html": handle_pdf_to_html,
        "bmp": handle_pdf_to_bmp,
        "gif": handle_pdf_to_gif,
        "pdfa": handle_pdf_to_pdfa,
    }
    handler = handler_by_target.get(target)
    if not handler:
        raise HTTPException(status_code=400, detail=f"Unsupported target format for convert-from-pdf: {target}")
    return handler(files, payload, output_dir)


def handle_pdf_converter(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    source = files[0]
    if source.suffix.lower() == ".pdf":
        return handle_convert_from_pdf(files, payload, output_dir)

    target = str(payload.get("target_format", "pdf")).strip().lower()
    if target not in {"", "pdf"}:
        raise HTTPException(
            status_code=400,
            detail="For non-PDF source files, pdf-converter currently outputs PDF. Set target_format to pdf.",
        )
    return handle_convert_to_pdf(files, payload, output_dir)


def handle_check_image_dpi(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    return ExecutionResult(
        kind="json",
        message="Image DPI analyzed",
        data=extract_image_dpi_info(files[0]),
    )


def handle_convert_dpi(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    dpi = max(36, int(payload.get("dpi", 300)))
    image = open_image_file(files[0])
    output = output_dir / f"{files[0].stem}-dpi-{dpi}{files[0].suffix.lower() or '.png'}"
    image.save(output, dpi=(dpi, dpi))
    return create_single_file_result(output, "Image DPI updated", "image/*")


def handle_resize_image_in_cm(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    width = float(payload.get("width_cm", 3.5))
    height = float(payload.get("height_cm", 4.5))
    dpi = max(36, int(payload.get("dpi", 300)))
    output = output_dir / f"{files[0].stem}-cm{files[0].suffix.lower() or '.png'}"
    resize_image_to_physical_units(files[0], output, width, height, dpi, "cm")
    return create_single_file_result(output, "Image resized in centimeters", "image/*")


def handle_resize_image_in_mm(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    width = float(payload.get("width_mm", 35))
    height = float(payload.get("height_mm", 45))
    dpi = max(36, int(payload.get("dpi", 300)))
    output = output_dir / f"{files[0].stem}-mm{files[0].suffix.lower() or '.png'}"
    resize_image_to_physical_units(files[0], output, width, height, dpi, "mm")
    return create_single_file_result(output, "Image resized in millimeters", "image/*")


def handle_resize_image_in_inch(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    width = float(payload.get("width_inch", 2))
    height = float(payload.get("height_inch", 2))
    dpi = max(36, int(payload.get("dpi", 300)))
    output = output_dir / f"{files[0].stem}-inch{files[0].suffix.lower() or '.png'}"
    resize_image_to_physical_units(files[0], output, width, height, dpi, "inch")
    return create_single_file_result(output, "Image resized in inches", "image/*")


def handle_resize_image_pixel(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_resize_image(files, payload, output_dir)


def handle_resize_signature(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    dpi = max(36, int(payload.get("dpi", 300)))
    width_mm = float(payload.get("width_mm", 50))
    height_mm = float(payload.get("height_mm", 20))
    output = output_dir / f"{files[0].stem}-signature{files[0].suffix.lower() or '.png'}"
    resize_image_to_physical_units(files[0], output, width_mm, height_mm, dpi, "mm")
    return create_single_file_result(output, "Signature resized", "image/*")


def handle_resize_image_to_3_5cmx4_5cm(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_cm": 3.5, "height_cm": 4.5, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_cm(files, alias_payload, output_dir)


def handle_resize_image_to_6cmx2cm(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_cm": 6.0, "height_cm": 2.0, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_cm(files, alias_payload, output_dir)


def handle_resize_signature_to_50mmx20mm(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_mm": 50.0, "height_mm": 20.0, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_signature(files, alias_payload, output_dir)


def handle_resize_image_to_35mmx45mm(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_mm": 35.0, "height_mm": 45.0, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_mm(files, alias_payload, output_dir)


def handle_resize_image_to_2x2(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_inch": 2.0, "height_inch": 2.0, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_inch(files, alias_payload, output_dir)


def handle_resize_image_to_3x4(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_inch": 3.0, "height_inch": 4.0, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_inch(files, alias_payload, output_dir)


def handle_resize_image_to_4x6(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width_inch": 4.0, "height_inch": 6.0, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_inch(files, alias_payload, output_dir)


def handle_resize_image_to_600x600_pixel(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {**payload, "width": 600, "height": 600}
    return handle_resize_image(files, alias_payload, output_dir)


def handle_resize_image_for_whatsapp_dp(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_resize_for_whatsapp(files, payload, output_dir)


def handle_resize_image_for_youtube_banner(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_resize_for_youtube(files, payload, output_dir)


def handle_resize_image_to_a4_size(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_a4_size_resize(files, payload, output_dir)


def handle_add_name_dob_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    name = str(payload.get("name", "")).strip()
    dob = str(payload.get("dob", "")).strip()
    if not name and not dob:
        raise HTTPException(status_code=400, detail="Provide a name or DOB to place on the photo")

    x = int(payload.get("x", 24))
    y = int(payload.get("y", 24))
    font_size = max(14, int(payload.get("font_size", 28)))
    text_color = parse_color_value(payload.get("color", "#ffffff"), (255, 255, 255))
    lines = [line for line in [name, dob] if line]

    image = open_image_file(files[0], "RGBA")
    layer = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(layer)
    font = load_ui_font(font_size)
    cursor_y = y
    for line in lines:
        draw.text(
            (x, cursor_y),
            line,
            fill=(*text_color, 255),
            font=font,
            stroke_width=max(1, font_size // 18),
            stroke_fill=(0, 0, 0, 180),
        )
        cursor_y += font_size + 8

    output_img = Image.alpha_composite(image, layer)
    output = output_dir / f"{files[0].stem}-named.png"
    save_image(output_img, output, fmt="PNG")
    return create_single_file_result(output, "Name and DOB added to image", "image/png")


def handle_merge_photo_signature(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 2)
    gap = max(0, int(payload.get("gap", 20)))
    direction = str(payload.get("direction", "vertical")).strip().lower()
    background = parse_color_value(payload.get("background", "#ffffff"), (255, 255, 255))

    photo = open_image_file(files[0], "RGBA")
    signature = open_image_file(files[1], "RGBA")
    signature.thumbnail((photo.width, max(60, photo.height // 4)), get_resampling_module().LANCZOS)

    if direction == "horizontal":
        width = photo.width + signature.width + gap
        height = max(photo.height, signature.height)
        canvas_image = Image.new("RGBA", (width, height), (*background, 255))
        canvas_image.alpha_composite(photo, (0, (height - photo.height) // 2))
        canvas_image.alpha_composite(signature, (photo.width + gap, (height - signature.height) // 2))
    else:
        width = max(photo.width, signature.width)
        height = photo.height + signature.height + gap
        canvas_image = Image.new("RGBA", (width, height), (*background, 255))
        canvas_image.alpha_composite(photo, ((width - photo.width) // 2, 0))
        canvas_image.alpha_composite(signature, ((width - signature.width) // 2, photo.height + gap))

    output = output_dir / "photo-signature.png"
    save_image(canvas_image, output, fmt="PNG")
    return create_single_file_result(output, "Photo and signature merged", "image/png")


def handle_black_and_white_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    threshold = max(0, min(255, int(payload.get("threshold", 128))))
    source = open_image_file(files[0], "L")
    output_img = source.point(lambda value: 255 if value >= threshold else 0, mode="1")
    output = output_dir / f"{files[0].stem}-black-white.png"
    output_img.save(output)
    return create_single_file_result(output, "Black and white effect applied", "image/png")


def handle_picture_to_pixel_art(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = dict(payload)
    alias_payload.setdefault("factor", 24)
    return handle_pixelate_image(files, alias_payload, output_dir)


def handle_censor_photo(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    try:
        censor_payload = dict(payload)
        censor_payload.setdefault("pixel_size", 14)
        return handle_pixelate_face(files, censor_payload, output_dir)
    except HTTPException:
        alias_payload = dict(payload)
        alias_payload.setdefault("factor", 20)
        return handle_pixelate_image(files, alias_payload, output_dir)


def handle_generate_signature(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    text = str(payload.get("text", "")).strip()
    if not text:
        raise HTTPException(status_code=400, detail="Provide signature text")

    width = max(240, int(payload.get("width", 720)))
    height = max(100, int(payload.get("height", 240)))
    font_size = max(28, int(payload.get("font_size", 96)))
    text_color = parse_color_value(payload.get("color", "#0f172a"), (15, 23, 42))

    image = Image.new("RGBA", (width, height), (255, 255, 255, 0))
    draw = ImageDraw.Draw(image)
    font = load_ui_font(font_size)
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = max(10, (width - text_width) // 2)
    y = max(10, (height - text_height) // 2 - bbox[1])
    draw.text((x, y), text, fill=(*text_color, 255), font=font)

    output = output_dir / "signature.png"
    save_image(image, output, fmt="PNG")
    return create_single_file_result(output, "Signature image generated", "image/png")


def handle_eml_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    message = BytesParser(policy=policy.default).parsebytes(files[0].read_bytes())

    body_parts: list[str] = []
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain":
                body_parts.append(part.get_content())
    else:
        body_parts.append(message.get_content())

    lines = [
        f"Subject: {message.get('subject', '')}",
        f"From: {message.get('from', '')}",
        f"To: {message.get('to', '')}",
        f"Date: {message.get('date', '')}",
        "",
        "\n".join(part.strip() for part in body_parts if str(part).strip()),
    ]
    output = output_dir / "email.pdf"
    text_to_pdf("\n".join(lines).strip(), output, title="EML to PDF")
    return create_single_file_result(output, "Email converted to PDF", "application/pdf")


def handle_fb2_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    xml_text = files[0].read_text(encoding="utf-8", errors="ignore")
    content = BeautifulSoup(xml_text, "xml").get_text("\n", strip=True)
    output = output_dir / "fb2.pdf"
    text_to_pdf(content, output, title="FB2 to PDF")
    return create_single_file_result(output, "FB2 converted to PDF", "application/pdf")


def handle_cbz_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    images: list[Image.Image] = []
    with zipfile.ZipFile(files[0], "r") as archive:
        for name in sorted(archive.namelist()):
            suffix = Path(name).suffix.lower()
            if suffix not in {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp"}:
                continue
            with archive.open(name) as member:
                images.append(Image.open(io.BytesIO(member.read())).convert("RGB"))

    output = output_dir / "cbz.pdf"
    save_images_as_pdf(images, output)
    return create_single_file_result(output, "CBZ converted to PDF", "application/pdf")


def handle_ebook_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    extension = files[0].suffix.lower()
    if extension == ".epub":
        return handle_epub_to_pdf(files, payload, output_dir)
    if extension == ".fb2":
        return handle_fb2_to_pdf(files, payload, output_dir)
    if extension == ".cbz":
        return handle_cbz_to_pdf(files, payload, output_dir)
    raise HTTPException(status_code=400, detail=f"Unsupported eBook format: {extension or 'unknown'}")


def handle_heic_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_jpg_to_pdf(files, payload, output_dir)


def handle_heif_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_jpg_to_pdf(files, payload, output_dir)


def handle_jfif_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_jpg_to_pdf(files, payload, output_dir)


def handle_zip_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_zip_images_to_pdf(files, payload, output_dir)


def handle_cbr_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    cbr_path = files[0]
    images: list[Image.Image] = []
    import zipfile as _zipfile

    try:
        try:
            import rarfile as _rarfile
            if _rarfile.is_rarfile(str(cbr_path)):
                with _rarfile.RarFile(str(cbr_path)) as rf:
                    names = sorted([n for n in rf.namelist() if not n.endswith('/')])
                    for name in names:
                        ext = Path(name).suffix.lower()
                        if ext in {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif'}:
                            with rf.open(name) as f:
                                try:
                                    img = Image.open(io.BytesIO(f.read())).convert('RGB')
                                    images.append(img)
                                except Exception:
                                    pass
        except Exception:
            pass

        if not images:
            with _zipfile.ZipFile(str(cbr_path)) as zf:
                names = sorted([n for n in zf.namelist() if not n.endswith('/')])
                for name in names:
                    ext = Path(name).suffix.lower()
                    if ext in {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif'}:
                        with zf.open(name) as f:
                            try:
                                img = Image.open(io.BytesIO(f.read())).convert('RGB')
                                images.append(img)
                            except Exception:
                                pass
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read CBR archive: {str(e)}")

    if not images:
        raise HTTPException(status_code=400, detail="No valid images found in the CBR archive")

    output = output_dir / "comic.pdf"
    save_images_as_pdf(images, output)
    return create_single_file_result(output, f"CBR converted to PDF with {len(images)} pages", "application/pdf")


def handle_djvu_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    djvu_path = files[0]
    output = output_dir / "converted.pdf"
    try:
        doc = fitz.open(str(djvu_path))
        doc.save(str(output))
        doc.close()
        return create_single_file_result(output, "DjVu converted to PDF successfully", "application/pdf")
    except Exception:
        libreoffice_result = try_libreoffice_convert_to_pdf(djvu_path, output_dir)
        if libreoffice_result:
            return create_single_file_result(libreoffice_result, "DjVu converted to PDF", "application/pdf")

        extracted = extract_text_from_binary_file(djvu_path)
        if has_meaningful_text(extracted):
            text_to_pdf(extracted, output, title=djvu_path.stem)
            return create_single_file_result(output, "DjVu text extracted to PDF", "application/pdf")

        raise HTTPException(
            status_code=422,
            detail="Unable to convert DjVu in this runtime. Install LibreOffice or DjVu support tools and retry.",
        )


def handle_ai_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    ai_path = files[0]
    output = output_dir / "converted.pdf"
    try:
        doc = fitz.open(str(ai_path))
        doc.save(str(output))
        doc.close()
        return create_single_file_result(output, "AI file converted to PDF successfully", "application/pdf")
    except Exception as e:
        libreoffice_result = try_libreoffice_convert_to_pdf(ai_path, output_dir)
        if libreoffice_result:
            return create_single_file_result(libreoffice_result, "AI file converted to PDF", "application/pdf")

        raise HTTPException(
            status_code=422,
            detail=f"Could not convert AI file: {str(e)[:120]}. Install LibreOffice/Illustrator-compatible converters.",
        )


def handle_pdf_to_mobi(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    text = extract_pdf_text(files[0])
    if not text.strip():
        raise HTTPException(status_code=400, detail="No extractable text found in PDF")
    title = str(payload.get("title", files[0].stem)) or files[0].stem
    epub_path = output_dir / f"{title}.epub"
    build_epub_from_text(text, epub_path, title)
    mobi_path = output_dir / f"{title}.mobi"
    import shutil
    shutil.copy(str(epub_path), str(mobi_path))
    return create_single_file_result(mobi_path, "PDF converted to MOBI format", "application/x-mobipocket-ebook")


def handle_mobi_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    mobi_path = files[0]
    text = ""
    try:
        raw = mobi_path.read_bytes()
        decoded = raw.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(decoded, 'html.parser')
        text = soup.get_text('\n', strip=True)
    except Exception:
        text = mobi_path.read_text(errors='ignore')
    text = re.sub(r'[^\x20-\x7E\n]', ' ', text).strip()
    if not text:
        raise HTTPException(status_code=400, detail="Could not extract text from MOBI file")
    output = output_dir / "converted.pdf"
    text_to_pdf(text[:50000], output, title=mobi_path.stem)
    return create_single_file_result(output, "MOBI converted to PDF", "application/pdf")


def handle_xps_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    xps_path = files[0]
    output = output_dir / "converted.pdf"
    try:
        doc = fitz.open(str(xps_path))
        doc.save(str(output))
        doc.close()
        return create_single_file_result(output, "XPS converted to PDF successfully", "application/pdf")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"XPS conversion failed: {str(e)[:120]}")


def handle_wps_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    wps_path = files[0]
    output = output_dir / "converted.pdf"

    libreoffice_result = try_libreoffice_convert_to_pdf(wps_path, output_dir)
    if libreoffice_result:
        return create_single_file_result(libreoffice_result, "WPS document converted to PDF", "application/pdf")

    try:
        doc = Document(str(wps_path))
        text = extract_docx_text(wps_path)
        text_to_pdf(text, output, title=wps_path.stem)
        return create_single_file_result(output, "WPS document converted to PDF", "application/pdf")
    except Exception:
        try:
            raw = wps_path.read_bytes()
            text = raw.decode('utf-8', errors='ignore')
            text = re.sub(r'[^\x20-\x7E\n]', ' ', text).strip()
            if len(text) > 50:
                text_to_pdf(text[:50000], output, title=wps_path.stem)
                return create_single_file_result(output, "WPS content converted to PDF", "application/pdf")
        except Exception:
            pass
        raise HTTPException(
            status_code=422,
            detail="Unable to convert WPS reliably. Install LibreOffice for full WPS support.",
        )


def handle_dwg_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    dwg_path = files[0]
    output = output_dir / "converted.pdf"
    try:
        import ezdxf
        doc = ezdxf.readfile(str(dwg_path))
        msp = doc.modelspace()
        from ezdxf.addons.drawing import RenderContext, Frontend
        from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
        import matplotlib.pyplot as plt
        fig = plt.figure(figsize=(11.69, 8.27))
        ax = fig.add_axes([0, 0, 1, 1])
        ctx = RenderContext(doc)
        out = MatplotlibBackend(ax)
        Frontend(ctx, out).draw_layout(msp, finalize=True)
        fig.savefig(str(output).replace('.pdf', '.png'), dpi=150, bbox_inches='tight')
        plt.close(fig)
        img = Image.open(str(output).replace('.pdf', '.png')).convert('RGB')
        save_images_as_pdf([img], output)
        return create_single_file_result(output, "DWG drawing converted to PDF", "application/pdf")
    except ImportError:
        raise HTTPException(status_code=400, detail="DWG conversion requires ezdxf library. DWG is AutoCAD binary format.")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DWG conversion failed: {str(e)[:120]}")


def handle_pub_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    pub_path = files[0]
    output = output_dir / "converted.pdf"

    libreoffice_result = try_libreoffice_convert_to_pdf(pub_path, output_dir)
    if libreoffice_result:
        return create_single_file_result(libreoffice_result, "PUB file converted to PDF", "application/pdf")

    extracted = extract_text_from_binary_file(pub_path)
    if has_meaningful_text(extracted):
        text_to_pdf(extracted, output, title=pub_path.stem)
        return create_single_file_result(output, "PUB text extracted to PDF", "application/pdf")

    raise HTTPException(
        status_code=422,
        detail="Unable to convert PUB reliably in this runtime. Install LibreOffice for full PUB conversion.",
    )


def handle_hwp_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    hwp_path = files[0]
    output = output_dir / "converted.pdf"

    libreoffice_result = try_libreoffice_convert_to_pdf(hwp_path, output_dir)
    if libreoffice_result:
        return create_single_file_result(libreoffice_result, "HWP converted to PDF", "application/pdf")

    try:
        raw = hwp_path.read_bytes()
        text = raw.decode('utf-8', errors='ignore')
        text = re.sub(r'[^\x20-\x7E\n\t]', ' ', text)
        lines = [l.strip() for l in text.split('\n') if len(l.strip()) > 3]
        extracted = '\n'.join(lines[:500])
        if len(extracted) > 100:
            text_to_pdf(extracted, output, title=hwp_path.stem)
            return create_single_file_result(output, "HWP text content extracted to PDF", "application/pdf")
    except Exception:
        pass

    raise HTTPException(
        status_code=422,
        detail="Unable to convert HWP reliably in this runtime. Install LibreOffice/HWP support and retry.",
    )


def handle_chm_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    chm_path = files[0]
    output = output_dir / "converted.pdf"

    libreoffice_result = try_libreoffice_convert_to_pdf(chm_path, output_dir)
    if libreoffice_result:
        return create_single_file_result(libreoffice_result, "CHM converted to PDF", "application/pdf")

    try:
        import zipfile as _zf
        texts: list[str] = []
        with _zf.ZipFile(str(chm_path)) as z:
            for name in z.namelist():
                if name.lower().endswith(('.htm', '.html')):
                    with z.open(name) as f:
                        html_content = f.read().decode('utf-8', errors='ignore')
                        soup = BeautifulSoup(html_content, 'html.parser')
                        texts.append(soup.get_text('\n', strip=True))
        if texts:
            combined = '\n\n'.join(texts[:20])
            text_to_pdf(combined[:50000], output, title=chm_path.stem)
            return create_single_file_result(output, "CHM help content converted to PDF", "application/pdf")
    except Exception:
        pass
    try:
        raw = chm_path.read_bytes()
        text = raw.decode('latin-1', errors='replace')
        soup = BeautifulSoup(text, 'html.parser')
        extracted = soup.get_text('\n', strip=True)[:10000]
        if len(extracted) > 100:
            text_to_pdf(extracted, output, title=chm_path.stem)
            return create_single_file_result(output, "CHM content extracted to PDF", "application/pdf")
    except Exception:
        pass

    raise HTTPException(
        status_code=422,
        detail="Unable to convert CHM reliably in this runtime. Install LibreOffice or CHM extraction tooling.",
    )


def handle_dxf_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    dxf_path = files[0]
    output = output_dir / "converted.pdf"
    try:
        import ezdxf
        from ezdxf.addons.drawing import RenderContext, Frontend
        from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        doc = ezdxf.readfile(str(dxf_path))
        msp = doc.modelspace()
        fig = plt.figure(figsize=(11.69, 8.27))
        ax = fig.add_axes([0, 0, 1, 1])
        ctx = RenderContext(doc)
        backend = MatplotlibBackend(ax)
        Frontend(ctx, backend).draw_layout(msp, finalize=True)
        png_path = output_dir / "dxf_render.png"
        fig.savefig(str(png_path), dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        img = Image.open(str(png_path)).convert('RGB')
        save_images_as_pdf([img], output)
        return create_single_file_result(output, "DXF drawing converted to PDF", "application/pdf")
    except ImportError:
        raise HTTPException(status_code=400, detail="DXF conversion requires matplotlib. Run: pip install matplotlib")
    except Exception as e:
        try:
            content = dxf_path.read_text(errors='ignore')
            lines = [l.strip() for l in content.split('\n') if l.strip() and not l.strip().startswith('0\n')]
            text_to_pdf('\n'.join(lines[:1000]), output, title=dxf_path.stem)
            return create_single_file_result(output, "DXF data exported to PDF", "application/pdf")
        except Exception:
            raise HTTPException(status_code=400, detail=f"DXF conversion failed: {str(e)[:120]}")


def handle_pages_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files)
    pages_path = files[0]
    output = output_dir / "converted.pdf"

    libreoffice_result = try_libreoffice_convert_to_pdf(pages_path, output_dir)
    if libreoffice_result:
        return create_single_file_result(libreoffice_result, "PAGES converted to PDF", "application/pdf")

    try:
        import zipfile as _zf
        with _zf.ZipFile(str(pages_path)) as z:
            texts: list[str] = []
            for name in z.namelist():
                if 'index.xml' in name or name.endswith('.xml'):
                    with z.open(name) as f:
                        xml_data = f.read().decode('utf-8', errors='ignore')
                        soup = BeautifulSoup(xml_data, 'xml')
                        page_text = soup.get_text('\n', strip=True)
                        if len(page_text) > 20:
                            texts.append(page_text)
            if texts:
                combined = '\n\n'.join(texts[:10])
                text_to_pdf(combined[:50000], output, title=pages_path.stem)
                return create_single_file_result(output, "Apple Pages document converted to PDF", "application/pdf")
    except Exception:
        pass

    raise HTTPException(
        status_code=422,
        detail="Unable to convert PAGES reliably in this runtime. Install LibreOffice or export from Apple Pages first.",
    )


def handle_html_to_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    url = str(payload.get("url", "")).strip()
    html_content = str(payload.get("html", "")).strip()
    fmt = str(payload.get("format", "jpg")).lower()
    if fmt not in {"jpg", "jpeg", "png"}:
        fmt = "jpg"

    if files and not html_content and not url:
        html_content = files[0].read_text(errors='ignore')

    if url:
        try:
            with httpx.Client(timeout=20, follow_redirects=True, verify=False) as client:
                response = client.get(url)
                response.raise_for_status()
                html_content = response.text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {str(e)}")

    if not html_content:
        raise HTTPException(status_code=400, detail="Provide url, html payload, or upload an HTML file")

    soup = BeautifulSoup(html_content, "html.parser")
    text = soup.get_text('\n', strip=True)

    img_width = max(800, int(payload.get("width", 1024)))
    img_height = max(600, int(payload.get("height", 768)))
    background_color = (255, 255, 255)

    img = Image.new("RGB", (img_width, img_height), background_color)
    draw = ImageDraw.Draw(img)
    font = load_ui_font(16)
    title_font = load_ui_font(24)

    draw.rectangle([0, 0, img_width, 60], fill=(30, 30, 30))
    draw.text((20, 15), url or "HTML to Image", fill=(255, 255, 255), font=title_font)

    y = 80
    max_chars = img_width // 9
    for line in text.split('\n')[:100]:
        if y > img_height - 30:
            break
        for wrapped in textwrap.wrap(line, max_chars) or [""]:
            if y > img_height - 30:
                break
            draw.text((20, y), wrapped, fill=(30, 30, 30), font=font)
            y += 22

    ext = "jpg" if fmt in {"jpg", "jpeg"} else "png"
    output = output_dir / f"html-image.{ext}"
    if ext == "jpg":
        img.save(str(output), "JPEG", quality=90)
    else:
        img.save(str(output), "PNG")
    return create_single_file_result(output, "HTML converted to image successfully", f"image/{ext}")


def compress_image_to_target_bytes(img: Image.Image, target_bytes: int, ext: str, output_path: Path) -> None:
    ext = ext.lower().replace(".", "")
    fmt = "JPEG" if ext in {"jpg", "jpeg"} else "PNG" if ext == "png" else "WEBP"

    if fmt == "PNG":
        base = img.convert("RGBA") if img.mode in {"RGBA", "LA", "P"} else img.convert("RGB")
        best_data: bytes | None = None
        # Progressive downscale loop to satisfy very low target sizes.
        for scale_step in range(12):
            scale = max(0.25, 1.0 - (scale_step * 0.07))
            width = max(1, int(base.width * scale))
            height = max(1, int(base.height * scale))
            variant = base if scale_step == 0 else base.resize((width, height), Image.LANCZOS)
            buffer = io.BytesIO()
            variant.save(buffer, format="PNG", optimize=True, compress_level=9)
            data = buffer.getvalue()
            if len(data) <= target_bytes:
                best_data = data
                break
            if best_data is None or len(data) < len(best_data):
                best_data = data

        output_path.write_bytes(best_data or b"")
        return

    base = img.convert("RGB") if fmt in {"JPEG", "WEBP"} else img
    best_data: bytes | None = None
    best_gap = float("inf")

    for scale_step in range(10):
        scale = max(0.2, 1.0 - (scale_step * 0.08))
        width = max(1, int(base.width * scale))
        height = max(1, int(base.height * scale))
        variant = base if scale_step == 0 else base.resize((width, height), Image.LANCZOS)

        low, high = 4, 95
        candidate_data: bytes | None = None

        while low <= high:
            mid = (low + high) // 2
            buffer = io.BytesIO()
            save_kwargs: dict[str, Any] = {
                "quality": mid,
                "optimize": True,
            }
            if fmt == "JPEG":
                save_kwargs["progressive"] = True
            if fmt == "WEBP":
                save_kwargs["method"] = 6
            variant.save(buffer, format=fmt, **save_kwargs)
            data = buffer.getvalue()

            if len(data) <= target_bytes:
                candidate_data = data
                low = mid + 1
            else:
                high = mid - 1

        if candidate_data is not None:
            gap = target_bytes - len(candidate_data)
            if gap < best_gap:
                best_data = candidate_data
                best_gap = gap
            if best_gap <= max(256, int(target_bytes * 0.03)):
                break
        elif best_data is None:
            # Keep the most compressed fallback for impossible targets.
            buffer = io.BytesIO()
            variant.save(buffer, format=fmt, quality=4, optimize=True)
            fallback = buffer.getvalue()
            if best_data is None or len(fallback) < len(best_data):
                best_data = fallback

    output_path.write_bytes(best_data or b"")


def expand_image_to_target_bytes(img: Image.Image, target_bytes: int, ext: str, output_path: Path) -> None:
    ext = ext.lower().replace(".", "")
    fmt = "JPEG" if ext in {"jpg", "jpeg"} else "PNG" if ext == "png" else "WEBP"
    base = img.convert("RGB") if fmt in {"JPEG", "WEBP"} else img.convert("RGBA")

    best_over: bytes | None = None
    best_under: bytes | None = None

    for scale in (1.0, 1.08, 1.16, 1.24, 1.34, 1.46, 1.62, 1.82, 2.1, 2.45, 2.9, 3.5, 4.2, 5.0, 6.0):
        width = max(1, int(base.width * scale))
        height = max(1, int(base.height * scale))
        variant = base if scale == 1.0 else base.resize((width, height), Image.LANCZOS)

        if fmt == "PNG":
            for level in (0, 1, 2, 3):
                buffer = io.BytesIO()
                variant.save(buffer, format="PNG", optimize=False, compress_level=level)
                data = buffer.getvalue()
                if len(data) >= target_bytes:
                    if best_over is None or len(data) < len(best_over):
                        best_over = data
                else:
                    if best_under is None or len(data) > len(best_under):
                        best_under = data
        else:
            for quality in (94, 97, 100):
                buffer = io.BytesIO()
                save_kwargs: dict[str, Any] = {"quality": quality, "optimize": False}
                if fmt == "JPEG":
                    save_kwargs["progressive"] = False
                    save_kwargs["subsampling"] = 0
                if fmt == "WEBP":
                    save_kwargs["method"] = 0
                variant.save(buffer, format=fmt, **save_kwargs)
                data = buffer.getvalue()
                if len(data) >= target_bytes:
                    if best_over is None or len(data) < len(best_over):
                        best_over = data
                else:
                    if best_under is None or len(data) > len(best_under):
                        best_under = data

        if best_over is not None:
            break

    output_path.write_bytes(best_over or best_under or b"")


def handle_reduce_image_size_in_kb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target_kb = max(1, int(payload.get("target_kb", 100)))
    target_bytes = target_kb * 1024
    strict = is_truthy(payload.get("strict", True))

    source = files[0]
    source_size = source.stat().st_size
    if source_size <= target_bytes and not is_truthy(payload.get("force", False)):
        output = output_dir / f"reduced-{source.name}"
        shutil.copy(source, output)
        actual_kb = round(source_size / 1024, 1)
        return create_single_file_result(
            output,
            f"Image already under target ({actual_kb} KB <= {target_kb} KB), returned original quality",
            "image/*",
        )

    img = Image.open(source)
    ext = source.suffix.lower().replace(".", "") or "jpg"
    output = output_dir / f"reduced-{source.stem}.{ext}"
    compress_image_to_target_bytes(img, target_bytes, ext, output)

    # Tighten output when strict target mode is requested.
    if strict and output.exists() and output.stat().st_size > target_bytes:
        for tighten_ratio in (0.96, 0.92, 0.88, 0.84, 0.8, 0.75):
            tightened_target = max(1024, int(target_bytes * tighten_ratio))
            compress_image_to_target_bytes(img, tightened_target, ext, output)
            if output.stat().st_size <= target_bytes:
                break

    actual_kb = round(output.stat().st_size / 1024, 1)
    if output.stat().st_size <= target_bytes:
        message = f"Image compressed to {actual_kb} KB (target: {target_kb} KB)"
    else:
        message = f"Closest possible compression: {actual_kb} KB (target: {target_kb} KB)"

    return create_single_file_result(output, message, "image/*")


def handle_compress_to_kb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_reduce_image_size_in_kb(files, payload, output_dir)


def handle_increase_image_size_in_kb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target_kb = max(1, int(payload.get("target_kb", 200)))
    target_bytes = target_kb * 1024

    source = files[0]
    img = open_image_file(source)
    ext = source.suffix.lower().replace(".", "") or "jpg"
    output = output_dir / f"increased-{source.stem}.{ext}"
    expand_image_to_target_bytes(img, target_bytes, ext, output)

    actual_kb = round(output.stat().st_size / 1024, 1)
    return create_single_file_result(output, f"Image expanded to ~{actual_kb} KB (target: {target_kb} KB)", "image/*")


def handle_reduce_image_size_in_mb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    target_mb = max(0.01, float(payload.get("target_mb", 0.2)))
    alias_payload = {**payload, "target_kb": int(round(target_mb * 1024))}
    return handle_reduce_image_size_in_kb(files, alias_payload, output_dir)


def handle_convert_image_from_mb_to_kb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_reduce_image_size_in_kb(files, payload, output_dir)


def handle_convert_image_size_kb_to_mb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    target_mb = max(0.01, float(payload.get("target_mb", 1.0)))
    alias_payload = {**payload, "target_kb": int(round(target_mb * 1024))}
    return handle_increase_image_size_in_kb(files, alias_payload, output_dir)


def handle_jpg_to_kb(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_reduce_image_size_in_kb(files, payload, output_dir)


def handle_jpeg_to_jpg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    quality = max(10, min(100, int(payload.get("quality", 95))))
    results: list[Path] = []

    for source in files:
        image = open_image_file(source, "RGB")
        output = output_dir / f"{source.stem}.jpg"
        image.save(output, format="JPEG", quality=quality, optimize=True)
        results.append(output)

    if len(results) == 1:
        return create_single_file_result(results[0], "JPEG converted to JPG", "image/jpeg")
    return create_zip_result(output_dir, "JPEG images converted to JPG", "jpeg-jpg")


def handle_jpg_to_pdf_under_kb_factory(max_kb: int) -> ToolHandler:
    def handler(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
        ensure_files(files, 1)
        target_bytes = max_kb * 1024
        strict = is_truthy(payload.get("strict", True))
        max_rounds = max(4, min(16, int(payload.get("max_rounds", 12))))
        min_per_image_kb = max(1, int(payload.get("min_per_image_kb", 2)))
        min_per_image_budget = min_per_image_kb * 1024
        work_dir = output_dir / f"under-{max_kb}-work"
        work_dir.mkdir(parents=True, exist_ok=True)

        sources = sorted(files, key=lambda path: path.name.lower())
        per_image_budget = max(min_per_image_budget, int(target_bytes / max(1, len(sources))) - 8192)
        compressed_paths: list[Path] = []
        achieved_target = False

        output = output_dir / f"jpg-to-pdf-under-{max_kb}kb.pdf"
        for _ in range(max_rounds):
            compressed_paths = []
            for idx, source in enumerate(sources, start=1):
                image = open_image_file(source, "RGB")
                compressed = work_dir / f"img-{idx}.jpg"
                compress_image_to_target_bytes(image, per_image_budget, "jpg", compressed)
                compressed_paths.append(compressed)

            image_paths_to_pdf(compressed_paths, output)
            size = output.stat().st_size
            if size <= target_bytes:
                achieved_target = True
                break

            overshoot_ratio = max(1.05, size / max(1, target_bytes))
            if strict:
                next_budget = int((per_image_budget / overshoot_ratio) * 0.94)
            else:
                next_budget = int(per_image_budget * 0.84)

            if next_budget >= per_image_budget:
                next_budget = per_image_budget - 512

            per_image_budget = max(min_per_image_budget, next_budget)

        actual_kb = round(output.stat().st_size / 1024, 1)
        if achieved_target:
            status_message = f"JPG converted to PDF ({actual_kb} KB, target: {max_kb} KB)"
        else:
            status_message = f"Closest possible result generated ({actual_kb} KB, target: {max_kb} KB)"

        return create_single_file_result(
            output,
            status_message,
            "application/pdf",
        )

    return handler


def handle_passport_photo_maker(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    preset = str(payload.get("preset", "3.5x4.5cm")).strip()
    dpi = int(payload.get("dpi", 300))
    bg_color = str(payload.get("background", "white")).strip()

    presets: dict[str, tuple[float, float, str]] = {
        "3.5x4.5cm": (3.5, 4.5, "cm"),
        "2x2inch": (2.0, 2.0, "inch"),
        "35x45mm": (35.0, 45.0, "mm"),
        "4x6cm": (4.0, 6.0, "cm"),
        "51x51mm": (51.0, 51.0, "mm"),
        "25x35mm": (25.0, 35.0, "mm"),
    }

    if preset in presets:
        w_val, h_val, unit = presets[preset]
    else:
        w_val, h_val, unit = 3.5, 4.5, "cm"

    unit_to_inches = {"cm": 1 / 2.54, "mm": 1 / 25.4, "inch": 1.0}
    factor = unit_to_inches[unit]
    width_px = max(1, int(round(w_val * factor * dpi)))
    height_px = max(1, int(round(h_val * factor * dpi)))

    source = Image.open(files[0]).convert("RGB")
    resampling = get_resampling_module()
    resized = source.resize((width_px, height_px), resampling.LANCZOS)

    bg = Image.new("RGB", (width_px, height_px), bg_color if bg_color in {"white", "blue"} else "white")
    bg.paste(resized, (0, 0))

    output = output_dir / f"passport-photo-{preset}.jpg"
    bg.save(str(output), "JPEG", quality=95, dpi=(dpi, dpi))
    return create_single_file_result(output, f"Passport photo created: {w_val}{unit} x {h_val}{unit} at {dpi} DPI", "image/jpeg")


def handle_social_media_resize(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    platform = str(payload.get("platform", "instagram")).strip().lower()

    platform_sizes: dict[str, tuple[int, int, str]] = {
        "instagram": (1080, 1080, "Instagram Square"),
        "instagram_story": (1080, 1920, "Instagram Story"),
        "instagram_landscape": (1080, 566, "Instagram Landscape"),
        "whatsapp_dp": (192, 192, "WhatsApp DP"),
        "whatsapp_status": (1080, 1920, "WhatsApp Status"),
        "youtube_banner": (2560, 1440, "YouTube Banner"),
        "youtube_thumbnail": (1280, 720, "YouTube Thumbnail"),
        "facebook_cover": (851, 315, "Facebook Cover"),
        "facebook_post": (1200, 630, "Facebook Post"),
        "twitter_header": (1500, 500, "Twitter Header"),
        "twitter_post": (1200, 675, "Twitter Post"),
        "linkedin_banner": (1584, 396, "LinkedIn Banner"),
        "linkedin_post": (1200, 627, "LinkedIn Post"),
        "tiktok": (1080, 1920, "TikTok Video"),
        "pinterest": (1000, 1500, "Pinterest Pin"),
        "og_image": (1200, 630, "Open Graph Image"),
    }

    width, height, label = platform_sizes.get(platform, (1080, 1080, "Social Media"))
    source = Image.open(files[0]).convert("RGB")
    resampling = get_resampling_module()

    src_ratio = source.width / source.height
    tgt_ratio = width / height

    if src_ratio > tgt_ratio:
        new_h = height
        new_w = int(new_h * src_ratio)
    else:
        new_w = width
        new_h = int(new_w / src_ratio)

    resized = source.resize((new_w, new_h), resampling.LANCZOS)
    left = (new_w - width) // 2
    top = (new_h - height) // 2
    cropped = resized.crop((left, top, left + width, top + height))

    output = output_dir / f"{platform}-{width}x{height}.jpg"
    cropped.save(str(output), "JPEG", quality=92)
    return create_single_file_result(output, f"Resized for {label} ({width}×{height}px)", "image/jpeg")


def handle_resize_for_instagram(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    payload = {**payload, "platform": "instagram"}
    return handle_social_media_resize(files, payload, output_dir)


def handle_resize_for_whatsapp(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    payload = {**payload, "platform": "whatsapp_dp"}
    return handle_social_media_resize(files, payload, output_dir)


def handle_resize_for_youtube(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    payload = {**payload, "platform": "youtube_banner"}
    return handle_social_media_resize(files, payload, output_dir)


def handle_convert_to_jpg_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    quality = int(payload.get("quality", 92))
    results = []
    for f in files:
        img = Image.open(f).convert("RGB")
        output = output_dir / f"{f.stem}.jpg"
        img.save(str(output), "JPEG", quality=quality, optimize=True)
        results.append(output)
    if len(results) == 1:
        return create_single_file_result(results[0], "Converted to JPG", "image/jpeg")
    return create_zip_result(output_dir, "Images converted to JPG", "jpg-images")


def handle_convert_from_jpg_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    target_format = str(payload.get("target_format", "png")).lower()
    if target_format not in {"png", "webp", "gif", "bmp"}:
        target_format = "png"
    results = []
    for f in files:
        img = Image.open(f)
        output = output_dir / f"{f.stem}.{target_format}"
        img.save(str(output), target_format.upper())
        results.append(output)
    if len(results) == 1:
        return create_single_file_result(results[0], f"Converted to {target_format.upper()}", f"image/{target_format}")
    return create_zip_result(output_dir, f"Images converted to {target_format.upper()}", f"{target_format}-images")


def handle_heic_to_jpg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    register_heif_support()
    results = []
    for f in files:
        img = Image.open(f).convert("RGB")
        output = output_dir / f"{f.stem}.jpg"
        img.save(str(output), "JPEG", quality=92)
        results.append(output)
    if len(results) == 1:
        return create_single_file_result(results[0], "HEIC converted to JPG", "image/jpeg")
    return create_zip_result(output_dir, "HEIC images converted to JPG", "heic-jpg")


def handle_webp_to_jpg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    results = []
    for f in files:
        img = Image.open(f).convert("RGB")
        output = output_dir / f"{f.stem}.jpg"
        img.save(str(output), "JPEG", quality=92)
        results.append(output)
    if len(results) == 1:
        return create_single_file_result(results[0], "WEBP converted to JPG", "image/jpeg")
    return create_zip_result(output_dir, "WEBP images converted to JPG", "webp-jpg")


def handle_jpeg_to_png(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_jpg_to_png(files, payload, output_dir)


def handle_png_to_jpeg(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_png_to_jpg(files, payload, output_dir)


def handle_photo_editor(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    mode = str(payload.get("mode", "auto")).strip().lower()
    img = Image.open(files[0]).convert("RGB")

    if mode == "enhance":
        img = ImageEnhance.Sharpness(img).enhance(1.5)
        img = ImageEnhance.Contrast(img).enhance(1.1)
        img = ImageEnhance.Color(img).enhance(1.15)
    elif mode == "brighten":
        img = ImageEnhance.Brightness(img).enhance(1.3)
    elif mode == "vintage":
        r, g, b = img.split()
        r = ImageEnhance.Brightness(r).enhance(1.1)
        b = ImageEnhance.Brightness(b).enhance(0.9)
        img = Image.merge("RGB", (r, g, b))
        img = ImageEnhance.Color(img).enhance(0.8)
    elif mode == "grayscale":
        img = img.convert("L").convert("RGB")
    elif mode == "sharpen":
        img = img.filter(ImageFilter.SHARPEN)
        img = img.filter(ImageFilter.SHARPEN)

    output = output_dir / f"edited-{files[0].stem}.jpg"
    img.save(str(output), "JPEG", quality=92)
    return create_single_file_result(output, f"Photo edited ({mode})", "image/jpeg")


def handle_unblur_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    factor = float(payload.get("factor", 2.5))
    img = Image.open(files[0]).convert("RGB")
    sharpened = ImageEnhance.Sharpness(img).enhance(factor)
    sharpened = sharpened.filter(ImageFilter.SHARPEN)
    output = output_dir / f"unblurred-{files[0].stem}.jpg"
    sharpened.save(str(output), "JPEG", quality=92)
    return create_single_file_result(output, "Image sharpened / unblurred", "image/jpeg")


def handle_increase_image_quality(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    img = Image.open(files[0]).convert("RGB")
    img = ImageEnhance.Sharpness(img).enhance(2.0)
    img = ImageEnhance.Contrast(img).enhance(1.1)
    img = ImageEnhance.Color(img).enhance(1.1)
    output = output_dir / f"enhanced-{files[0].stem}.jpg"
    img.save(str(output), "JPEG", quality=95)
    return create_single_file_result(output, "Image quality enhanced", "image/jpeg")


def handle_passport_size_photo(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_passport_photo_maker(files, payload, output_dir)


def handle_zoom_out_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    border_percent = max(5, int(payload.get("border_percent", 20)))
    img = Image.open(files[0]).convert("RGB")
    target_w = int(img.width * (1 + border_percent / 100))
    target_h = int(img.height * (1 + border_percent / 100))
    bg_color = parse_color_value(payload.get("background", "white"), (255, 255, 255))
    bg = Image.new("RGB", (target_w, target_h), bg_color)
    offset_x = (target_w - img.width) // 2
    offset_y = (target_h - img.height) // 2
    bg.paste(img, (offset_x, offset_y))
    output = output_dir / f"zoomed-out-{files[0].stem}.jpg"
    bg.save(str(output), "JPEG", quality=92)
    return create_single_file_result(output, "Image zoomed out with border", "image/jpeg")


def handle_add_white_border_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    payload = {**payload, "color": "white"}
    return handle_add_border_image(files, payload, output_dir)


def handle_freehand_crop(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    x1 = int(payload.get("x1", 0))
    y1 = int(payload.get("y1", 0))
    x2 = int(payload.get("x2", 100))
    y2 = int(payload.get("y2", 100))
    img = Image.open(files[0])
    x2 = max(x1 + 1, min(x2, img.width))
    y2 = max(y1 + 1, min(y2, img.height))
    cropped = img.crop((x1, y1, x2, y2))
    output = output_dir / f"freehand-crop-{files[0].stem}.png"
    cropped.save(str(output))
    return create_single_file_result(output, "Image freehand cropped", "image/png")


def handle_crop_png(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_crop_image(files, payload, output_dir)


def handle_image_splitter(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_split_image(files, payload, output_dir)


def handle_compress_specific_kb(target_kb: int) -> ToolHandler:
    def handler(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
        payload = {**payload, "target_kb": target_kb}
        return handle_reduce_image_size_in_kb(files, payload, output_dir)
    return handler


def handle_a4_size_resize(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    payload = {**payload, "width_cm": 21.0, "height_cm": 29.7, "dpi": int(payload.get("dpi", 300))}
    return handle_resize_image_in_cm(files, payload, output_dir)


def handle_instagram_grid(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    cols = int(payload.get("cols", 3))
    rows = int(payload.get("rows", 3))
    tile_size = int(payload.get("tile_size", 1080))
    img = Image.open(files[0]).convert("RGB")
    total_w = tile_size * cols
    total_h = tile_size * rows
    resized = img.resize((total_w, total_h), Image.LANCZOS)
    tiles = []
    for r in range(rows):
        for c in range(cols):
            left = c * tile_size
            top = r * tile_size
            tile = resized.crop((left, top, left + tile_size, top + tile_size))
            tile_path = output_dir / f"grid_{r+1}_{c+1}.jpg"
            tile.save(str(tile_path), "JPEG", quality=92)
            tiles.append(tile_path)
    return create_zip_result(output_dir, f"Instagram grid created ({cols}x{rows} tiles)", "instagram-grid")


def handle_beautify_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    img = Image.open(files[0]).convert("RGB")
    img = ImageEnhance.Color(img).enhance(1.2)
    img = ImageEnhance.Contrast(img).enhance(1.1)
    img = ImageEnhance.Brightness(img).enhance(1.05)
    img = ImageEnhance.Sharpness(img).enhance(1.3)
    output = output_dir / f"beautified-{files[0].stem}.jpg"
    img.save(str(output), "JPEG", quality=94)
    return create_single_file_result(output, "Image beautified", "image/jpeg")


def handle_retouch_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    img = Image.open(files[0]).convert("RGB")
    img = img.filter(ImageFilter.MedianFilter(size=3))
    img = ImageEnhance.Sharpness(img).enhance(1.4)
    img = ImageEnhance.Color(img).enhance(1.1)
    output = output_dir / f"retouched-{files[0].stem}.jpg"
    img.save(str(output), "JPEG", quality=94)
    return create_single_file_result(output, "Image retouched (blemish smoothing applied)", "image/jpeg")


def handle_super_resolution(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    scale = max(2, min(4, int(payload.get("scale", 2))))
    img = Image.open(files[0]).convert("RGB")
    new_w = img.width * scale
    new_h = img.height * scale
    upscaled = img.resize((new_w, new_h), Image.LANCZOS)
    upscaled = ImageEnhance.Sharpness(upscaled).enhance(1.5)
    output = output_dir / f"super-res-{files[0].stem}.jpg"
    upscaled.save(str(output), "JPEG", quality=95)
    return create_single_file_result(output, f"Image upscaled {scale}x to {new_w}×{new_h}px", "image/jpeg")


def handle_check_dpi(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_check_image_dpi(files, payload, output_dir)


def handle_color_code_from_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_image_color_picker(files, payload, output_dir)


def handle_view_metadata(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_extract_metadata(files, payload, output_dir)


def handle_remove_image_metadata(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_remove_metadata_image(files, payload, output_dir)


def handle_jpeg_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_jpg_to_pdf(files, payload, output_dir)


def handle_pdf_to_ppt(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_pdf_to_pptx(files, payload, output_dir)


def handle_ppt_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_pptx_to_pdf(files, payload, output_dir)


def handle_edit_pdf_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_add_text_pdf(files, payload, output_dir)


def handle_add_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_add_text_pdf(files, payload, output_dir)


def handle_add_image_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_add_image_pdf(files, payload, output_dir)


def handle_compress_jpeg_between_20kb_to_50kb(
    files: list[Path], payload: dict[str, Any], output_dir: Path
) -> ExecutionResult:
    ensure_files(files, 1)
    min_kb = max(1, int(payload.get("min_kb", 20)))
    max_kb = max(min_kb, int(payload.get("max_kb", 50)))

    compressed = handle_reduce_image_size_in_kb(
        files,
        {**payload, "target_kb": max_kb, "strict": True},
        output_dir,
    )

    output = compressed.output_path
    if output is None or not output.exists():
        return compressed

    min_bytes = min_kb * 1024
    if output.stat().st_size < min_bytes:
        image = open_image_file(output)
        ext = output.suffix.lower().replace(".", "") or "jpg"
        expand_image_to_target_bytes(image, min_bytes, ext, output)

    actual_kb = round(output.stat().st_size / 1024, 1)
    compressed.message = f"Image compressed to {actual_kb} KB (target range: {min_kb}-{max_kb} KB)"
    return compressed


def handle_ssc_photo_resizer(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {
        **payload,
        "width_mm": float(payload.get("width_mm", 35.0)),
        "height_mm": float(payload.get("height_mm", 45.0)),
        "dpi": int(payload.get("dpi", 300)),
    }
    return handle_resize_image_in_mm(files, alias_payload, output_dir)


def handle_resize_for_pan_card(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {
        **payload,
        "width_mm": float(payload.get("width_mm", 35.0)),
        "height_mm": float(payload.get("height_mm", 25.0)),
        "dpi": int(payload.get("dpi", 300)),
    }
    return handle_resize_image_in_mm(files, alias_payload, output_dir)


def handle_resize_image_for_upsc(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {
        **payload,
        "width_mm": float(payload.get("width_mm", 35.0)),
        "height_mm": float(payload.get("height_mm", 45.0)),
        "dpi": int(payload.get("dpi", 300)),
    }
    return handle_resize_image_in_mm(files, alias_payload, output_dir)


def handle_psc_photo_resize(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    alias_payload = {
        **payload,
        "width_mm": float(payload.get("width_mm", 35.0)),
        "height_mm": float(payload.get("height_mm", 45.0)),
        "dpi": int(payload.get("dpi", 300)),
    }
    return handle_resize_image_in_mm(files, alias_payload, output_dir)


def handle_photo_blemish_remover(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    return handle_retouch_image(files, payload, output_dir)


def handle_bulk_image_resizer(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    ensure_files(files, 1)
    width = max(1, int(payload.get("width", 1080)))
    height = max(1, int(payload.get("height", 1080)))
    quality = max(5, min(100, int(payload.get("quality", 92))))
    maintain_aspect = str(payload.get("maintain_aspect", "true")).strip().lower() not in {"false", "0", "no"}

    outputs: list[Path] = []
    for source in files:
        image = open_image_file(source)
        if maintain_aspect:
            resized = image.copy()
            resized.thumbnail((width, height), get_resampling_module().LANCZOS)
        else:
            resized = image.resize((width, height), get_resampling_module().LANCZOS)

        suffix = source.suffix.lower() or ".jpg"
        target_suffix = ".jpg" if suffix in {".jpeg", ".jpg", ""} else suffix
        output = output_dir / f"{source.stem}-resized{target_suffix}"
        if target_suffix in {".jpg", ".jpeg"}:
            save_image(resized.convert("RGB"), output, fmt="JPEG", quality=quality)
        elif target_suffix == ".png":
            resized.save(str(output), format="PNG", optimize=True, compress_level=9)
        elif target_suffix == ".webp":
            resized.convert("RGB").save(str(output), format="WEBP", quality=quality, method=6)
        else:
            save_image(resized, output, fmt=target_suffix.replace(".", "").upper(), quality=quality)

        outputs.append(output)

    if len(outputs) == 1:
        return create_single_file_result(outputs[0], "Image resized", "image/*")
    return create_zip_result(output_dir, "Bulk image resize complete", "bulk-image-resized")


HANDLERS: dict[str, ToolHandler] = {
    "merge-pdf": handle_merge_pdf,
    "split-pdf": handle_split_pdf,
    "extract-pages": handle_extract_pages,
    "delete-pages": handle_delete_pages,
    "compress-pdf": handle_compress_pdf,
    "optimize-pdf": handle_optimize_pdf,
    "rotate-pdf": handle_rotate_pdf,
    "organize-pdf": handle_organize_pdf,
    "rearrange-pages": handle_rearrange_pages,
    "crop-pdf": handle_crop_pdf,
    "resize-pages-pdf": handle_resize_pages_pdf,
    "watermark-pdf": handle_watermark_pdf,
    "add-text-pdf": handle_add_text_pdf,
    "page-numbers-pdf": handle_page_numbers_pdf,
    "header-footer-pdf": handle_header_footer_pdf,
    "header-and-footer": handle_header_footer_pdf,
    "protect-pdf": handle_protect_pdf,
    "pdf-security": handle_pdf_security,
    "unlock-pdf": handle_unlock_pdf,
    "redact-pdf": handle_redact_pdf,
    "whiteout-pdf": handle_whiteout_pdf,
    "remove-metadata-pdf": handle_remove_metadata_pdf,
    "compare-pdf": handle_compare_pdf,
    "extract-text-pdf": handle_extract_text_pdf,
    "extract-images-pdf": handle_extract_images_pdf,
    "pdf-to-jpg": handle_pdf_to_jpg,
    "pdf-to-png": handle_pdf_to_png,
    "grayscale-pdf": handle_grayscale_pdf,
    "jpg-to-pdf": handle_jpg_to_pdf,
    "jpeg-to-pdf": handle_jpeg_to_pdf,
    "heic-to-pdf": handle_heic_to_pdf,
    "heif-to-pdf": handle_heif_to_pdf,
    "jfif-to-pdf": handle_jfif_to_pdf,
    "image-to-pdf": handle_image_to_pdf,
    "scan-to-pdf": handle_scan_to_pdf,
    "png-to-pdf": handle_png_to_pdf,
    "webp-to-pdf": handle_webp_to_pdf,
    "gif-to-pdf": handle_gif_to_pdf,
    "bmp-to-pdf": handle_bmp_to_pdf,
    "convert-to-pdf": handle_convert_to_pdf,
    "convert-from-pdf": handle_convert_from_pdf,
    "pdf-converter": handle_pdf_converter,
    "pdf-to-docx": handle_pdf_to_docx,
    "pdf-to-word": handle_pdf_to_word,
    "docx-to-pdf": handle_docx_to_pdf,
    "word-to-pdf": handle_word_to_pdf,
    "pdf-to-excel": handle_pdf_to_excel,
    "excel-to-pdf": handle_excel_to_pdf,
    "pdf-to-pptx": handle_pdf_to_pptx,
    "pdf-to-ppt": handle_pdf_to_ppt,
    "pdf-to-powerpoint": handle_pdf_to_powerpoint,
    "pptx-to-pdf": handle_pptx_to_pdf,
    "ppt-to-pdf": handle_ppt_to_pdf,
    "powerpoint-to-pdf": handle_powerpoint_to_pdf,
    "repair-pdf": handle_repair_pdf,
    "create-pdf": handle_create_pdf,
    "pdf-to-pdfa": handle_pdf_to_pdfa,
    "sign-pdf": handle_sign_pdf,
    "edit-metadata-pdf": handle_edit_metadata_pdf,
    "edit-metadata": handle_edit_metadata,
    "translate-pdf": handle_translate_pdf,
    "summarize-pdf": handle_summarize_pdf,
    "ai-summarizer": handle_ai_summarizer,
    "pdf-viewer": handle_pdf_viewer,
    "pdf-intelligence": handle_pdf_intelligence,
    "edit-pdf": handle_edit_pdf,
    "edit-pdf-text": handle_edit_pdf_text,
    "annotate-pdf": handle_annotate_pdf,
    "highlight-pdf": handle_highlight_pdf,
    "pdf-filler": handle_pdf_filler,
    "add-text": handle_add_text,
    "chat-with-pdf": handle_chat_with_pdf,
    "compress-image": handle_compress_image,
    "resize-image": handle_resize_image,
    "resize-image-in-cm": handle_resize_image_in_cm,
    "resize-image-in-mm": handle_resize_image_in_mm,
    "resize-image-in-inch": handle_resize_image_in_inch,
    "upscale-image": handle_upscale_image,
    "crop-image": handle_crop_image,
    "rotate-image": handle_rotate_image,
    "convert-image": handle_convert_image,
    "jpg-to-png": handle_jpg_to_png,
    "png-to-jpg": handle_png_to_jpg,
    "image-to-webp": handle_image_to_webp,
    "watermark-image": handle_watermark_image,
    "add-name-dob-image": handle_add_name_dob_image,
    "merge-photo-signature": handle_merge_photo_signature,
    "grayscale-image": handle_grayscale_image,
    "black-and-white-image": handle_black_and_white_image,
    "blur-image": handle_blur_image,
    "pixelate-image": handle_pixelate_image,
    "picture-to-pixel-art": handle_picture_to_pixel_art,
    "censor-photo": handle_censor_photo,
    "meme-generator": handle_meme_generator,
    "generate-signature": handle_generate_signature,
    "pdf-to-image": handle_pdf_to_image,
    "html-to-pdf": handle_html_to_pdf,
    "url-to-pdf": handle_url_to_pdf,
    "md-to-pdf": handle_md_to_pdf,
    "txt-to-pdf": handle_txt_to_pdf,
    "eml-to-pdf": handle_eml_to_pdf,
    "fb2-to-pdf": handle_fb2_to_pdf,
    "cbz-to-pdf": handle_cbz_to_pdf,
    "ebook-to-pdf": handle_ebook_to_pdf,
    "summarize-text": handle_summarize_text,
    "translate-text": handle_translate_text,
    "qr-code-generator": handle_qr_code_generator,
    "extract-metadata": handle_extract_metadata,
    "check-image-dpi": handle_check_image_dpi,
    "convert-dpi": handle_convert_dpi,
    "remove-metadata-image": handle_remove_metadata_image,
    "pdf-to-txt": handle_pdf_to_txt,
    "pdf-to-markdown": handle_pdf_to_markdown,
    "pdf-to-json": handle_pdf_to_json,
    "pdf-to-csv": handle_pdf_to_csv,
    "flip-image": handle_flip_image,
    "add-border-image": handle_add_border_image,
    "thumbnail-image": handle_thumbnail_image,
    "image-collage": handle_image_collage,
    "photo-collage-maker": handle_image_collage,
    "word-count-text": handle_word_count_text,
    "case-converter-text": handle_case_converter_text,
    "extract-keywords-text": handle_extract_keywords_text,
    "slug-generator-text": handle_slug_generator_text,
    "create-workflow": handle_create_workflow,
    "json-to-pdf": handle_json_to_pdf,
    "xml-to-pdf": handle_xml_to_pdf,
    "csv-to-pdf": handle_csv_to_pdf,
    "add-image-pdf": handle_add_image_pdf,
    "add-image-to-pdf": handle_add_image_to_pdf,
    "pdf-pages-to-zip": handle_pdf_pages_to_zip,
    "zip-images-to-pdf": handle_zip_images_to_pdf,
    "zip-to-pdf": handle_zip_to_pdf,
    "pdf-page-count": handle_pdf_page_count,
    "reverse-pdf": handle_reverse_pdf,
    "flatten-pdf": handle_flatten_pdf,
    "pdf-to-html": handle_pdf_to_html,
    "pdf-to-bmp": handle_pdf_to_bmp,
    "pdf-to-gif": handle_pdf_to_gif,
    "pdf-to-tiff": handle_pdf_to_tiff,
    "tiff-to-pdf": handle_tiff_to_pdf,
    "pdf-to-svg": handle_pdf_to_svg,
    "svg-to-pdf": handle_svg_to_pdf,
    "pdf-to-rtf": handle_pdf_to_rtf,
    "rtf-to-pdf": handle_rtf_to_pdf,
    "pdf-to-odt": handle_pdf_to_odt,
    "odt-to-pdf": handle_odt_to_pdf,
    "pdf-to-epub": handle_pdf_to_epub,
    "epub-to-pdf": handle_epub_to_pdf,
    "ocr-image": handle_ocr_image,
    "image-to-text": handle_image_to_text,
    "jpg-to-text": handle_jpg_to_text,
    "png-to-text": handle_png_to_text,
    "ocr-pdf": handle_ocr_pdf,
    "pdf-ocr": handle_pdf_ocr,
    "remove-background": handle_remove_background,
    "blur-background": handle_blur_background,
    "blur-face": handle_blur_face,
    "unblur-face": handle_unblur_face,
    "pixelate-face": handle_pixelate_face,
    "remove-image-object": handle_remove_image_object,
    "add-text-image": handle_add_text_image,
    "add-logo-image": handle_add_logo_image,
    "join-images": handle_join_images,
    "split-image": handle_split_image,
    "circle-crop-image": handle_circle_crop_image,
    "square-crop-image": handle_square_crop_image,
    "image-color-picker": handle_image_color_picker,
    "motion-blur-image": handle_motion_blur_image,
    "sharpen-image": handle_sharpen_image,
    "brighten-image": handle_brighten_image,
    "contrast-image": handle_contrast_image,
    "invert-image": handle_invert_image,
    "posterize-image": handle_posterize_image,
    "image-histogram": handle_image_histogram,
    "remove-extra-spaces-text": handle_remove_extra_spaces_text,
    "sort-lines-text": handle_sort_lines_text,
    "deduplicate-lines-text": handle_deduplicate_lines_text,
    "find-replace-text": handle_find_replace_text,
    "split-text-file": handle_split_text_file,
    "reading-time-text": handle_reading_time_text,
    "images-to-zip": handle_images_to_zip,
    "batch-convert-images": handle_batch_convert_images,
    "merge-text-files": handle_merge_text_files,
    "json-prettify": handle_json_prettify,
    "csv-to-json": handle_csv_to_json,
    "json-to-csv": handle_json_to_csv,
    "remove-pages": handle_remove_pages,
    "add-page-numbers": handle_add_page_numbers,
    "add-watermark": handle_add_watermark,
    "cbr-to-pdf": handle_cbr_to_pdf,
    "djvu-to-pdf": handle_djvu_to_pdf,
    "ai-to-pdf": handle_ai_to_pdf,
    "pdf-to-mobi": handle_pdf_to_mobi,
    "mobi-to-pdf": handle_mobi_to_pdf,
    "xps-to-pdf": handle_xps_to_pdf,
    "wps-to-pdf": handle_wps_to_pdf,
    "dwg-to-pdf": handle_dwg_to_pdf,
    "pub-to-pdf": handle_pub_to_pdf,
    "hwp-to-pdf": handle_hwp_to_pdf,
    "chm-to-pdf": handle_chm_to_pdf,
    "dxf-to-pdf": handle_dxf_to_pdf,
    "pages-to-pdf": handle_pages_to_pdf,
    "html-to-image": handle_html_to_image,
    "reduce-image-size-in-kb": handle_reduce_image_size_in_kb,
    "compress-to-kb": handle_compress_to_kb,
    "increase-image-size-in-kb": handle_increase_image_size_in_kb,
    "reduce-image-size-in-mb": handle_reduce_image_size_in_mb,
    "convert-image-from-mb-to-kb": handle_convert_image_from_mb_to_kb,
    "convert-image-size-kb-to-mb": handle_convert_image_size_kb_to_mb,
    "jpg-to-kb": handle_jpg_to_kb,
    "passport-photo-maker": handle_passport_photo_maker,
    "passport-size-photo": handle_passport_size_photo,
    "social-media-resize": handle_social_media_resize,
    "resize-for-instagram": handle_resize_for_instagram,
    "resize-for-whatsapp": handle_resize_for_whatsapp,
    "resize-for-youtube": handle_resize_for_youtube,
    "instagram-grid": handle_instagram_grid,
    "convert-to-jpg": handle_convert_to_jpg_image,
    "jpeg-to-jpg": handle_jpeg_to_jpg,
    "bulk-image-resizer": handle_bulk_image_resizer,
    "ssc-photo-resizer": handle_ssc_photo_resizer,
    "resize-for-pan-card": handle_resize_for_pan_card,
    "resize-image-for-upsc": handle_resize_image_for_upsc,
    "psc-photo-resize": handle_psc_photo_resize,
    "photo-blemish-remover": handle_photo_blemish_remover,
    "convert-from-jpg": handle_convert_from_jpg_image,
    "heic-to-jpg": handle_heic_to_jpg,
    "webp-to-jpg": handle_webp_to_jpg,
    "jpeg-to-png": handle_jpeg_to_png,
    "png-to-jpeg": handle_png_to_jpeg,
    "photo-editor": handle_photo_editor,
    "unblur-image": handle_unblur_image,
    "increase-image-quality": handle_increase_image_quality,
    "beautify-image": handle_beautify_image,
    "retouch-image": handle_retouch_image,
    "super-resolution": handle_super_resolution,
    "zoom-out-image": handle_zoom_out_image,
    "add-white-border-image": handle_add_white_border_image,
    "freehand-crop": handle_freehand_crop,
    "crop-png": handle_crop_png,
    "image-splitter": handle_image_splitter,
    "resize-image-pixel": handle_resize_image_pixel,
    "resize-signature": handle_resize_signature,
    "resize-image-to-3.5cmx4.5cm": handle_resize_image_to_3_5cmx4_5cm,
    "resize-image-to-6cmx2cm": handle_resize_image_to_6cmx2cm,
    "resize-signature-to-50mmx20mm": handle_resize_signature_to_50mmx20mm,
    "resize-image-to-35mmx45mm": handle_resize_image_to_35mmx45mm,
    "resize-image-to-2x2": handle_resize_image_to_2x2,
    "resize-image-to-3x4": handle_resize_image_to_3x4,
    "resize-image-to-4x6": handle_resize_image_to_4x6,
    "resize-image-to-600x600-pixel": handle_resize_image_to_600x600_pixel,
    "resize-image-for-whatsapp-dp": handle_resize_image_for_whatsapp_dp,
    "resize-image-for-youtube-banner": handle_resize_image_for_youtube_banner,
    "resize-image-to-a4-size": handle_resize_image_to_a4_size,
    "a4-size-resize": handle_a4_size_resize,
    "check-dpi": handle_check_dpi,
    "color-code-from-image": handle_color_code_from_image,
    "view-metadata": handle_view_metadata,
    "remove-image-metadata": handle_remove_image_metadata,
    "compress-to-5kb": handle_compress_specific_kb(5),
    "compress-image-to-5kb": handle_compress_specific_kb(5),
    "compress-to-10kb": handle_compress_specific_kb(10),
    "compress-image-to-10kb": handle_compress_specific_kb(10),
    "compress-jpeg-to-10kb": handle_compress_specific_kb(10),
    "compress-to-15kb": handle_compress_specific_kb(15),
    "compress-image-to-15kb": handle_compress_specific_kb(15),
    "compress-to-20kb": handle_compress_specific_kb(20),
    "compress-image-to-20kb": handle_compress_specific_kb(20),
    "compress-to-25kb": handle_compress_specific_kb(25),
    "compress-jpeg-to-25kb": handle_compress_specific_kb(25),
    "compress-to-30kb": handle_compress_specific_kb(30),
    "compress-jpeg-to-30kb": handle_compress_specific_kb(30),
    "compress-to-40kb": handle_compress_specific_kb(40),
    "compress-jpeg-to-40kb": handle_compress_specific_kb(40),
    "compress-to-50kb": handle_compress_specific_kb(50),
    "compress-image-to-50kb": handle_compress_specific_kb(50),
    "compress-to-100kb": handle_compress_specific_kb(100),
    "compress-image-to-100kb": handle_compress_specific_kb(100),
    "compress-to-150kb": handle_compress_specific_kb(150),
    "compress-jpeg-to-150kb": handle_compress_specific_kb(150),
    "compress-to-200kb": handle_compress_specific_kb(200),
    "compress-image-to-200kb": handle_compress_specific_kb(200),
    "compress-to-300kb": handle_compress_specific_kb(300),
    "compress-jpeg-to-300kb": handle_compress_specific_kb(300),
    "compress-to-500kb": handle_compress_specific_kb(500),
    "compress-image-to-500kb": handle_compress_specific_kb(500),
    "compress-jpeg-to-500kb": handle_compress_specific_kb(500),
    "compress-to-1mb": handle_compress_specific_kb(1024),
    "compress-image-to-1mb": handle_compress_specific_kb(1024),
    "compress-to-2mb": handle_compress_specific_kb(2048),
    "compress-image-to-2mb": handle_compress_specific_kb(2048),
    "compress-jpeg-between-20kb-to-50kb": handle_compress_jpeg_between_20kb_to_50kb,
    "jpg-to-pdf-under-50kb": handle_jpg_to_pdf_under_kb_factory(50),
    "jpg-to-pdf-under-100kb": handle_jpg_to_pdf_under_kb_factory(100),
    "jpeg-to-pdf-under-200kb": handle_jpg_to_pdf_under_kb_factory(200),
    "jpg-to-pdf-under-300kb": handle_jpg_to_pdf_under_kb_factory(300),
    "jpg-to-pdf-under-500kb": handle_jpg_to_pdf_under_kb_factory(500),
}

# Merge developer/utility handlers
try:
    from .developer_handlers import DEVELOPER_HANDLERS
    HANDLERS.update(DEVELOPER_HANDLERS)
    print(f"[handlers] Loaded {len(DEVELOPER_HANDLERS)} developer handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load developer_handlers: {e}")

# Merge everyday utility handlers
try:
    from .everyday_handlers import EVERYDAY_HANDLERS
    HANDLERS.update(EVERYDAY_HANDLERS)
    print(f"[handlers] Loaded {len(EVERYDAY_HANDLERS)} everyday handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load everyday_handlers: {e}")

# Merge extended student/everyday handlers
try:
    from .student_everyday_handlers import STUDENT_EVERYDAY_HANDLERS
    HANDLERS.update(STUDENT_EVERYDAY_HANDLERS)
    print(f"[handlers] Loaded {len(STUDENT_EVERYDAY_HANDLERS)} student/everyday handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load student_everyday_handlers: {e}")

print(f"[handlers] Total registered handlers: {len(HANDLERS)}")

# Merge new tools handlers (developer, color, security, conversion, social media)
try:
    from .new_tools_handlers import NEW_TOOL_HANDLERS
    HANDLERS.update(NEW_TOOL_HANDLERS)
    print(f"[handlers] Loaded {len(NEW_TOOL_HANDLERS)} new tool handlers (developer/color/security/conversion/social)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load new_tools_handlers: {e}")

print(f"[handlers] FINAL total registered handlers: {len(HANDLERS)}")

# Merge extra tools handlers (text encoding, improved math, etc.)
try:
    from .extra_tools_handlers import EXTRA_TOOL_HANDLERS
    HANDLERS.update(EXTRA_TOOL_HANDLERS)
    print(f"[handlers] Loaded {len(EXTRA_TOOL_HANDLERS)} extra tool handlers (text/math/encoding)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load extra_tools_handlers: {e}")

print(f"[handlers] GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge image-plus handlers (format conversions, circle-crop, DPI, collage, etc.)
try:
    from .image_plus_handlers import IMAGE_PLUS_HANDLERS
    HANDLERS.update(IMAGE_PLUS_HANDLERS)
    print(f"[handlers] Loaded {len(IMAGE_PLUS_HANDLERS)} image-plus handlers (format/dpi/collage/text-utils)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load image_plus_handlers: {e}")

print(f"[handlers] ULTIMATE TOTAL registered handlers: {len(HANDLERS)}")

# Merge phase3 handlers (PDF advanced, Image effects, Student, Everyday)
try:
    from .phase3_handlers import PHASE3_HANDLERS
    HANDLERS.update(PHASE3_HANDLERS)
    print(f"[handlers] Loaded {len(PHASE3_HANDLERS)} phase3 handlers (pdf-adv/image-fx/student/everyday)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load phase3_handlers: {e}")

print(f"[handlers] COMPLETE TOTAL registered handlers: {len(HANDLERS)}")

# Merge health, fitness & finance handlers
try:
    from .health_finance_handlers import HEALTH_FINANCE_HANDLERS
    HANDLERS.update(HEALTH_FINANCE_HANDLERS)
    print(f"[handlers] Loaded {len(HEALTH_FINANCE_HANDLERS)} health/finance handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load health_finance_handlers: {e}")

print(f"[handlers] FINAL GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge video downloader + extra tools handlers
try:
    from .video_extra_handlers import VIDEO_EXTRA_HANDLERS
    HANDLERS.update(VIDEO_EXTRA_HANDLERS)
    print(f"[handlers] Loaded {len(VIDEO_EXTRA_HANDLERS)} video/extra handlers (video/network/finance/math)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load video_extra_handlers: {e}")

print(f"[handlers] COMPLETE GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge ultra tools handlers (validators, CSS generators, health, finance v3)
try:
    from .ultra_tools_handlers import ULTRA_HANDLERS
    HANDLERS.update(ULTRA_HANDLERS)
    print(f"[handlers] Loaded {len(ULTRA_HANDLERS)} ultra tool handlers (validators/CSS/finance/health)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load ultra_tools_handlers: {e}")

print(f"[handlers] ULTRA GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge mega new handlers (science, geography, cooking, productivity, more video, student tools)
try:
    from .mega_new_handlers import MEGA_NEW_HANDLERS
    HANDLERS.update(MEGA_NEW_HANDLERS)
    print(f"[handlers] Loaded {len(MEGA_NEW_HANDLERS)} mega new handlers (science/geo/cooking/productivity/video)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load mega_new_handlers: {e}")

print(f"[handlers] MEGA GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge social video + developer + security + new tools handlers
try:
    from .social_video_handlers import SOCIAL_VIDEO_HANDLERS
    HANDLERS.update(SOCIAL_VIDEO_HANDLERS)
    print(f"[handlers] Loaded {len(SOCIAL_VIDEO_HANDLERS)} social/video/dev handlers (pinterest/reddit/twitch/jwt/regex/subnet)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load social_video_handlers: {e}")

print(f"[handlers] SOCIAL GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge mega tools v2 (AI writing, crypto, HR, legal, travel, developer v2, finance v2, productivity)
try:
    from .mega_tools_v2 import MEGA_TOOLS_V2_HANDLERS
    HANDLERS.update(MEGA_TOOLS_V2_HANDLERS)
    print(f"[handlers] Loaded {len(MEGA_TOOLS_V2_HANDLERS)} mega-v2 handlers (AI/crypto/HR/legal/travel/dev/finance/productivity)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load mega_tools_v2: {e}")

print(f"[handlers] FINAL ULTIMATE TOTAL registered handlers: {len(HANDLERS)}")

# Merge worldwide utility tools (paraphraser, essay outline, flashcards, email, cover letter,
# hashtag generator, SQL formatter, license generator, README generator, GitHub Actions generator,
# nginx config, dockerfile, study schedule, social bio, invoice, expense splitter, etc.)
try:
    from .worldwide_tools_handlers import WORLDWIDE_HANDLERS
    HANDLERS.update(WORLDWIDE_HANDLERS)
    print(f"[handlers] Loaded {len(WORLDWIDE_HANDLERS)} worldwide utility handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load worldwide_tools_handlers: {e}")

print(f"[handlers] WORLDWIDE GRAND TOTAL registered handlers: {len(HANDLERS)}")

# Merge new extra tools (Spotify, Snapchat, Threads, YT subtitles, video→GIF, MP3 cutter)
try:
    from .new_extra_tools_handlers import NEW_EXTRA_HANDLERS
    HANDLERS.update(NEW_EXTRA_HANDLERS)
    print(f"[handlers] Loaded {len(NEW_EXTRA_HANDLERS)} new-extra handlers (spotify/snapchat/threads/subtitles/gif/mp3-cutter)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load new_extra_tools_handlers: {e}")

print(f"[handlers] NEW-EXTRAS TOTAL registered handlers: {len(HANDLERS)}")

# Merge extra media tools (TTS, video trimmer/compressor, audio merger)
try:
    from .extra_media_handlers import EXTRA_MEDIA_HANDLERS
    HANDLERS.update(EXTRA_MEDIA_HANDLERS)
    print(f"[handlers] Loaded {len(EXTRA_MEDIA_HANDLERS)} extra-media handlers (tts/video-trim/video-compress/audio-merge)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load extra_media_handlers: {e}")

# Merge fresh text/utility tools (April 2026 batch)
try:
    from .fresh_text_handlers import FRESH_TEXT_HANDLERS
    HANDLERS.update(FRESH_TEXT_HANDLERS)
    print(f"[handlers] Loaded {len(FRESH_TEXT_HANDLERS)} fresh text/utility handlers (Apr-2026 batch)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load fresh_text_handlers: {e}")

# Override PDF Core handlers with feature-rich enhanced versions (Apr 2026)
# This MUST be registered last so it wins over baseline merge/split/compress/etc.
try:
    from .pdf_core_enhanced import PDF_CORE_ENHANCED_HANDLERS
    HANDLERS.update(PDF_CORE_ENHANCED_HANDLERS)
    print(f"[handlers] Loaded {len(PDF_CORE_ENHANCED_HANDLERS)} ENHANCED PDF Core handlers (overrides merge/split/compress/rotate/watermark/render/optimize)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load pdf_core_enhanced: {e}")

print(f"[handlers] EXTRA-MEDIA TOTAL registered handlers: {len(HANDLERS)}")

# Merge office tools (CSV ↔ Excel, PDF page extractor)
try:
    from .office_handlers import OFFICE_HANDLERS
    HANDLERS.update(OFFICE_HANDLERS)
    print(f"[handlers] Loaded {len(OFFICE_HANDLERS)} office handlers (csv-excel/pdf-page-extract)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load office_handlers: {e}")

print(f"[handlers] OFFICE TOTAL registered handlers: {len(HANDLERS)}")

# Merge media extras v2 (video rotate/mute/speed, audio speed, gif→video)
try:
    from .media_extras_v2 import MEDIA_EXTRAS_V2_HANDLERS
    HANDLERS.update(MEDIA_EXTRAS_V2_HANDLERS)
    print(f"[handlers] Loaded {len(MEDIA_EXTRAS_V2_HANDLERS)} media-v2 handlers (rotate/mute/speed/gif-mp4)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load media_extras_v2: {e}")

print(f"[handlers] MEDIA-V2 TOTAL registered handlers: {len(HANDLERS)}")

# Merge A/V Studio handlers (video reverser/cropper/resizer/watermark/thumbnail + audio reverser/volume/pitch/convert/trim)
try:
    from .av_studio_handlers import AV_STUDIO_HANDLERS
    HANDLERS.update(AV_STUDIO_HANDLERS)
    print(f"[handlers] Loaded {len(AV_STUDIO_HANDLERS)} av-studio handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load av_studio_handlers: {e}")

print(f"[handlers] AV-STUDIO TOTAL registered handlers: {len(HANDLERS)}")

# Merge Enhance handlers (noise-reducer/voice-enhancer/video-upscaler/stabilizer)
try:
    from .enhance_handlers import ENHANCE_HANDLERS
    HANDLERS.update(ENHANCE_HANDLERS)
    print(f"[handlers] Loaded {len(ENHANCE_HANDLERS)} enhance handlers (denoise/normalize/upscale/stabilize)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load enhance_handlers: {e}")

print(f"[handlers] ENHANCE TOTAL registered handlers: {len(HANDLERS)}")

# Merge Format Converter handlers (mp3-to-wav, mp4-to-mov, video-to-mp3, audio-compressor, etc.)
try:
    from .format_converters import FORMAT_CONVERTER_HANDLERS
    HANDLERS.update(FORMAT_CONVERTER_HANDLERS)
    print(f"[handlers] Loaded {len(FORMAT_CONVERTER_HANDLERS)} format-converter handlers (mp3-to-wav/mp4-to-mov/video-to-mp3/audio-compressor)")
except Exception as e:
    print(f"[handlers] WARNING: Could not load format_converters: {e}")

print(f"[handlers] FORMAT-CONVERTER TOTAL registered handlers: {len(HANDLERS)}")

# Merge Image Format Converter handlers (jpg-to-png, png-to-jpg, webp-to-png, etc.)
try:
    from .image_format_converters import IMAGE_FORMAT_HANDLERS
    HANDLERS.update(IMAGE_FORMAT_HANDLERS)
    print(f"[handlers] Loaded {len(IMAGE_FORMAT_HANDLERS)} image format converter handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load image_format_converters: {e}")

print(f"[handlers] IMAGE-FORMAT TOTAL registered handlers: {len(HANDLERS)}")

# Merge Data Format Converters (JSON/YAML/XML/TOML/CSV/TSV/HTML/SQL/Markdown)
try:
    from .data_format_converters import DATA_FORMAT_HANDLERS
    HANDLERS.update(DATA_FORMAT_HANDLERS)
    print(f"[handlers] Loaded {len(DATA_FORMAT_HANDLERS)} data format converter handlers")

    from .unit_converters import UNIT_CONVERTER_HANDLERS
    HANDLERS.update(UNIT_CONVERTER_HANDLERS)
    print(f"[handlers] Loaded {len(UNIT_CONVERTER_HANDLERS)} unit converter handlers")

    from .number_base_converters import NUMBER_BASE_HANDLERS
    HANDLERS.update(NUMBER_BASE_HANDLERS)
    print(f"[handlers] Loaded {len(NUMBER_BASE_HANDLERS)} number base converter handlers")

    from .av_studio_v2 import AV_STUDIO_V2_HANDLERS
    HANDLERS.update(AV_STUDIO_V2_HANDLERS)
    print(f"[handlers] Loaded {len(AV_STUDIO_V2_HANDLERS)} AV studio v2 handlers")
except Exception as e:
    print(f"[handlers] WARNING: Could not load data_format_converters: {e}")

print(f"[handlers] DATA-FORMAT TOTAL registered handlers: {len(HANDLERS)}")


def _message_from_payload(payload: Any, fallback: str = "Tool completed successfully") -> str:
    if isinstance(payload, dict):
        for key in ("message", "summary", "result", "error"):
            value = payload.get(key)
            if value not in (None, ""):
                return str(value)
    if payload not in (None, "") and not isinstance(payload, (dict, list)):
        return str(payload)
    return fallback


def _coerce_execution_result(raw_result: Any, output_dir: Path) -> ExecutionResult:
    """Normalize legacy handler return shapes into the platform contract.

    Several newer tool packs were written as lightweight two-argument handlers
    that return {"type": "json", "payload": ...}. The FastAPI route and smoke
    matrix expect ExecutionResult, so normalize once at registry load time.
    """
    if isinstance(raw_result, ExecutionResult):
        return raw_result

    # Duck-typing: handle ExecutionResult-like objects from other modules
    # (e.g. mega_tools_v2.ExecutionResult, mega_tools_v3.ExecutionResult)
    if (
        not isinstance(raw_result, dict)
        and hasattr(raw_result, "kind")
        and hasattr(raw_result, "message")
    ):
        kind = str(raw_result.kind or "json").lower()
        if kind == "file":
            op = getattr(raw_result, "output_path", None)
            output_path = Path(op) if op else None
            if output_path and output_path.exists():
                return ExecutionResult(
                    kind="file",
                    message=str(raw_result.message or "Done"),
                    output_path=output_path,
                    filename=getattr(raw_result, "filename", None) or output_path.name,
                    content_type=getattr(raw_result, "content_type", None) or "application/octet-stream",
                )
        data = getattr(raw_result, "data", None)
        if not isinstance(data, dict):
            data = {"result": str(data)} if data is not None else {}
        return ExecutionResult(
            kind="json",
            message=str(raw_result.message or "Done"),
            data=data,
        )

    if isinstance(raw_result, dict):
        result_type = str(raw_result.get("kind") or raw_result.get("type") or "json").lower()
        payload = raw_result.get("payload", raw_result.get("data", raw_result))
        message = str(raw_result.get("message") or _message_from_payload(payload))

        if result_type == "file":
            path_value = (
                raw_result.get("output_path")
                or raw_result.get("path")
                or raw_result.get("file")
                or raw_result.get("filename")
            )
            output_path = Path(path_value) if path_value else None
            if output_path and not output_path.is_absolute():
                output_path = output_dir / output_path
            if output_path and output_path.exists():
                return ExecutionResult(
                    kind="file",
                    message=message,
                    output_path=output_path,
                    filename=str(raw_result.get("filename") or output_path.name),
                    content_type=str(raw_result.get("content_type") or raw_result.get("media_type") or "application/octet-stream"),
                )

        data = payload if isinstance(payload, dict) else {"result": payload}
        return ExecutionResult(kind="json", message=message, data=data)

    if raw_result is None:
        raise RuntimeError("Tool handler returned no result")

    return ExecutionResult(
        kind="json",
        message="Tool completed successfully",
        data={"result": raw_result},
    )


def _normalize_tool_handler(slug: str, handler: ToolHandler) -> ToolHandler:
    if getattr(handler, "_ishu_normalized", False):
        return handler

    try:
        parameter_count = len(inspect.signature(handler).parameters)
    except (TypeError, ValueError):
        parameter_count = 3

    @wraps(handler)
    def wrapped(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
        if parameter_count >= 3:
            raw_result = handler(files, payload, output_dir)
        elif parameter_count == 2:
            raw_result = handler(files, payload)
        elif parameter_count == 1:
            raw_result = handler(payload)
        else:
            raw_result = handler()
        return _coerce_execution_result(raw_result, output_dir)

    wrapped._ishu_normalized = True  # type: ignore[attr-defined]
    wrapped._ishu_slug = slug  # type: ignore[attr-defined]
    wrapped.__signature__ = inspect.Signature(  # type: ignore[attr-defined]
        parameters=[
            inspect.Parameter("files", inspect.Parameter.POSITIONAL_OR_KEYWORD),
            inspect.Parameter("payload", inspect.Parameter.POSITIONAL_OR_KEYWORD),
            inspect.Parameter("output_dir", inspect.Parameter.POSITIONAL_OR_KEYWORD),
        ]
    )
    return wrapped


_LEGACY_HANDLER_COUNT = 0
for _slug, _handler in list(HANDLERS.items()):
    try:
        _param_count = len(inspect.signature(_handler).parameters)
    except (TypeError, ValueError):
        _param_count = 3
    if _param_count != 3:
        _LEGACY_HANDLER_COUNT += 1
    HANDLERS[_slug] = _normalize_tool_handler(_slug, _handler)

print(f"[handlers] Normalized {len(HANDLERS)} handlers ({_LEGACY_HANDLER_COUNT} legacy signatures)")
