# Production-Ready Tool Handlers - Complete Implementation
# All tools fully functional with proper error handling

from __future__ import annotations

import io
import json
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any

import fitz
from docx import Document
from fastapi import HTTPException
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageDraw, ImageEnhance, ImageFilter
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import A4, letter
from reportlab.pdfgen import canvas as pdf_canvas

from .handlers import (
    ExecutionResult,
    create_single_file_result,
    create_zip_result,
    ensure_files,
    extract_pdf_text,
    open_image_file,
    save_image,
    text_to_pdf,
    try_libreoffice_convert_to_pdf,
)


# ============================================================================
# OFFICE DOCUMENT CONVERSIONS (PDF ↔ Word, Excel, PowerPoint)
# ============================================================================

def handle_pdf_to_word(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert PDF to DOCX with proper text extraction and formatting"""
    ensure_files(files, 1)
    
    try:
        # Try pdf2docx first (best quality)
        from pdf2docx import Converter
        
        output = output_dir / f"{files[0].stem}.docx"
        cv = Converter(str(files[0]))
        cv.convert(str(output))
        cv.close()
        
        return create_single_file_result(output, "PDF converted to Word", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    except Exception:
        # Fallback: Extract text and create DOCX
        text = extract_pdf_text(files[0])
        if not text:
            raise HTTPException(status_code=400, detail="No text content found in PDF")
        
        doc = Document()
        doc.add_heading(files[0].stem, 0)
        
        for paragraph in text.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        
        output = output_dir / f"{files[0].stem}.docx"
        doc.save(str(output))
        
        return create_single_file_result(output, "PDF converted to Word", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


def handle_word_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert Word (DOCX/DOC) to PDF"""
    ensure_files(files, 1)
    
    # Try LibreOffice first
    result = try_libreoffice_convert_to_pdf(files[0], output_dir, timeout_seconds=120)
    if result and result.exists():
        return create_single_file_result(result, "Word converted to PDF", "application/pdf")
    
    # Fallback: Extract text and create PDF
    try:
        doc = Document(str(files[0]))
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        text = '\n\n'.join(text_parts)
        if not text:
            raise HTTPException(status_code=400, detail="No content found in Word document")
        
        output = output_dir / f"{files[0].stem}.pdf"
        text_to_pdf(text, output, title=files[0].stem)
        
        return create_single_file_result(output, "Word converted to PDF", "application/pdf")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to convert Word document: {str(e)}")


def handle_pdf_to_powerpoint(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert PDF to PowerPoint"""
    ensure_files(files, 1)
    
    # Extract text from PDF
    text = extract_pdf_text(files[0])
    if not text:
        raise HTTPException(status_code=400, detail="No text content found in PDF")
    
    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = files[0].stem
    subtitle.text = "Converted from PDF"
    
    # Content slides - split text into chunks
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    
    for i in range(0, len(paragraphs), 5):
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"Page {i//5 + 1}"
        content.text = '\n\n'.join(paragraphs[i:i+5])
    
    output = output_dir / f"{files[0].stem}.pptx"
    prs.save(str(output))
    
    return create_single_file_result(output, "PDF converted to PowerPoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation")


def handle_powerpoint_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert PowerPoint to PDF"""
    ensure_files(files, 1)
    
    # Try LibreOffice first
    result = try_libreoffice_convert_to_pdf(files[0], output_dir, timeout_seconds=120)
    if result and result.exists():
        return create_single_file_result(result, "PowerPoint converted to PDF", "application/pdf")
    
    # Fallback: Extract text and create PDF
    try:
        prs = Presentation(str(files[0]))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        text = '\n\n'.join(text_parts)
        if not text:
            raise HTTPException(status_code=400, detail="No content found in PowerPoint")
        
        output = output_dir / f"{files[0].stem}.pdf"
        text_to_pdf(text, output, title=files[0].stem)
        
        return create_single_file_result(output, "PowerPoint converted to PDF", "application/pdf")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to convert PowerPoint: {str(e)}")


def handle_pdf_to_excel(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert PDF to Excel"""
    ensure_files(files, 1)
    
    try:
        # Try tabula for table extraction
        import tabula
        
        dfs = tabula.read_pdf(str(files[0]), pages='all', multiple_tables=True)
        
        if not dfs:
            raise Exception("No tables found")
        
        output = output_dir / f"{files[0].stem}.xlsx"
        
        with pd.ExcelWriter(str(output), engine='openpyxl') as writer:
            for i, df in enumerate(dfs, 1):
                sheet_name = f"Table_{i}"
                df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        return create_single_file_result(output, "PDF converted to Excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    except Exception:
        # Fallback: Extract text and create simple Excel
        text = extract_pdf_text(files[0])
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        
        for row_num, line in enumerate(lines, 1):
            ws.cell(row=row_num, column=1, value=line)
        
        output = output_dir / f"{files[0].stem}.xlsx"
        wb.save(str(output))
        
        return create_single_file_result(output, "PDF converted to Excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def handle_excel_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert Excel to PDF"""
    ensure_files(files, 1)
    
    # Try LibreOffice first
    result = try_libreoffice_convert_to_pdf(files[0], output_dir, timeout_seconds=120)
    if result and result.exists():
        return create_single_file_result(result, "Excel converted to PDF", "application/pdf")
    
    # Fallback: Extract data and create PDF
    try:
        wb = load_workbook(str(files[0]), data_only=True)
        
        output = output_dir / f"{files[0].stem}.pdf"
        c = pdf_canvas.Canvas(str(output), pagesize=letter)
        width, height = letter
        
        y = height - 40
        c.setFont("Helvetica-Bold", 14)
        c.drawString(40, y, files[0].stem)
        y -= 30
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            c.setFont("Helvetica-Bold", 12)
            c.drawString(40, y, f"Sheet: {sheet_name}")
            y -= 20
            
            c.setFont("Helvetica", 10)
            
            for row in ws.iter_rows(values_only=True):
                if y < 40:
                    c.showPage()
                    y = height - 40
                    c.setFont("Helvetica", 10)
                
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    c.drawString(40, y, row_text[:100])
                    y -= 15
            
            y -= 10
        
        c.save()
        
        return create_single_file_result(output, "Excel converted to PDF", "application/pdf")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to convert Excel: {str(e)}")


# ============================================================================
# ADVANCED IMAGE TOOLS
# ============================================================================

def handle_remove_background(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Remove background from image using rembg"""
    ensure_files(files, 1)
    
    try:
        from rembg import remove
        
        input_image = Image.open(files[0])
        output_image = remove(input_image)
        
        output = output_dir / f"no-bg-{files[0].stem}.png"
        output_image.save(str(output), format='PNG')
        
        return create_single_file_result(output, "Background removed", "image/png")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to remove background: {str(e)}")


def handle_blur_face(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Blur faces in image using OpenCV"""
    ensure_files(files, 1)
    
    try:
        import cv2
        import numpy as np
        
        # Load image
        image = cv2.imread(str(files[0]))
        if image is None:
            raise Exception("Failed to load image")
        
        # Load face cascade
        face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        
        # Convert to grayscale for detection
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        
        # Detect faces
        faces = face_cascade.detectMultiScale(gray, scaleFactor=1.1, minNeighbors=5, minSize=(30, 30))
        
        # Blur each face
        blur_factor = int(payload.get('blur_factor', 50))
        for (x, y, w, h) in faces:
            # Extract face region
            face = image[y:y+h, x:x+w]
            # Apply strong blur
            face = cv2.GaussianBlur(face, (blur_factor|1, blur_factor|1), 0)
            # Put blurred face back
            image[y:y+h, x:x+w] = face
        
        # Save result
        output = output_dir / f"blur-face-{files[0].name}"
        cv2.imwrite(str(output), image)
        
        return create_single_file_result(output, f"Blurred {len(faces)} face(s)", "image/*")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to blur faces: {str(e)}")


def handle_flip_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Flip image horizontally or vertically"""
    ensure_files(files, 1)
    
    direction = str(payload.get('direction', 'horizontal')).lower()
    
    image = Image.open(files[0])
    
    if direction == 'vertical':
        flipped = image.transpose(Image.FLIP_TOP_BOTTOM)
    else:  # horizontal
        flipped = image.transpose(Image.FLIP_LEFT_RIGHT)
    
    output = output_dir / f"flip-{files[0].name}"
    save_image(flipped, output)
    
    return create_single_file_result(output, f"Image flipped {direction}ly", "image/*")


def handle_add_border_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Add border to image"""
    ensure_files(files, 1)
    
    border_size = int(payload.get('border_size', 20))
    border_color = str(payload.get('border_color', '#000000'))
    
    from PIL import ImageColor
    
    image = Image.open(files[0])
    
    # Create new image with border
    new_width = image.width + 2 * border_size
    new_height = image.height + 2 * border_size
    
    bordered = Image.new('RGB', (new_width, new_height), ImageColor.getrgb(border_color))
    bordered.paste(image, (border_size, border_size))
    
    output = output_dir / f"border-{files[0].name}"
    save_image(bordered, output)
    
    return create_single_file_result(output, "Border added to image", "image/*")


def handle_thumbnail_image(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Generate thumbnail from image"""
    ensure_files(files, 1)
    
    max_size = int(payload.get('max_size', 200))
    
    image = Image.open(files[0])
    image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
    
    output = output_dir / f"thumb-{files[0].name}"
    save_image(image, output, quality=85)
    
    return create_single_file_result(output, "Thumbnail generated", "image/*")


def handle_image_collage(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Create collage from multiple images"""
    ensure_files(files, 2)
    
    cols = int(payload.get('cols', 2))
    spacing = int(payload.get('spacing', 10))
    
    # Load all images
    images = [Image.open(f) for f in files]
    
    # Calculate grid dimensions
    rows = (len(images) + cols - 1) // cols
    
    # Find max dimensions
    max_width = max(img.width for img in images)
    max_height = max(img.height for img in images)
    
    # Create collage
    collage_width = cols * max_width + (cols + 1) * spacing
    collage_height = rows * max_height + (rows + 1) * spacing
    
    collage = Image.new('RGB', (collage_width, collage_height), 'white')
    
    for idx, img in enumerate(images):
        row = idx // cols
        col = idx % cols
        x = col * max_width + (col + 1) * spacing
        y = row * max_height + (row + 1) * spacing
        collage.paste(img, (x, y))
    
    output = output_dir / "collage.jpg"
    collage.save(str(output), quality=90)
    
    return create_single_file_result(output, "Collage created", "image/jpeg")


# ============================================================================
# TEXT OPERATIONS
# ============================================================================

def handle_word_count_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Count words, sentences, and characters in text"""
    from .handlers import read_text_input
    
    text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide text to analyze")
    
    words = re.findall(r'\b\w+\b', text)
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if s.strip()]
    
    stats = {
        "characters": len(text),
        "characters_no_spaces": len(text.replace(' ', '')),
        "words": len(words),
        "sentences": len(sentences),
        "paragraphs": len([p for p in text.split('\n\n') if p.strip()]),
        "lines": len(text.split('\n')),
    }
    
    return ExecutionResult(
        kind="json",
        message="Text statistics generated",
        data=stats
    )


def handle_case_converter_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert text case"""
    from .handlers import read_text_input
    
    text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide text to convert")
    
    case_type = str(payload.get('case_type', 'lower')).lower()
    
    if case_type == 'upper':
        result = text.upper()
    elif case_type == 'lower':
        result = text.lower()
    elif case_type == 'title':
        result = text.title()
    elif case_type == 'sentence':
        sentences = re.split(r'([.!?]\s+)', text)
        result = ''.join(s.capitalize() if i % 2 == 0 else s for i, s in enumerate(sentences))
    else:
        result = text
    
    output = output_dir / "converted-text.txt"
    output.write_text(result, encoding='utf-8')
    
    return create_single_file_result(output, f"Text converted to {case_type} case", "text/plain")


def handle_extract_keywords_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Extract keywords from text"""
    from .handlers import read_text_input, extract_top_keywords
    
    text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide text to analyze")
    
    limit = int(payload.get('limit', 10))
    keywords = extract_top_keywords(text, limit)
    
    return ExecutionResult(
        kind="json",
        message="Keywords extracted",
        data={"keywords": keywords}
    )


def handle_slug_generator_text(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Generate URL-friendly slug from text"""
    from .handlers import read_text_input
    
    text = read_text_input(files, payload)
    if not text:
        raise HTTPException(status_code=400, detail="Provide text to convert to slug")
    
    # Convert to lowercase and replace spaces/special chars with hyphens
    slug = text.lower()
    slug = re.sub(r'[^\w\s-]', '', slug)
    slug = re.sub(r'[-\s]+', '-', slug)
    slug = slug.strip('-')
    
    return ExecutionResult(
        kind="json",
        message="Slug generated",
        data={"original": text[:100], "slug": slug}
    )


# ============================================================================
# ARCHIVE AND BATCH OPERATIONS
# ============================================================================

def handle_pdf_pages_to_zip(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Export PDF pages as images in ZIP"""
    ensure_files(files, 1)
    
    document = fitz.open(str(files[0]))
    
    for index, page in enumerate(document, 1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        output = output_dir / f"page-{index}.png"
        pix.save(str(output))
    
    document.close()
    
    return create_zip_result(output_dir, "PDF pages exported to ZIP", "pdf-pages")


def handle_zip_images_to_pdf(files: list[Path], payload: dict[str, Any], output_dir: Path) -> ExecutionResult:
    """Convert images from ZIP to PDF"""
    ensure_files(files, 1)
    
    if not files[0].suffix.lower() == '.zip':
        raise HTTPException(status_code=400, detail="Upload a ZIP file containing images")
    
    # Extract ZIP
    extract_dir = output_dir / "extracted"
    extract_dir.mkdir(exist_ok=True)
    
    with zipfile.ZipFile(files[0], 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
    
    # Find all images
    image_extensions = {'.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif', '.tif', '.tiff'}
    image_files = sorted([
        f for f in extract_dir.rglob('*')
        if f.is_file() and f.suffix.lower() in image_extensions
    ])
    
    if not image_files:
        raise HTTPException(status_code=400, detail="No images found in ZIP file")
    
    # Convert to PDF
    from .handlers import image_paths_to_pdf
    
    output = output_dir / "images.pdf"
    image_paths_to_pdf(image_files, output)
    
    return create_single_file_result(output, f"Converted {len(image_files)} images to PDF", "application/pdf")


# Register all handlers
PRODUCTION_HANDLERS = {
    "pdf-to-word": handle_pdf_to_word,
    "word-to-pdf": handle_word_to_pdf,
    "pdf-to-powerpoint": handle_pdf_to_powerpoint,
    "powerpoint-to-pdf": handle_powerpoint_to_pdf,
    "pdf-to-excel": handle_pdf_to_excel,
    "excel-to-pdf": handle_excel_to_pdf,
    "remove-background": handle_remove_background,
    "blur-face": handle_blur_face,
    "flip-image": handle_flip_image,
    "add-border-image": handle_add_border_image,
    "thumbnail-image": handle_thumbnail_image,
    "image-collage": handle_image_collage,
    "word-count-text": handle_word_count_text,
    "case-converter-text": handle_case_converter_text,
    "extract-keywords-text": handle_extract_keywords_text,
    "slug-generator-text": handle_slug_generator_text,
    "pdf-pages-to-zip": handle_pdf_pages_to_zip,
    "zip-images-to-pdf": handle_zip_images_to_pdf,
}
